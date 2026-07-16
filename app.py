import sys
import os

# v9.2: required BEFORE anything else once the .spec is built with
# console=False (windowed mode, no visible console window). In that mode
# PyInstaller sets sys.stdout/sys.stderr to None (there is no console to
# write to) -- but this app calls print() extensively throughout for
# logging/debugging. Without this guard, the very first print() call would
# crash the app instantly with "AttributeError: 'NoneType' object has no
# attribute 'write'", before the GUI even has a chance to open. Redirecting
# to os.devnull makes every print() a harmless no-op instead. Must run
# before any other import that might print something at import time.
if sys.stdout is None:
    sys.stdout = open(os.devnull, "w")
if sys.stderr is None:
    sys.stderr = open(os.devnull, "w")

import math
import tkinter as tk
from tkinter import messagebox, ttk, filedialog
import webbrowser
import subprocess
import re
import sqlite3
import threading
import string
import unicodedata
import time
from datetime import datetime
import pandas as pd 
import docx # pip install python-docx
import pypdf # pip install pypdf
from pptx import Presentation # pip install python-pptx
import logging
import warnings
warnings.filterwarnings("ignore")
logging.basicConfig(level=logging.CRITICAL)
logging.getLogger().setLevel(logging.CRITICAL)
logging.getLogger("pypdf").setLevel(logging.CRITICAL)
logging.getLogger("pypdf._cmap").setLevel(logging.CRITICAL)
try:
    import xlrd
    xlrd.xlsx.ensure_elementtree_imported(False, False)
except:
    pass
logging.getLogger('xlrd').setLevel(logging.CRITICAL)

# --- CONFIGURATION ---
# Always find DB next to the exe (or .py script), regardless of working directory
if getattr(sys, 'frozen', False):
    _BASE_DIR = os.path.dirname(sys.executable)  # when running as exe
else:
    _BASE_DIR = os.path.dirname(os.path.abspath(__file__))  # when running as .py

DB_FILE = os.path.join(_BASE_DIR, 'search_data.db')
# v7.10 FIX: Search History used to be logged straight into DB_FILE
# (search_data.db) itself. That meant just typing a query -- even with
# --update data NEVER run -- created/grew search_data.db (a few KB, just
# the 'history' table). On next launch the ramp light saw that file
# "exists" and jumped from Red to Yellow, even though nothing was ever
# actually indexed. Search History now lives in its own separate file so
# merely searching can never make the ramp light think indexing happened.
HISTORY_DB_FILE = os.path.join(_BASE_DIR, 'search_history.db')

# v9.0: AI libraries (torch/transformers/sentence-transformers/einops) are
# now bundled directly into the exe at build time -- no more runtime
# 'ai_libs' folder, no more sys.path surgery, no more heal-attempt
# counters. See the top-of-file note near _load_semantic_model for why.

# v5.8: Stopword filter -- fixes irrelevant OR-fallback matches.
# Without this, a query like "simpack realtime relevant to HILS" splits into
# keywords ["simpack","realtime","relevant","to","hils"]. The 2-letter word
# "to" then substring-matches almost anything ending in ".toc"
# (Analysis-00.toc, EXE-00.toc, PYZ-00.toc, ...) via `"to" in filename`,
# flooding results with junk that has nothing to do with the real query.
# Filler words (English + Vietnamese) carry no filename-matching signal and
# are stripped before AND/OR keyword matching. If stripping would remove
# every keyword (rare edge case: the whole query was filler words), the
# original keyword list is used instead so the search never returns empty.
STOPWORDS = {
    # English fillers
    "a", "an", "the", "to", "of", "in", "on", "for", "and", "or", "is", "are",
    "was", "were", "be", "been", "with", "at", "by", "from", "as", "that",
    "this", "it", "its", "into", "about", "regarding", "re", "vs", "via",
    "relevant", "related", "concerning",
    # Vietnamese fillers
    "và", "của", "cho", "là", "các", "những", "với", "về", "đến", "tới",
    "cái", "này", "đó", "cùng", "hoặc", "trong", "trên", "dưới", "theo",
    "liên", "quan",
}


def _strip_stopwords(keywords):
    """Remove filler words from a keyword list; fall back to the original
    list if that would leave nothing (so a query never becomes empty)."""
    filtered = [k for k in keywords if k not in STOPWORDS]
    return filtered if filtered else keywords


def _kw_matches(kw, text_lower, whole_word=False):
    """Does `kw` match somewhere inside `text_lower`?
    whole_word=False (default): plain substring match -- original behavior,
        e.g. "adas" matches inside "readasync.xml" (buried mid-word).
    whole_word=True: `kw` only counts as a match when both the character
        immediately before and immediately after it are NOT ascii letters
        (i.e. digit/underscore/hyphen/dot/space/parenthesis/start-of-string/
        end-of-string all count as a boundary). This keeps matches like
        "ADAS_systems.pdf" or "VDIM_0ADAS_ESP" (ADAS sits next to a digit/
        underscore/dot) while rejecting "readasync.xml" or
        "ReadAStringExample.mlx" (adas sits buried between two letters)."""
    if not whole_word:
        return kw in text_lower
    start = 0
    n = len(text_lower)
    klen = len(kw)
    while True:
        idx = text_lower.find(kw, start)
        if idx == -1:
            return False
        before_ok = (idx == 0) or (not text_lower[idx - 1].isalpha())
        after_idx = idx + klen
        after_ok = (after_idx >= n) or (not text_lower[after_idx].isalpha())
        if before_ok and after_ok:
            return True
        start = idx + 1  # keep scanning -- an earlier occurrence might fail while a later one succeeds

SMALL_SIZE = "85x35" 
LARGE_SIZE = "390x35"
RESULT_SIZE = "1250x700"
HELP_SIZE = "630x630"
BG_COLOR = "#f4f5f7"
ENTRY_BG = "#ffffff"
TEXT_COLOR = "#1c1e21"
PLACE_COLOR = "#8a8d93"

READY_PH = "SR/NTSR/QA/BR/ER/IR/CRIT/JIRA/Trigram/SiteID/ContactID..."
READY_PH1 = "File/Folder/File content..."
READY_PH2 = "ESC, --exit, --quit"
WAIT_PH = "Updating Database, please wait..."

MAX_CONTENT_SIZE = 2 * 1024 * 1024  
MAX_CHARS_TO_INDEX = 60000 

# URLs
BASE_SR = "https://dsext001-eu1-215dsi0708-ifwe.3dexperience.3ds.com/#app:MAP-HEYJWSUBD/content:redirect=true&objtype=SR&objectid="
BASE_NTSR = "https://dsext001-eu1-215dsi0708-ifwe.3dexperience.3ds.com/#app:MAP-HEYJWSUBD/content:redirect=true&objtype=NTSR&objectid="
BASE_QA = "https://support.3ds.com/knowledge-base/?q=docid:"
BASE_BR = "https://support.3ds.com/knowledge-base/?q=docid:"
BASE_IR = "https://dsxdev-online.dsone.3ds.com/enovia/common/dsxRedirect.jsp?type=IR&name="
BASE_CRIT = "http://dsxcli/DB/&Class=CRITNAME&ID="
BASE_JIRA = "https://spck-jira.ux.dsone.3ds.com:8443/browse/"
BASE_DS = "https://eu1-215dsi0708-ifwe.3dexperience.3ds.com/#app:MAP-IQPLHVVAF/content:people="
BASE_ER = "https://dsxdev-online.dsone.3ds.com/enovia/common/dsxRedirect.jsp?type=*&name="
BASE_SIT = "http://dsxcli/DB/&Class=SIT&ID="
BASE_CTC = "http://dsxcli/DB/&Class=CTC&ID="
CMD_PATH = r"C:\Windows\System32\cmd.exe /k"

SKIP_FOLDERS = {
    # Windows system
    "windows", "$recycle.bin", "system volume information", "appdata",
    # Program folders
    "program files", "program files (x86)", "programdata",
    # Python installations
    "python39", "python311", "python312", "python314", "python3", "python",
    # Database / middleware
    "oracle19c", "oracle", "mysql", "postgresql",
    # Cloud sync — files may be cloud-only placeholders, opening triggers download
    "onedrive", "onedrive - personal", "sharepoint", "google drive", "dropbox", "box",
    # Dev / build artifacts
    "node_modules", ".git", "__pycache__", ".venv", "venv", "env",
    # 3rd party app data
    "3ds",
    # IT / Security
    "_it", "it", "security", "credentials", "vault", "secrets",
    "password", "passwords", "private", "confidential", "restricted",
}
# Files whose names match these patterns are never opened for content indexing
import re as _re_sec
_SENSITIVE_FILE_RE = _re_sec.compile(
    r'(?i)(password|passwd|credential|secret|private[_\-]?key|'
    r'id_rsa|vault|keystore|\.kdbx|keepass|lastpass|api[_\-]?key)'
)
BINARY_EXT = {".exe", ".dll", ".lib", ".obj", ".pyc", ".bin", ".jpg", ".png", ".gif", ".zip", ".7z", ".rar"}

# v2.5: Recognize a downloaded HuggingFace / SentenceTransformers model repo by
# its file signature (config.json alongside tokenizer/weight files) — this way
# ANY local model cache (jina, bge, e5-large, e5-base, or any future model,
# regardless of folder name) gets skipped during --update data, instead of
# being walked and partially indexed as if it were user documents.
_MODEL_REPO_WEIGHT_OR_TOKENIZER_FILES = {
    "tokenizer.json", "tokenizer_config.json", "pytorch_model.bin",
    "model.safetensors", "modules.json", "sentence_bert_config.json",
    "adapter_config.json",
}
def _looks_like_model_repo(dir_path):
    try:
        names = set(os.listdir(dir_path))
    except Exception:
        return False
    return "config.json" in names and bool(names & _MODEL_REPO_WEIGHT_OR_TOKENIZER_FILES)

# ── Semantic Search config ────────────────────────────────────────────────────
# ── Multi-model AI registry ───────────────────────────────────────────────────
# Each model has: weight folder (auto-detected under D:\mySearch\models or next to .py/.exe),
# its own table in search_data.db (BM25/content still share ONE DB — only
# embedding tables are split because each model produces different-dimension/space vectors),
# and its own encode convention (each model uses different prefix/task/instruction).
SEMANTIC_MODELS = {
    "jina_v3": {
        "label":        "Jina-v3",
        "dir_names":    ["jina-embeddings-v3", "jina-v3", "jina_v3"],
        "table":        "semantic_index_jina_v3",
        "trust_remote": True,   # jina-v3 requires trust_remote_code=True for LoRA task adapter
        "hf_repo":      "jinaai/jina-embeddings-v3",
    },
    "bge_gemma2": {
        "label":        "BGE-Gemma2",
        "dir_names":    ["bge-multilingual-gemma2", "bge-gemma2", "bge_gemma2"],
        "table":        "semantic_index_bge_gemma2",
        "trust_remote": True,
        # v3.4: BGE-Gemma2 is Gemma2-based and too heavy for small/laptop GPUs
        # (e.g. 4GB VRAM cards routinely OOM on it) — force it onto CPU always,
        # regardless of what GPU is available. jina_v3 still uses GPU normally.
        "force_cpu":    True,
        "hf_repo":      "BAAI/bge-multilingual-gemma2",
    },
}
DEFAULT_SEMANTIC_MODEL = "jina_v3"

# Root folder for models — relative to wherever the app itself lives
# (next to the .py / .exe), so it works the same on any machine/user
# without editing the source. Models must sit in a "models" subfolder
# next to the app; _BASE_DIR falls back to the app's own folder for
# older layouts where the model folders sit alongside the script directly.
_MODEL_ROOT_CANDIDATES = [
    os.path.join(_BASE_DIR, "models"),
    _BASE_DIR,
]

def _find_model_dir_for(model_key):
    """Auto-detect the weight folder for a model from _MODEL_ROOT_CANDIDATES."""
    info = SEMANTIC_MODELS.get(model_key, {})
    names = info.get("dir_names", [])
    candidates = []
    for root in _MODEL_ROOT_CANDIDATES:
        for n in names:
            candidates.append(os.path.join(root, n))
    for p in candidates:
        if os.path.isdir(p):
            return p
    return candidates[0] if candidates else ""  # return first candidate so load error is explicit

# Pre-resolve each model's path at startup (lazy but cached once)
SEMANTIC_MODEL_DIRS = {k: _find_model_dir_for(k) for k in SEMANTIC_MODELS}

SEMANTIC_MIN_WORDS = 3      # query needs at least N words to trigger semantic search
SEMANTIC_TOP_K     = 50     # retrieve top K semantic results

# v9.3: per-model threshold instead of one shared value. jina_v3 and
# bge_gemma2 are different architectures with different cosine-similarity
# score distributions -- the same cutoff can behave very differently
# between them (this is why bge_gemma2 was returning ZERO results for
# some queries where jina_v3 still returned some: bge_gemma2's best score
# for that query was likely just under 0.25). If the debug print in
# _semantic_search shows a model's best score consistently sitting just
# below its threshold, lower that model's entry here.
SEMANTIC_THRESHOLD_BY_MODEL = {
    "jina_v3":    0.25,
    # v9.3.1: lowered from 0.25 based on real observed data -- a genuinely
    # relevant result for a Vietnamese-no-diacritics query scored 0.231
    # with bge_gemma2 and was being filtered out entirely by the old shared
    # 0.25 threshold (while jina_v3 cleared it fine for the same query).
    # 0.20 leaves a little headroom below that observed score. If you see
    # bge_gemma2 results that are clearly NOT relevant sneaking in now,
    # raise this back up a bit; if good results are still being cut off,
    # check the [Semantic] debug log for the new best-score numbers and
    # lower it further.
    "bge_gemma2": 0.20,
}
SEMANTIC_THRESHOLD_DEFAULT = 0.25  # fallback for any model not listed above

# ── Currently active model for AI Search ────────────────────────────────────
# Changed via UI dropdown (see RealtimeSmartSearchApp._on_ai_model_change)
_sem_model_key = DEFAULT_SEMANTIC_MODEL

_sem_model  = None
_sem_ready  = False
_sem_device = "cpu"
_sem_loaded_key = None   # model key currently loaded in _sem_model (to detect when reload is needed)

# ── Vietnamese diacritics restoration (v9.4) ────────────────────────────────
# Fixes AI Search returning nothing/wrong results for Vietnamese queries typed
# without diacritics (e.g. "he thong kiem soat" instead of "hệ thống kiểm
# soát") -- jina_v3/bge_gemma2 were trained almost entirely on text WITH
# diacritics, so diacritics-less Vietnamese tokenizes into something close to
# gibberish for them, regardless of how "smart" the model otherwise is. This
# restores the diacritics on the QUERY text before it gets embedded.
#
# Model: yammdd/vietnamese-error-correction (LoRA adapter over
# vinai/bartpho-syllable, MIT license, ~0.4B params). v9.5: like jina_v3/
# bge_gemma2, this is now downloaded ON DEMAND into models/vi-diacritics/
# next to the exe (via the "Install online..." button in the Update DB
# window, using the same safe snapshot_download-into-a-folder mechanism
# already used for the embedding models), NOT baked into the exe at build
# time. An earlier version baked it into _internal, but bartpho-syllable's
# weights alone are ~1.5-2GB (it inherits mBART's huge multilingual
# vocabulary/embedding table) even after stripping the redundant TF/Flax/
# ONNX copies, which made the "small download" idea (~500MB) unrealistic to
# bake in -- moving it to the same on-demand download flow as the other
# models keeps the base exe small again and treats this as what it really
# is: an optional download, not core infrastructure.
#
# If it's not installed, diacritics restoration is just silently skipped
# (AI Search still works, just without this enhancement) -- it's not
# treated as a fatal error.
_vi_dia_pipe = None
_vi_dia_ready = False
_vi_dia_load_attempted = False

VI_DIA_BASE_REPO = "vinai/bartpho-syllable"
VI_DIA_ADAPTER_REPO = "yammdd/vietnamese-error-correction"
# Files we never need at runtime regardless of which of the two repos above
# they came from -- see scripts/download_vi_diacritics_model.py, which
# applies the same filtering when actually downloading these.
VI_DIA_IGNORE_PATTERNS = [
    "*.h5", "tf_model*", "*.msgpack", "flax_model*", "*.onnx", "*.ot",
    "*.md", "*.png", "*.jpg",
]

def _vi_dia_local_dir():
    """Local folder these get downloaded into -- same models/ root as
    jina_v3/bge_gemma2, NOT inside the exe/_internal."""
    return os.path.join(_MODEL_ROOT_CANDIDATES[0], "vi-diacritics")

def _vi_dia_installed():
    """v9.13.1 fix: search recursively for the marker files instead of a
    single hardcoded path. Different huggingface_hub versions lay out
    local_dir downloads differently (flat vs. nested inside a hashed
    snapshot subfolder) — a hardcoded exact path made this report "not
    installed" even right after a successful install, depending on which
    huggingface_hub version happened to get installed in the build."""
    d = _vi_dia_local_dir()
    def _has_file_named(root, filename):
        if not os.path.isdir(root):
            return False
        for r, _, files in os.walk(root):
            if filename in files:
                return True
        return False
    return (_has_file_named(os.path.join(d, "base"), "config.json") and
            _has_file_named(os.path.join(d, "adapter"), "adapter_config.json"))

def _vi_dia_resolve_dir(root, marker_filename):
    """Return the folder that actually CONTAINS marker_filename under root
    (searching recursively — see _vi_dia_installed for why), or `root`
    itself if not found (from_pretrained will then raise its own clear
    error rather than us silently guessing wrong)."""
    if os.path.isfile(os.path.join(root, marker_filename)):
        return root
    for r, _, files in os.walk(root):
        if marker_filename in files:
            return r
    return root

def _load_vi_diacritics_model():
    """Lazily load the Vietnamese diacritics-restoration model from its
    local models/vi-diacritics/ folder. Safe to call repeatedly -- only
    does real work once. Returns True if the model is ready to use, False
    if unavailable (not installed / failed to load), in which case callers
    should just skip the restoration step rather than error out -- this is
    a nice-to-have enhancement, not a hard requirement for AI Search to
    function.

    Loaded explicitly via PEFT (base model + LoRA adapter as two separate
    local folders) rather than relying on transformers' automatic
    repo-ID-based PEFT resolution -- that auto-resolution only reliably
    works when loading by repo ID against a proper HF hub cache, not from
    two arbitrary local folders."""
    global _vi_dia_pipe, _vi_dia_ready, _vi_dia_load_attempted
    if _vi_dia_ready:
        return True
    if _vi_dia_load_attempted:
        return False   # already tried and failed this run -- don't retry every keystroke
    _vi_dia_load_attempted = True
    try:
        if not _vi_dia_installed():
            print(f"[VI-diacritics] Not installed (see {_vi_dia_local_dir()}) — "
                  f"restoration disabled. Install it from the Update DB window "
                  f"if you want this. AI Search still works normally otherwise.")
            return False
        from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline
        from peft import PeftModel
        base_dir = _vi_dia_resolve_dir(os.path.join(_vi_dia_local_dir(), "base"), "config.json")
        adapter_dir = _vi_dia_resolve_dir(os.path.join(_vi_dia_local_dir(), "adapter"), "adapter_config.json")
        tok = AutoTokenizer.from_pretrained(adapter_dir)
        base_model = AutoModelForSeq2SeqLM.from_pretrained(base_dir)
        mdl = PeftModel.from_pretrained(base_model, adapter_dir)
        _vi_dia_pipe = pipeline("text2text-generation", model=mdl, tokenizer=tok)
        _vi_dia_ready = True
        print("[VI-diacritics] Model loaded OK.")
        return True
    except Exception as _e:
        print(f"[VI-diacritics] Load failed ({_e}) — restoration disabled, "
              f"AI Search still works normally otherwise.")
        return False

# Matches any character Vietnamese diacritics actually use (tone marks +
# accented vowels + đ/Đ). If a query already contains any of these, it's
# either already fully accented or not Vietnamese at all -- either way,
# running it through the restoration model is more likely to make it worse
# than better, so we only trigger restoration when NONE of these appear.
_VI_DIACRITIC_CHARS = re.compile(
    "[àáảãạăằắẳẵặâầấẩẫậèéẻẽẹêềếểễệìíỉĩịòóỏõọôồốổỗộơờớởỡợ"
    "ùúủũụưừứửữựỳýỷỹỵđ"
    "ÀÁẢÃẠĂẰẮẲẴẶÂẦẤẨẪẬÈÉẺẼẸÊỀẾỂỄỆÌÍỈĨỊÒÓỎÕỌÔỒỐỔỖỘƠỜỚỞỠỢ"
    "ÙÚỦŨỤƯỪỨỬỮỰỲÝỶỸỴĐ]"
)

def _maybe_restore_diacritics(query):
    """If `query` looks like it might be Vietnamese typed without diacritics
    (no diacritic characters at all, but does contain letters), try
    restoring it via the bundled model. Always returns a usable query --
    falls back to the original text unchanged on any failure/skip, so this
    can never make search worse than not calling it at all, only better or
    neutral."""
    if not query or not query.strip():
        return query
    if _VI_DIACRITIC_CHARS.search(query):
        return query   # already has diacritics (or isn't Vietnamese) -- leave it alone
    if not re.search(r"[a-zA-Z]", query):
        return query   # no letters at all (e.g. pure numbers/symbols) -- nothing to restore
    if not _load_vi_diacritics_model():
        return query   # model unavailable -- silent no-op, not an error
    try:
        restored = _vi_dia_pipe(query, max_new_tokens=128)[0]["generated_text"]
        if restored and restored.strip():
            print(f"[VI-diacritics] '{query}' -> '{restored}'")
            return restored
    except Exception as _e:
        print(f"[VI-diacritics] Restoration failed ({_e}), using original query.")
    return query

def _find_model_dir(base_dir):
    """Walk up to 4 subdirectory levels to find config.json (HuggingFace model)."""
    if base_dir and os.path.isfile(os.path.join(base_dir, "config.json")):
        return base_dir
    if not base_dir or not os.path.isdir(base_dir):
        return base_dir
    for root, dirs, files in os.walk(base_dir):
        depth = root.replace(base_dir, "").count(os.sep)
        if depth > 4: break
        if "config.json" in files:
            return root
    return base_dir

def _semantic_table_for(model_key=None):
    """Return the embedding table name for the current (or specified) model."""
    k = model_key if model_key is not None else _sem_model_key
    return SEMANTIC_MODELS.get(k, SEMANTIC_MODELS[DEFAULT_SEMANTIC_MODEL])["table"]

def _encode_query(text, model_key):
    """Encode a single query string using the correct convention for each model."""
    if model_key == "jina_v3":
        return _sem_model.encode([text], task="retrieval.query",
                                  convert_to_numpy=True, normalize_embeddings=True)[0]
    if model_key == "bge_gemma2":
        instruction = "Given a web search query, retrieve relevant passages that answer the query."
        return _sem_model.encode([text], prompt=instruction,
                                  convert_to_numpy=True, normalize_embeddings=True)[0]
    return _sem_model.encode([text], convert_to_numpy=True, normalize_embeddings=True)[0]

def _encode_passages(texts, model_key, batch_size=32):
    """Encode a list of passages/documents using the correct convention for each model."""
    if model_key == "jina_v3":
        return _sem_model.encode(texts, task="retrieval.passage", batch_size=batch_size,
                                  convert_to_numpy=True, normalize_embeddings=True,
                                  show_progress_bar=False)
    if model_key == "bge_gemma2":
        # bge-gemma2 does not need instruction for passages, only for queries
        return _sem_model.encode(texts, batch_size=batch_size, convert_to_numpy=True,
                                  normalize_embeddings=True, show_progress_bar=False)
    return _sem_model.encode(texts, batch_size=batch_size, convert_to_numpy=True,
                              normalize_embeddings=True, show_progress_bar=False)

# v9.0: Removed the entire runtime pip-install / self-heal / self-relaunch
# system that used to live here (_bootstrap_embeddable_python,
# _is_real_python, _find_pip_python, _pip_install, _heal_into_shadow_dir,
# _spawn_healer_and_exit). AI libraries (torch/transformers/
# sentence-transformers/einops) are now baked into the exe at BUILD time
# (see requirements.txt + the GitHub Actions workflow), not downloaded at
# runtime on the user's machine. This was necessary because self-
# relaunching / spawning hidden child processes / downloading GB-scale
# files at runtime kept getting silently blocked on locked-down corporate
# machines (privilege management / EDR software), with no usable error
# message. A plain `import sentence_transformers` below now just works,
# same as importing any other bundled dependency.


def _load_semantic_model(model_key=None):
    """Load embedding model by key ('jina_v3'/'bge_gemma2'). Returns immediately
    if the correct model is already loaded. If the user switched to a different model,
    unloads the old model first to avoid holding two heavy models in RAM."""
    global _sem_model, _sem_ready, _sem_device, _sem_loaded_key, _sem_model_key
    want_key = model_key if model_key is not None else _sem_model_key
    if want_key not in SEMANTIC_MODELS:
        want_key = DEFAULT_SEMANTIC_MODEL
    if _sem_ready and _sem_loaded_key == want_key:
        return True
    try:
        # v9.0: sentence_transformers/transformers/torch/einops are baked
        # into the exe at build time (see requirements.txt, pinned to
        # transformers<5 there specifically because jina_v3's custom
        # trust_remote_code modeling file, XLMRobertaLoRA, is written
        # against the transformers v4.x tied-weights API and crashes with
        # "'XLMRobertaLoRA' object has no attribute 'all_tied_weights_keys'"
        # on transformers v5+). No install/version-check/self-heal needed
        # here anymore -- if this import ever fails, it means the exe was
        # built wrong (missing dependency at build time), not something
        # this running process can fix on its own; see the error message
        # below for what to do in that case.
        try:
            import torch
            from sentence_transformers import SentenceTransformer
            import importlib
            if SEMANTIC_MODELS.get(want_key, {}).get("trust_remote"):
                importlib.import_module("einops")
        except Exception as _imp_e:
            print(f"[Semantic] Thu vien AI bi thieu hoac loi ({_imp_e}). "
                  f"Day la loi build, khong phai loi cau hinh may nay -- "
                  f"ban can build lai exe voi requirements.txt day du roi "
                  f"phat hanh lai, khong the tu sua trong luc chay.")
            raise

        _cuda_ok = torch.cuda.is_available()
        model_info_preview = SEMANTIC_MODELS.get(want_key, {})
        if model_info_preview.get("force_cpu"):
            # This model (e.g. bge_gemma2) is too heavy for small GPUs —
            # always run it on CPU even if CUDA is available.
            _sem_device = "cpu"
            if _cuda_ok:
                print(f"[Semantic] {want_key}: forcing CPU (GPU too small for this model)")
            else:
                print("[Semantic] No CUDA GPU, using CPU")
        elif _cuda_ok:
            _sem_device = "cuda"
            print(f"[Semantic] GPU: {torch.cuda.get_device_name(0)}")
        else:
            _sem_device = "cpu"
            print("[Semantic] No CUDA GPU, using CPU")
        # jina-v3 / bge-gemma2 are heavy on CPU — still allowed, just slower
        if want_key == "bge_gemma2" and _sem_device == "cpu":
            print(f"[Semantic] Note: {want_key} always runs on CPU by design (GPU too small) — slower than GPU but stable.")

        # Free old model before loading new one (saves RAM, especially with bge-gemma2)
        if _sem_model is not None:
            try:
                del _sem_model
                _sem_model = None
                import gc; gc.collect()
                try:
                    if torch.cuda.is_available():
                        torch.cuda.empty_cache()
                except Exception:
                    pass
            except Exception:
                pass

        model_info = SEMANTIC_MODELS[want_key]
        model_base_dir = SEMANTIC_MODEL_DIRS.get(want_key, "")
        model_path = _find_model_dir(model_base_dir)
        print(f"[Semantic] Loading model ({want_key}) from: {model_path}  device={_sem_device}")
        if not model_path or not os.path.isfile(os.path.join(model_path, "config.json")):
            raise FileNotFoundError(f"config.json not found for '{want_key}' in: {model_path}")
        load_kwargs = {"device": _sem_device}
        if model_info.get("trust_remote"):
            load_kwargs["trust_remote_code"] = True
        _sem_model = SentenceTransformer(model_path, **load_kwargs)
        _sem_ready = True
        _sem_loaded_key = want_key
        _sem_model_key = want_key
        print(f"[Semantic] Model ({want_key}) loaded OK on {_sem_device.upper()}")
        return True
    except Exception as _e:
        print(f"[Semantic] Load failed: {_e}")
        _sem_ready = False
        # v3.4: a failed load (e.g. missing trust_remote_code file, OOM mid-init)
        # can leave partially-allocated GPU tensors behind even though _sem_model
        # was never assigned. Left uncleaned, this "leaks" VRAM into the NEXT
        # model's load attempt (e.g. jina_v3 fails -> bge_gemma2 then OOMs on a
        # GPU that should have had plenty of free memory). Force-clear here too.
        _sem_model = None
        _sem_loaded_key = None
        try:
            import gc; gc.collect()
            import torch
            if torch.cuda.is_available():
                torch.cuda.empty_cache()
                torch.cuda.ipc_collect()
        except Exception:
            pass
        return False
# ─────────────────────────────────────────────────────────────────────────────



def is_text_file(filepath):
    if not filepath: return False
    if os.path.splitext(filepath)[1].lower() in BINARY_EXT: return False
    try:
        with open(filepath, 'rb') as f:
            chunk = f.read(1024)
            if b'\0' in chunk: return False
        return True
    except: return False

# ── OCR (v9.13) ──────────────────────────────────────────────────────────
# Optional, off by default (see the "OCR images" checkbox in Update DB).
# Images have no literal embedded text -- unlike .txt/.py (plain read) or
# .docx/.pdf (format-specific parser), a screenshot/scan needs an actual
# computer-vision model to "read" the pixels. EasyOCR is used here (not
# Tesseract) because it needs no separate native binary installed
# alongside Python -- just a pip package -- and has solid out-of-the-box
# support for Japanese + Vietnamese + English together, matching this
# app's document mix. It downloads its own model weights (~a few hundred
# MB per language) on first use.
_OCR_IMAGE_EXTS = {'.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp', '.tiff', '.tif'}
OCR_ENABLED = False   # set from the Update DB dialog's "OCR images" checkbox
_ocr_readers = {}          # {'ja_en': Reader, 'vi_en': Reader} once loaded
_ocr_load_attempted = False

def _load_ocr_readers():
    """Lazily load both EasyOCR readers. Safe to call repeatedly -- only
    does real work once. Returns True if at least one reader loaded OK."""
    global _ocr_load_attempted
    if _ocr_readers:
        return True
    if _ocr_load_attempted:
        return False
    _ocr_load_attempted = True
    try:
        import easyocr
        print("[OCR] Loading EasyOCR readers (first run downloads model weights)...")
        _ocr_readers['ja_en'] = easyocr.Reader(['ja', 'en'], gpu=False)
        _ocr_readers['vi_en'] = easyocr.Reader(['vi', 'en'], gpu=False)
        print("[OCR] EasyOCR readers loaded OK.")
        return True
    except Exception as _e:
        print(f"[OCR] Failed to load EasyOCR ({_e}) — OCR disabled for this run, "
              f"indexing continues normally for everything else.")
        return False

def _run_ocr(filepath):
    """Try both JA/EN and VI/EN readers on `filepath`, keep whichever gives
    the higher average confidence — lets one code path handle screenshots
    in either language without knowing ahead of time which one a given
    image is in. Returns "" (not an error) on any failure — a bad/corrupt
    image should never stop the rest of indexing."""
    if not _load_ocr_readers():
        return ""
    best_text, best_conf = "", -1.0
    for reader in _ocr_readers.values():
        try:
            results = reader.readtext(filepath, detail=1)
            if not results:
                continue
            avg_conf = sum(r[2] for r in results) / len(results)
            if avg_conf > best_conf:
                best_conf = avg_conf
                best_text = " ".join(r[1] for r in results)
        except Exception as _e:
            print(f"[OCR] reader failed on {filepath}: {_e}")
    return best_text

def get_file_icon(filepath, is_folder=False):
    if is_folder: return "📁 "
    if not filepath: return "📎 "
    try:
        ext = os.path.splitext(str(filepath))[1].lower()
        if ext == ".pdf": return "📕 "
        elif ext in [".doc", ".docx"]: return "📝 "
        elif ext in [".xls", ".xlsx", ".csv"]: return "📊 "
        elif ext in [".ppt", ".pptx"]: return "📉 "
        elif ext == ".one": return "🔮 "
        elif ext in [".txt", ".log", ".ini", ".json", ".xml"]: return "📄 "
    except: pass
    return "📎 "

# ─────────────────────────────────────────────────────────────────────────
# v4.9: REAL Windows shell icons (same icons Explorer/Everything show) ──
# Extracted once per file-extension via the Windows Shell API and cached
# as Tk PhotoImage objects, so the (slow-ish) icon lookup only happens once
# per unique extension for the whole app session, not once per row.
#
# Requires on Windows:  pip install pywin32 pillow
# If either package is missing, or we're not running on Windows, every
# lookup silently returns None and callers fall back to the old emoji
# icons -- the app keeps working either way, it just looks plainer.
# ─────────────────────────────────────────────────────────────────────────
_ICON_PHOTO_CACHE = {}   # key: ext (".pdf") or "__folder__" / "__file__" -> ImageTk.PhotoImage
_ICON_BACKEND_OK = False
_com_local = threading.local()   # tracks whether CoInitializeEx has run on THIS thread
if sys.platform == "win32":
    try:
        import ctypes
        from ctypes import wintypes
        import win32gui
        import win32ui
        import win32con
        from PIL import Image, ImageTk
        _ICON_BACKEND_OK = True

        class _SHFILEINFOW(ctypes.Structure):
            _fields_ = [
                ("hIcon", wintypes.HANDLE),
                ("iIcon", ctypes.c_int),
                ("dwAttributes", wintypes.DWORD),
                ("szDisplayName", ctypes.c_wchar * 260),
                ("szTypeName", ctypes.c_wchar * 80),
            ]

        _SHGFI_ICON = 0x000000100
        _SHGFI_SMALLICON = 0x000000001
        _SHGFI_USEFILEATTRIBUTES = 0x000000010
        _FILE_ATTRIBUTE_NORMAL = 0x80
        _FILE_ATTRIBUTE_DIRECTORY = 0x10

        class _BITMAPINFOHEADER(ctypes.Structure):
            _fields_ = [
                ("biSize", wintypes.DWORD), ("biWidth", ctypes.c_long),
                ("biHeight", ctypes.c_long), ("biPlanes", wintypes.WORD),
                ("biBitCount", wintypes.WORD), ("biCompression", wintypes.DWORD),
                ("biSizeImage", wintypes.DWORD), ("biXPelsPerMeter", ctypes.c_long),
                ("biYPelsPerMeter", ctypes.c_long), ("biClrUsed", wintypes.DWORD),
                ("biClrImportant", wintypes.DWORD),
            ]
        _DI_NORMAL = 0x0003

        def _hicon_to_photoimage(hicon, size=16):
            """Convert a Win32 HICON handle to a Tk PhotoImage (real RGBA alpha).
            v5.2 fix: the old version used win32ui's high-level DrawIcon() into a
            device-dependent bitmap (CreateCompatibleBitmap). Two bugs came from
            that: (1) a DDB has no real per-pixel alpha channel, so reading it
            back as RGBA gave alpha=0 almost everywhere -- icons rendered flat
            and washed-out ("no color"). (2) DrawIcon() paints the icon at its
            *native* resolution with no stretching, so on any DPI scale where
            the system small-icon isn't exactly `size` px, the icon came out
            cropped/misaligned ("wrong size" vs Explorer/Everything).
            Fix: draw into a real 32-bit top-down DIB section via DrawIconEx,
            which both stretches to the exact requested size AND preserves the
            icon's real per-pixel alpha (needed for colored folder/app icons)."""
            try:
                hdc_screen = win32gui.GetDC(0)
                hdc_mem = win32gui.CreateCompatibleDC(hdc_screen)
                bmi = _BITMAPINFOHEADER()
                bmi.biSize = ctypes.sizeof(_BITMAPINFOHEADER)
                bmi.biWidth = size
                bmi.biHeight = -size          # negative = top-down (matches PIL row order)
                bmi.biPlanes = 1
                bmi.biBitCount = 32
                bmi.biCompression = 0         # BI_RGB
                ppv_bits = ctypes.c_void_p()
                hbmp = ctypes.windll.gdi32.CreateDIBSection(
                    hdc_mem, ctypes.byref(bmi), 0, ctypes.byref(ppv_bits), None, 0
                )
                if not hbmp or not ppv_bits.value:
                    win32gui.DeleteDC(hdc_mem); win32gui.ReleaseDC(0, hdc_screen)
                    return None
                old_obj = win32gui.SelectObject(hdc_mem, hbmp)
                # DI_NORMAL draws mask+color honoring per-pixel alpha on 32-bit
                # icons, and (unlike DrawIcon) stretches to exactly size x size.
                ctypes.windll.user32.DrawIconEx(
                    hdc_mem, 0, 0, hicon, size, size, 0, None, _DI_NORMAL
                )
                buf = ctypes.string_at(ppv_bits, size * size * 4)
                img = Image.frombuffer("RGBA", (size, size), buf, "raw", "BGRA", 0, 1)
                win32gui.SelectObject(hdc_mem, old_obj)
                win32gui.DeleteObject(hbmp)
                win32gui.DeleteDC(hdc_mem)
                win32gui.ReleaseDC(0, hdc_screen)
                return ImageTk.PhotoImage(img)
            except Exception:
                return None

        def _ensure_com_initialized():
            """SHGetFileInfoW drives the Shell's icon machinery, which needs
            COM initialized (CoInitialize) on whichever thread calls it. The
            Tk main thread usually already has COM up from Tkinter/pywin32
            startup -- but real-time search inserts rows from background
            worker threads that never call CoInitialize, so the FIRST icon
            lookup from each new thread silently failed. The previous fix
            "only extract icons on the main thread" was treating that
            symptom, not the cause -- it made background-thread rows never
            get real icons at all (only Advanced/AI Search, which happen to
            run on the main thread, showed them). The actual fix: initialize
            COM once per thread, lazily, right before that thread's first
            Shell icon call -- then let every thread extract icons freely.
            Safe to call repeatedly: it's a no-op after the first successful
            call on a given thread."""
            if getattr(_com_local, "initialized", False):
                return
            try:
                # COINIT_APARTMENTTHREADED (0x2) -- Shell icon APIs are
                # documented as STA-only. RPC_E_CHANGED_MODE / S_FALSE here
                # just mean this thread already has COM up in some form
                # (e.g. the main thread, via Tkinter/pywin32) -- harmless,
                # the Shell call still works either way.
                ctypes.windll.ole32.CoInitializeEx(None, 0x2)
            except Exception:
                pass
            _com_local.initialized = True

        def _extract_shell_icon(path, is_folder, size=16):
            """Ask the Windows Shell for the icon it would show in Explorer for
            this file/folder, using the real path when it exists on disk so we
            get the exact registered app icon (Word/Excel/PDF reader/etc.)."""
            shfi = _SHFILEINFOW()
            flags = _SHGFI_ICON | _SHGFI_SMALLICON
            use_path = path if (path and os.path.exists(path)) else None
            if use_path:
                target = use_path
                attr = 0
            else:
                # File no longer on disk (offline/cloud/deleted) -- fall back
                # to a generic lookup purely from the extension/attributes.
                flags |= _SHGFI_USEFILEATTRIBUTES
                attr = _FILE_ATTRIBUTE_DIRECTORY if is_folder else _FILE_ATTRIBUTE_NORMAL
                target = path if path else ("folder" if is_folder else "file.txt")
            res = ctypes.windll.shell32.SHGetFileInfoW(
                target, attr, ctypes.byref(shfi), ctypes.sizeof(shfi), flags
            )
            if not res or not shfi.hIcon:
                return None
            photo = _hicon_to_photoimage(shfi.hIcon, size=size)
            win32gui.DestroyIcon(shfi.hIcon)
            return photo
    except Exception as _icon_backend_err:
        _ICON_BACKEND_OK = False
        print(f"[icon] Real Windows icons NOT available -- falling back to emoji icons. Reason: {_icon_backend_err}")
        print("[icon] Fix: run  pip install pywin32 pillow  then restart the app.")
else:
    print("[icon] Real Windows icons only work on Windows -- falling back to emoji icons.")



def get_tree_icon_image(filepath, is_folder=False, size=16):
    """Return a cached ImageTk.PhotoImage with the REAL Windows icon for this
    file/folder (same icon Explorer/Everything display), or None if the
    Shell icon backend isn't available (non-Windows, or pywin32/Pillow not
    installed) -- callers should fall back to get_file_icon() emoji text."""
    if not _ICON_BACKEND_OK:
        return None
    if is_folder:
        key = "__folder__"
    else:
        key = os.path.splitext(str(filepath))[1].lower() or "__file__"
    if key in _ICON_PHOTO_CACHE:
        return _ICON_PHOTO_CACHE[key]
    # v5.6 fix: previously restricted to the main thread only, as a
    # workaround for a failure that was actually caused by missing
    # per-thread COM init (see _ensure_com_initialized() above), not by
    # which thread was calling. That restriction was the reason icons never
    # showed up in the File Name / Folder Name / File Content tabs on first
    # load (real-time search fills those from background threads) while
    # Advanced/AI Search -- which run on the main thread -- looked fine.
    # Now every thread initializes its own COM apartment once, lazily, and
    # is free to extract icons.
    if _ICON_BACKEND_OK:
        _ensure_com_initialized()
    try:
        photo = _extract_shell_icon(filepath, is_folder, size=size)
    except Exception as _e:
        photo = None
        if not getattr(get_tree_icon_image, "_warned", False):
            get_tree_icon_image._warned = True
            print(f"[icon] Icon extraction failed for '{filepath}': {_e}")
    if photo is None:
        if not getattr(get_tree_icon_image, "_warned_none", False):
            get_tree_icon_image._warned_none = True
            print(f"[icon] SHGetFileInfoW returned no icon for '{filepath}' (is_folder={is_folder}) -- check the file/folder actually exists on disk.")
        return None  # do NOT cache the miss -- retry next call
    _ICON_PHOTO_CACHE[key] = photo
    return photo

def format_size(bytes_val):
    """Fix 0KB display bug — convert bytes to human-readable unit accurately."""
    if bytes_val is None: return ""
    try:
        b_val = float(bytes_val)
    except:
        return ""
    if b_val <= 0: return "0 B"
    for unit in ['B', 'KB', 'MB', 'GB']:
        if b_val < 1024.0:
            return f"{b_val:.1f} {unit}".replace(".0 ", " ")
        b_val /= 1024.0
    return f"{b_val:.1f} TB"

def get_live_size(filepath, db_fallback=0):
    """Get the real file size directly from disk.
    Falls back to the DB value if the file is offline/cloud/deleted."""
    try:
        if filepath and os.path.isfile(filepath):
            real = os.path.getsize(filepath)
            return real  # Return actual size, including genuine 0-byte files
    except Exception:
        pass
    return db_fallback

def get_file_type(filepath):
    """Return file extension as type label, e.g. 'PDF', 'DOCX', 'Folder'"""
    if not filepath: return ""
    if os.path.isdir(filepath): return "Folder"
    ext = os.path.splitext(str(filepath))[1]
    return ext.lstrip(".").upper() if ext else ""

def get_live_mtime(filepath):
    """Get file modification time as formatted string. Returns '' on failure."""
    try:
        if filepath and os.path.exists(filepath):
            ts = os.path.getmtime(filepath)
            return datetime.fromtimestamp(ts).strftime("%Y-%m-%d %H:%M")
    except Exception:
        pass
    return ""

def parse_size_filter(query_str):
    """Fix: strip extra whitespace that would leave the SQL LIKE condition empty."""
    pattern = r'(?:size\s*)?(>=|<=|>|<|=)\s*([0-9.]+)\s*(B|KB|MB|GB)?'
    match = re.search(pattern, query_str, re.IGNORECASE)
    if not match:
        return query_str.strip(), None, None
    
    op, val_str, unit = match.groups()
    try:
        val = float(val_str)
    except:
        return query_str.strip(), None, None
        
    unit = (unit or 'B').upper()
    multiplier = {'B': 1, 'KB': 1024, 'MB': 1024**2, 'GB': 1024**3}.get(unit, 1)
    bytes_val = int(val * multiplier)
    
    # Clean the text after stripping the operator, collapse all extra whitespace
    cleaned_query = re.sub(pattern, '', query_str, flags=re.IGNORECASE).strip()
    cleaned_query = " ".join(cleaned_query.split())
    
    return cleaned_query, op, bytes_val

# ── CJK-aware keyword anchor helpers ─────────────────────────────────────────
# The search logic requires an "anchor" keyword of len(k) >= 3 before it will
# run a filename/content query, to avoid noise from short English words/stop
# words ('a', 'of', ...). But CJK text has no spaces, so a whole query like
# "解析" (2 characters, a complete/specific word meaning "analysis") becomes
# ONE token of length 2 and used to get silently rejected — while "解析条件"
# (4 chars) passed and returned results. Unlike English, 1-2 CJK characters
# are often already a complete, specific, meaningful word, so we treat any
# token containing CJK characters as an anchor regardless of length.
def _contains_cjk(s):
    for ch in s:
        cp = ord(ch)
        if (0x3040 <= cp <= 0x30FF   # Hiragana + Katakana
                or 0x3400 <= cp <= 0x9FFF   # CJK Unified Ideographs (+ Ext A)
                or 0xF900 <= cp <= 0xFAFF   # CJK Compatibility Ideographs
                or 0xAC00 <= cp <= 0xD7A3   # Hangul syllables
                or 0xFF00 <= cp <= 0xFFEF): # Fullwidth forms
            return True
    return False

def _norm_txt(s):
    """NFC-normalize text before comparing/substring-matching it.

    Fixes the Name Filter silently matching nothing for filenames containing
    Vietnamese/Japanese diacritics or accented characters: two strings that
    *look* identical on screen can still be different byte sequences --
    e.g. "a" + combining acute accent (NFD, decomposed) vs the single
    precomposed "á" character (NFC). Depending on the Vietnamese IME/input
    method used when the file was created vs. when the filter text is typed,
    the DB path and the typed filter text can end up in different forms, so
    `needle in haystack` silently fails even though a human reading both
    strings sees an exact match. Normalizing both sides to NFC before
    comparing makes the match reliable regardless of which form each side
    happened to arrive in.
    """
    try:
        return unicodedata.normalize('NFC', s)
    except Exception:
        return s

def _is_anchor_kw(k, min_len=3):
    """True if keyword k is specific enough to safely drive a search:
    either it meets the normal min_len for Latin-script text, or it
    contains any CJK character (in which case length doesn't matter)."""
    return len(k) >= min_len or _contains_cjk(k)
# ──────────────────────────────────────────────────────────────────────────

def setup_context_menu(target_widget, entry_widget, search_cmd):
    def show_popup_menu(event):
        entry_widget.focus_set()
        m = tk.Menu(target_widget, tearoff=0)
        m.add_command(label="Cut", command=lambda: entry_widget.event_generate("<<Cut>>"))
        m.add_command(label="Copy", command=lambda: entry_widget.event_generate("<<Copy>>"))
        m.add_command(label="Paste", command=lambda: entry_widget.event_generate("<<Paste>>"))
        m.add_separator()
        m.add_command(label="Search", command=search_cmd)
        m.tk_popup(event.x_root, event.y_root)
    target_widget.bind("<Button-3>", show_popup_menu)

def add_only_copy_menu(widget):
    m = tk.Menu(widget, tearoff=0)
    m.add_command(label="Copy", command=lambda: widget.event_generate("<<Copy>>"))
    widget.bind("<Button-3>", lambda e: m.post(e.x_root, e.y_root))

def add_tooltip(widget, text):
    """Small English tooltip shown on hover, e.g. for toolbar buttons.
    Flips to appear below the widget instead of above it when there isn't
    enough room above (e.g. the searchbox is docked near the top of the
    screen), so the tooltip is never clipped off-screen.
    `text` may be a plain string, or a zero-arg callable returning a string
    (evaluated fresh every time the tooltip is shown) -- used e.g. for the
    ramp-light status dot whose tooltip text depends on the current color."""
    tip = {"win": None}
    def show(_e=None):
        if tip["win"] is not None:
            return
        _text = text() if callable(text) else text
        w = tk.Toplevel(widget)
        w.overrideredirect(True)
        w.attributes("-topmost", True)
        lbl = tk.Label(w, text=_text, font=("Segoe UI", 8), bg="#ffffe0", fg="#111111",
                        relief="solid", bd=1, padx=5, pady=2)
        lbl.pack()
        w.update_idletasks()
        tip_h = w.winfo_reqheight(); tip_w = w.winfo_reqwidth()
        wx = widget.winfo_rootx(); wy = widget.winfo_rooty()
        ww = widget.winfo_width(); wh = widget.winfo_height()
        above_y = wy - tip_h - 4
        # Not enough room above (tooltip would go off the top of the screen,
        # or above the widget's containing monitor edge) -> show below instead.
        if above_y < 0:
            y = wy + wh + 4
        else:
            y = above_y
        x = wx + ww // 2 - tip_w // 2
        screen_w = widget.winfo_screenwidth()
        x = max(0, min(x, screen_w - tip_w))
        w.geometry(f"+{x}+{y}")
        tip["win"] = w
    def hide(_e=None):
        if tip["win"] is not None:
            tip["win"].destroy()
            tip["win"] = None
    widget.bind("<Enter>", show, add="+")
    widget.bind("<Leave>", hide, add="+")
    widget.bind("<ButtonPress>", hide, add="+")

class HistoryPanel:
    """Search History — used to be a separate popup opened via the small 'H'
    button below the ramp light. That button is gone now; this same UI is
    built directly into the right pane of the Help tab instead, so it's
    always visible (rebuilt fresh, like the other tabs, each time results
    are shown — so it's always up to date)."""
    def __init__(self, container, parent_app):
        self.parent_app = parent_app; self.last_selected = None; self.last_time = 0; self.edit_entry = None
        f = tk.Frame(container, bg=BG_COLOR); f.pack(fill="x", padx=8, pady=8)
        now = datetime.now()
        self.yr_v = tk.StringVar(value=str(now.year)); self.mo_v = tk.StringVar(value=now.strftime("%m")); self.da_v = tk.StringVar(value=now.strftime("%d"))
        ttk.OptionMenu(f, self.yr_v, self.yr_v.get(), *[str(y) for y in range(now.year, now.year-3, -1)]).pack(side="left")
        ttk.OptionMenu(f, self.mo_v, self.mo_v.get(), *[f"{m:02d}" for m in range(1, 13)]).pack(side="left", padx=2)
        ttk.OptionMenu(f, self.da_v, self.da_v.get(), *[f"{d:02d}" for d in range(1, 32)]).pack(side="left")
        tk.Button(f, text="Filter", command=self.refresh, bg="#444", fg="white").pack(side="left", padx=5)
        tk.Button(f, text="Show All", command=self.show_all, bg="#5c5c5c", fg="white").pack(side="left", padx=2)
        tk.Button(f, text="Export Excel", command=self.export_to_excel, bg="#2e7d32", fg="white").pack(side="left", padx=5)
        tk.Button(f, text="Clear All", command=self.clear_all, bg="#c62828", fg="white").pack(side="right", padx=5)
        t_f = tk.Frame(container, bg=BG_COLOR); t_f.pack(fill="both", expand=True, padx=5, pady=5)
        self.tree = ttk.Treeview(t_f, columns=("d", "q"), show="headings"); self.tree.heading("d", text="Date Time"); self.tree.heading("q", text="Search History")
        self.tree.column("d", width=140, stretch=False); self.tree.column("q", width=300, stretch=True)
        sb = ttk.Scrollbar(t_f, orient="vertical", command=self.tree.yview); self.tree.configure(yscrollcommand=sb.set); self.tree.pack(side="left", fill="both", expand=True); sb.pack(side="right", fill="y")
        self.tree.bind("<Button-1>", self.on_click); self.tree.bind("<Button-3>", self.show_menu); self.refresh()

    def on_click(self, e):
        now = time.time(); row = self.tree.identify_row(e.y); col = self.tree.identify_column(e.x)
        if row and col == "#2" and row == self.last_selected and (now - self.last_time) > 0.4: self.show_edit_box(row, col)
        self.last_selected = row; self.last_time = now
    def show_edit_box(self, row, col):
        if self.edit_entry: self.edit_entry.destroy()
        bbox = self.tree.bbox(row, col); x, y, w, h = bbox; val = self.tree.item(row, "values")[1]
        self.edit_entry = tk.Entry(self.tree, font=("Segoe UI", 9), bd=0); self.edit_entry.insert(0, val); self.edit_entry.place(x=x, y=y, width=w, height=h); self.edit_entry.focus_set()
        add_only_copy_menu(self.edit_entry); self.edit_entry.bind("<FocusOut>", lambda e: self.edit_entry.destroy())
    def show_all(self):
        for i in self.tree.get_children(): self.tree.delete(i)
        # v7.10 fix: Search History now lives in its own HISTORY_DB_FILE,
        # completely separate from search_data.db (DB_FILE) -- so opening/
        # reading history can never affect the indexing ramp light. Also
        # guarded against sqlite3.connect() silently creating an empty file.
        if not os.path.exists(HISTORY_DB_FILE): return
        try:
            conn = sqlite3.connect(HISTORY_DB_FILE); c = conn.cursor(); c.execute("SELECT date, query FROM history ORDER BY id DESC")
            [self.tree.insert("", tk.END, values=r) for r in c.fetchall()]; conn.close()
        except: pass
    def export_to_excel(self):
        data = [self.tree.item(i)['values'] for i in self.tree.get_children()]
        if data:
            f_path = filedialog.asksaveasfilename(defaultextension=".xlsx", initialfile="SearchHistory.xlsx")
            if f_path: pd.DataFrame(data, columns=["Date Time", "Keyword"]).to_excel(f_path, index=False)
    def show_menu(self, e):
        row = self.tree.identify_row(e.y)
        if row: val = self.tree.item(row, "values")[1]; m = tk.Menu(self.tree, tearoff=0); m.add_command(label="Search again", command=lambda: self.parent_app.use_query_from_hist(val)); m.post(e.x_root, e.y_root)
    def refresh(self):
        for i in self.tree.get_children(): self.tree.delete(i)
        # v7.10 fix: reads from HISTORY_DB_FILE now, not search_data.db --
        # this is the call that runs automatically at app startup
        # (__init__ calls self.refresh() before the user ever opens the
        # Help tab), so it was the main thing silently creating a phantom
        # search_data.db before this fix.
        if not os.path.exists(HISTORY_DB_FILE): return
        try:
            conn = sqlite3.connect(HISTORY_DB_FILE); c = conn.cursor(); dt = f"{self.yr_v.get()}-{self.mo_v.get()}-{self.da_v.get()}"
            c.execute("SELECT date, query FROM history WHERE date LIKE ? ORDER BY id DESC", (f"{dt}%",)); [self.tree.insert("", tk.END, values=r) for r in c.fetchall()]; conn.close()
        except: pass
    def clear_all(self):
        if not os.path.exists(HISTORY_DB_FILE): return  # v7.10 fix: nothing to clear, don't create a phantom DB
        if messagebox.askyesno("Confirm", "Clear all history?"):
            conn = sqlite3.connect(HISTORY_DB_FILE); c = conn.cursor(); c.execute("DELETE FROM history"); conn.commit(); conn.close(); self.refresh()

class RealtimeSmartSearchApp:
    def __init__(self, root):
        self.root = root; self.root.overrideredirect(True); self.root.attributes("-topmost", True)
        # v2.9: don't stay pinned above every other window forever -- release
        # "always on top" shortly after launch (same pattern already used by
        # HistoryWindow below). This still pops the search box to the front the
        # instant it appears, but once the user opens a file/folder/browser
        # from a result, that new window is free to sit above this one instead
        # of the search box permanently covering it.
        self.root.after(500, lambda: self.root.attributes("-topmost", False))
        self.root.configure(bg=BG_COLOR)
        self.args = sys.argv[1:]; self.has_args = len(self.args) > 0
        self.ph_list = [READY_PH, READY_PH1, READY_PH2]; self.ph_index = 0
        self.x_pos = (root.winfo_screenwidth() // 4) - 120
        self.root.geometry(f"{SMALL_SIZE}+{self.x_pos}+5")
        
        self.active_result_win = None
        self.search_timer = None 
        self.current_search_id = 0
        # v5.8: auto-save to Search History after 10s of no typing/Enter --
        # see on_key_release / _maybe_save_hist_idle for details.
        self._hist_idle_timer = None
        self._last_saved_hist_query = ""
        self.db_conn = None   # persistent connection — opened once, reused every keystroke
        # v2.6: guards window creation against races. Japanese/CJK IME input fires
        # several KeyRelease events in a burst (once per romaji key during
        # composition, then again on conversion/commit) — each one used to spawn
        # its own _mft_scan_search thread with no coordination, so more than one
        # thread could see "no result window yet" before the first one's
        # scheduled show_results() actually ran, creating two duplicate windows.
        self._win_create_lock = threading.Lock()
        self._opening_result_win = False
        # Smart size filter
        self.size_op_var   = tk.StringVar(value=">")
        self.size_num_var  = tk.StringVar(value="")
        self.size_unit_var = tk.StringVar(value="MB")
        # v7.10: "Whole word" toggle -- ON by default (v7.10b: switched from
        # OFF after testing showed it filters out buried-substring false
        # positives well, e.g. "adas" no longer matches "readasync.xml" or
        # "ReadAStringExample.mlx"). When ON, a keyword like "adas" only
        # counts as a match when it's bounded by non-letter characters
        # (start/end of name, or a digit/underscore/hyphen/dot/space/
        # parenthesis) on both sides -- so "ADAS_systems.pdf" or
        # "VDIM_0ADAS_ESP" still match, but names with the keyword buried
        # mid-word don't. Trade-off: a name like "ADASDemoInstructions.pdf"
        # (no separator between "ADAS" and "Demo") also won't match while
        # ON, since both sides are letters.
        self.whole_word_var = tk.BooleanVar(value=True)
        self._all_files_data  = []   # cache for client-side filter (File Name tab)
        self._all_content_data = []  # cache for client-side filter (File Content tab)
        self.filter_count_label = None
        self.content_filter_count_label = None
        self._last_query = ""
        self._last_bm25_cont_res = []   # BM25-only cont_res cache for AI merge
        self._last_bm25_file_res = []   # BM25-only file_res cache for toggle
        self._ai_search_btn = None      # reference to AI Search button
        self._ai_mode_active = False    # True = currently showing Hybrid results
        self._ai_cont_res  = []         # AI merged content results cache
        self._ai_file_res  = []         # AI merged file results cache
        # Track split-pane trees per tab so filters can update all panes
        # Each dict: {"main": tree_widget, "adv": tree_widget, "ai": tree_widget}
        self._c_pane_trees   = {}
        self._f_pane_trees   = {}
        self._fol_pane_trees = {}
        self._adv_search_btn = None      # reference to Advanced button
        self._update_db_btn  = None      # reference to Update DB button (result window only)
        self._update_db_running = False  # True while indexing_worker() is running (button or --update data)
        self._last_index_status_text = None  # v9.11: last real progress text (e.g. "AI 2/2: 28%"),
                                              # survives the button widget being recreated on minimize/reopen
        self._mft_file_res   = []        # v2.3: live MFT scan results (files)
        self._mft_folder_res = []        # v2.3: live MFT scan results (folders)
        self._mft_render_pending_f   = False  # v2.4: coalesced re-render flag (File Name tab)
        self._mft_render_pending_fol = False  # v2.4: coalesced re-render flag (Folder Name tab)
        # v7.7 FIX: race between the DB-backed search (_smart_search_realtime,
        # scans the FULL indexed corpus) and the live disk scan (_mft_scan_search,
        # only walks the user's home folder + non-C: drives -- a much narrower
        # scope). Whichever finished LAST used to blindly overwrite tree_f/
        # tree_fol, so a single generic keyword could flash a full DB result
        # set and then immediately get stomped down to the live scan's much
        # smaller subset the moment its (slower) os.walk finished. This flag
        # records which search id the DB results were last painted for, so
        # the live-scan renderer below can MERGE with them instead of wiping
        # them out.
        self._db_rendered_sid = -1
        self._db_rendered_file_res = []   # snapshot of DB file_res (files+folders) for merge
        self._adv_mode_active = False   # True = currently showing full (unfiltered) results
        self._adv_page = 0              # 0=realtime only, 1=all results shown
        self._adv_all_cont  = []        # full content result set from Advanced search
        self._adv_all_files = []        # full file result set from Advanced search
        self._adv_split_win = None      # Advanced split window (top/bottom)
        self._ai_split_win  = None      # AI Search split window (left/right)
        # Extension filter
        self.ext_filter_var = tk.StringVar(value="")
        # Content tab filter vars (separate from File Name tab)
        self.c_size_op_var   = tk.StringVar(value=">")
        self.c_size_num_var  = tk.StringVar(value="")
        self.c_size_unit_var = tk.StringVar(value="MB")
        self.c_ext_filter_var = tk.StringVar(value="")
        # Name filter vars (substring filter on filename) — both tabs
        self.name_filter_var   = tk.StringVar(value="")   # File Name tab
        self.c_name_filter_var = tk.StringVar(value="")   # File Content tab
        # AI model selector: key in SEMANTIC_MODELS ("jina_v3"/"bge_gemma2")
        self.ai_model_var = tk.StringVar(value=_sem_model_key)

        # v2.8: this top bar (self.bg_f) is now the ONE persistent widget for the
        # search box across every state -- idle, expanded/typing, and full results.
        # Previously, the moment the first result arrived, a brand-new Toplevel
        # window was created (with its own icon/entry/history/close row) while this
        # window got hidden -- visually a "popup swap" that felt like a stutter.
        # Now show_results() just resizes THIS window and adds a results frame
        # below; this bar itself is never destroyed/recreated.
        # v4.9: the outer bar itself has NO border (icon + ramp light + close
        # button live here too) -- the border belongs only to the Entry below,
        # see the Entry creation further down.
        self.bg_f = tk.Frame(self.root, bg=BG_COLOR, height=35)
        self.bg_f.pack(fill="x", side="top")
        # v2.8.1: lock the bar's height so it can NEVER grow, no matter what gets
        # packed into it later (the AI Search / Advanced / Update DB buttons are
        # visually a bit taller than the entry due to their border+padding, and
        # without this they were quietly stretching the whole bar taller the
        # moment results opened — the "pop" the user was seeing).
        self.bg_f.pack_propagate(False)
        self._draw_search_icon(self.bg_f, size=16).pack(side="left", padx=(6, 3), pady=5)

        # v5.9: r_p sits flush against the Entry with no gap (padx below) so
        # the Entry's border reads as reaching all the way to the ramp dot
        # instead of visibly stopping short with dead space before it. r_p
        # itself stays borderless (a border around the ramp dot alone looked
        # like a separate boxed-in element, which wasn't wanted).
        r_p = self._r_p = tk.Frame(self.bg_f, bg=BG_COLOR)
        r_p.pack(side="right", fill="y", padx=(1, 5), pady=5)
        # v3.4: the small "H" (Search History) button that used to live below
        # the ramp light is gone -- Search History now lives permanently in
        # the "Help" tab instead. With only the ramp light left in this
        # column, it no longer needs to hug the top -- center it vertically
        # and make it bigger since it's now the only thing here.
        self.status_label = tk.Label(r_p, text="●", fg="#444", bg=BG_COLOR, font=("Arial", 13))
        self.status_label.pack(expand=True)
        # Ramp-light tooltip: Red/Yellow -> DB not fully indexed yet, Green -> ready.
        # While an Update DB run is actively in progress, show "Updating DB"
        # regardless of the ramp's current color (it can already read Green/
        # Yellow mid-run depending on stage) instead of the misleading
        # "Need Update DB". Callable so the text is re-evaluated fresh every
        # time the mouse hovers.
        def _ramp_tooltip_text():
            if getattr(self, "_update_db_running", False):
                return "Updating DB..."
            _fg = self.status_label.cget("fg")
            if _fg == "#2196f3":
                return "Data + AI Updated"
            if _fg == "#4caf50":
                return "Data Updated (AI pending)"
            return "Need Update DB"
        add_tooltip(self.status_label, _ramp_tooltip_text)

        # "+" button -- manually add individual files to the index, without
        # needing a full Update DB rescan. Useful for one-off files outside
        # the normal scanned folders, or files you want searchable right
        # away instead of waiting for the next scheduled Update DB run.
        self.add_file_btn = tk.Button(self.bg_f, text="+", font=("Segoe UI", 11, "bold"),
                                       bg=BG_COLOR, fg="#7ec8e3", bd=0,
                                       activebackground=BG_COLOR, cursor="hand2",
                                       command=lambda: self._add_files_dialog())
        self.add_file_btn.pack(side="right", padx=(0, 4), pady=5, before=r_p)
        add_tooltip(self.add_file_btn, "Add file(s) to the index manually")

        # "✕" close button -- only shown once results are being displayed (there's
        # nothing to close in idle mode). Re-packed with before=r_p each time it's
        # shown so it always lands as the outermost-right widget on the bar.
        self.close_btn = tk.Button(self.bg_f, text="✕", font=("Arial", 9), bg=BG_COLOR, fg="#888",
                                    bd=0, activebackground=BG_COLOR)

        self.entry_var = tk.StringVar()
        # v7.3: the border used to be drawn via the Entry's own
        # highlightthickness ring -- at SMALL_SIZE (85px total window) that
        # ring's right-hand column was getting clipped/not painted (a Tk/
        # Win32 rendering quirk that only shows up once the Entry's own
        # allotted width gets very small; it was fine at LARGE_SIZE where
        # there's slack width). Drawing the border as an actual Frame
        # background color, with the Entry inset 1px inside it, means the
        # border is just a normal widget background -- it can't be clipped
        # the way an overlay ring can, at any window size.
        self.entry_border = tk.Frame(self.bg_f, bg="#8a8a8a")
        self.entry_border.pack(side="left", fill="both", expand=True, padx=(5, 0), pady=5)
        self.entry = tk.Entry(self.entry_border, textvariable=self.entry_var, font=("Segoe UI", 10),
                               bg=ENTRY_BG, fg=TEXT_COLOR, insertbackground=TEXT_COLOR,
                               bd=0, highlightthickness=0)
        self.entry.pack(fill="both", expand=True, padx=1, pady=1)
        self.placeholder = tk.Label(self.entry, text=WAIT_PH, fg=PLACE_COLOR, bg=ENTRY_BG, font=("Segoe UI", 9, "italic"))
        if not self.has_args: self.placeholder.place(x=2, y=2)
        
        setup_context_menu(self.entry, self.entry, self.handle_action)
        setup_context_menu(self.placeholder, self.entry, self.handle_action)
        # v2.9 fix: the placeholder Label sits ON TOP of the Entry (via .place())
        # whenever the box is empty, so a left-click there was hitting the Label
        # instead of the Entry underneath -- the box looked clickable but nothing
        # got focus, so typing did nothing until the user closed (✕) and reopened.
        # Forward left-clicks (and click-drag-select) on the placeholder straight
        # to the real entry so it always gets focus + the caret.
        def _focus_entry_from_placeholder(e=None):
            self.entry.focus_set()
        self.placeholder.bind("<Button-1>", _focus_entry_from_placeholder)
        self.entry.bind("<FocusIn>", self.on_expand); self.entry.bind("<FocusOut>", self.on_shrink)
        self.entry.bind("<Return>", lambda e: self.handle_action()); self.entry.bind("<Escape>", lambda e: self.root.destroy())
        self.entry_var.trace_add("write", self.toggle_placeholder)
        
        self.entry.bind("<KeyRelease>", self.on_key_release)
        self.bg_f.bind("<Button-1>", self.start_drag); self.bg_f.bind("<B1-Motion>", self.do_drag)

        # Results-mode state -- the notebook/trees/filter bars etc. built by
        # show_results() live in this frame, packed below self.bg_f only while
        # results are showing. self._results_extra_bar holds the AI Search /
        # Advanced / Update DB buttons that also live on the search bar row.
        self.results_frame = None
        self._results_extra_bar = None
        self._in_results_mode = False
        # v4.4: set True right before a "Search again" re-run from the
        # History tab so the results notebook jumps to the File Name tab
        # once results are (re)rendered, instead of silently staying on
        # whatever tab (Help) the user triggered the re-search from.
        self._force_file_tab = False

        self.root.after(10, lambda: self.root.geometry(f"{SMALL_SIZE}+{self.x_pos}+5"))
        # v5.9c: also force a repaint shortly after the very first paint --
        # previously _force_repaint only ran on later shrink events (focus
        # loss, closing results), so the initial launch frame could show the
        # same cut-off-border/missing-ramp glitch with nothing to fix it.
        self.root.after(90, self._force_repaint)
        self.root.after(100, self.start_logic)
        self.root.after(10000, self.rotate_placeholder)
        if self.has_args: self.entry_var.set(" ".join(self.args)); self.root.after(300, self.handle_action)
        else: self.entry.focus_set()

    def rotate_placeholder(self):
        if not self.entry_var.get() and self.status_label.cget("fg") == "#4caf50":
            self.ph_index = (self.ph_index + 1) % len(self.ph_list)
            self.placeholder.config(text=self.ph_list[self.ph_index])
        self.root.after(10000, self.rotate_placeholder)

    def use_query_from_hist(self, val):
        # v3.4 FIX: this used to call self.active_result_win.destroy() first to
        # "force a fresh window" -- but since v2.8, active_result_win IS
        # self.root (results are shown in the same window, not a separate
        # Toplevel anymore). Destroying it destroyed the entire app, which is
        # exactly the crash ("application has been destroyed") seen when
        # right-clicking a history row -> "Search again". handle_action()
        # already fully rebuilds the results tabs on every search (same as
        # typing a new query + Enter while results are open), so no destroy
        # step is needed at all.
        #
        # v7.4 FIX: _force_file_tab used to be the ONLY way the tab switch
        # happened, and it was only consumed inside update_or_show_results --
        # which only runs once the DB-backed _smart_search_realtime thread
        # finishes. If db_conn wasn't ready yet (or that thread was just
        # slow), the switch silently never happened and the view stayed on
        # whatever tab (Help) "Search again" was clicked from. Since this
        # function only ever runs from a right-click inside the already-open
        # results window, self.nb is guaranteed to exist here -- so switch
        # tabs directly instead of waiting on an async callback.
        self._force_file_tab = True
        try:
            self.nb.select(0)
        except Exception:
            pass
        self.entry_var.set(val)
        self.handle_action()
    def on_expand(self, e=None): 
        if not self.has_args and not self._in_results_mode: self.root.geometry(f"{LARGE_SIZE}+{self.x_pos}+5")
    def on_shrink(self, e=None):
        self.root.after(200, self._real_shrink)
    def _real_shrink(self):
        # v2.8: results mode manages its own window size now -- the entry losing
        # focus (e.g. clicking into the results tree) must NOT shrink the window
        # back down, since it's the same window the results are displayed in.
        if self._in_results_mode: return
        try:
            focused = self.root.focus_get()
        except KeyError:
            # Tkinter bug on Windows: Combobox popdown widget causes KeyError in focus_get()
            return
        if not self.has_args and focused != self.entry:
            self.root.geometry(f"{SMALL_SIZE}+{self.x_pos}+5")
            self.root.after(60, self._force_repaint)
    def _force_repaint(self):
        try:
            # v5.9e: THE confirmed actual fix (via the [RAMP DEBUG] dump) --
            # r_p's on-screen x-position was stale from a much wider layout
            # (e.g. left over from LARGE/results mode: rp_geom showed
            # x=1229 in a 115px-wide window, which is why it reported
            # unmapped -- that position is nowhere near this window at all).
            # Re-packing the PARENT (bg_f, as this function used to do) does
            # NOT force Tk to recompute a CHILD's position -- only
            # explicitly pack_forget()+pack()'ing r_p (and the Entry, same
            # risk) itself does that, recalculating against bg_f's CURRENT
            # (now narrow) width instead of whatever it was mid-results.
            try:
                self._r_p.pack_forget()
                self._r_p.pack(side="right", fill="y", padx=(1, 5), pady=5)
            except Exception:
                pass
            try:
                self.entry_border.pack_forget()
                self.entry_border.pack(side="left", fill="both", expand=True, padx=(5, 0), pady=5)
            except Exception:
                pass
            self.root.update_idletasks()
            # v5.9: previous attempts here (v5.2's alpha nudge, plus this
            # session's lift()/geometry nudges) all failed specifically for
            # the "user clicks empty Desktop" case -- the window loses OS
            # foreground focus, and Tk-side nudges apparently aren't enough
            # to make Windows actually recomposite it. Use the native Win32
            # RedrawWindow call instead: this directly tells Windows to
            # invalidate + immediately repaint the window (and all its
            # children), bypassing whatever DWM optimization was skipping
            # the repaint. Falls back to the old alpha nudge if this ever
            # fails (e.g. non-Windows, or winfo_id() unavailable).
            try:
                import ctypes
                hwnd = self.root.winfo_id()
                RDW_INVALIDATE, RDW_ERASE, RDW_UPDATENOW, RDW_ALLCHILDREN = 0x1, 0x4, 0x100, 0x80
                ok = ctypes.windll.user32.RedrawWindow(
                    hwnd, None, None,
                    RDW_INVALIDATE | RDW_ERASE | RDW_UPDATENOW | RDW_ALLCHILDREN)
                if not ok:
                    raise RuntimeError("RedrawWindow returned 0")
            except Exception:
                self.root.attributes("-alpha", 0.99)
                self.root.after(30, lambda: self.root.attributes("-alpha", 1.0))
            try:
                self._r_p.lift()
                self.status_label.lift()
                self.status_label.config(bg=self.status_label.cget("bg"))
                self.status_label.update_idletasks()
            except Exception:
                pass
        except Exception:
            pass
    def start_drag(self, e): self._offsetx = e.x; self._offsety = e.y
    def do_drag(self, e):
        # v2.8: active_result_win is now the same window as self.root (results
        # render inline instead of in a separate Toplevel), so there is no longer
        # a second window to keep in sync here -- moving root IS moving results.
        x = self.root.winfo_x() + e.x - self._offsetx; y = self.root.winfo_y() + e.y - self._offsety
        self.root.geometry(f"+{x}+{y}"); self.x_pos = x

    def toggle_placeholder(self, *args):
        if self.entry_var.get(): self.placeholder.place_forget()
        else: self.placeholder.place(x=2, y=2)

    def _open_db_conn(self):
        """Open (or reopen) the persistent read connection with performance PRAGMAs."""
        try:
            if self.db_conn:
                try: self.db_conn.close()
                except: pass
            self.db_conn = sqlite3.connect(DB_FILE, timeout=10, check_same_thread=False)
            c = self.db_conn.cursor()
            c.execute("PRAGMA journal_mode=WAL")
            c.execute("PRAGMA cache_size=-131072")   # 128 MB read cache
            c.execute("PRAGMA temp_store=MEMORY")
            c.execute("PRAGMA mmap_size=4294967296") # 4GB memory-mapped I/O
        except Exception as e:
            print(f"DB open error: {e}")
            self.db_conn = None

    def _ramp_watchdog_tick(self):
        """Safety net: re-apply _sync_ai_adv_lock() once a second for the
        life of the app. Cheap (a few cget/config calls, all guarded by
        winfo_exists) and makes the Advanced/AI Search button state
        self-healing — if any single call site that's supposed to trigger a
        re-sync ever gets missed (e.g. a future code path forgets to call
        it), the buttons still catch up within ~1s instead of staying wrong
        until the next unrelated sync happens to fire."""
        try:
            self._sync_ai_adv_lock()
        except Exception:
            pass
        try:
            self.root.after(1000, self._ramp_watchdog_tick)
        except Exception:
            pass

    def start_logic(self):
        # v3.5: 5-state lamp —
        #   Red    (#ff0000) : search_data.db doesn't exist yet — never indexed
        #   Yellow (#ffcc00) : search_data.db exists but content indexing never
        #                      finished (files table empty/missing)
        #   Yellow (blinking): an Update DB run is actively scanning/extracting
        #                      content right now (see indexing_worker/_ramp_blink_*)
        #   Green  (#4caf50) : BM25 content indexing finished — search/Advanced
        #                      are usable — but AI embeddings for the currently
        #                      selected model aren't fully built yet
        #   Blue   (#2196f3) : BM25 content AND AI embeddings for the currently
        #                      selected model are both fully up to date
        db_exists = os.path.exists(DB_FILE)
        has_data = False
        if db_exists:
            try:
                conn = sqlite3.connect(DB_FILE); c = conn.cursor()
                # Migration: add size column if old DB doesn't have it yet
                c.execute("PRAGMA table_info(files)")
                existing_cols = [row[1] for row in c.fetchall()]
                if existing_cols and 'size' not in existing_cols:
                    c.execute("ALTER TABLE files ADD COLUMN size INTEGER DEFAULT 0")
                    conn.commit()
                c.execute("SELECT count(*) FROM files")
                if c.fetchone()[0] > 0: has_data = True
                conn.close()
            except: pass

        if not db_exists:
            # Red: no database.db at all yet — user should run --update data
            self.status_label.config(fg="#ff0000")
            self.placeholder.config(text=READY_PH)
        elif not has_data:
            # Yellow: database.db exists but hasn't finished being indexed
            # (e.g. app was closed mid-update, or DB was created but never
            # populated) — MFT filename search still works, content search won't
            self.status_label.config(fg="#ffcc00")
            self.placeholder.config(text=READY_PH)
        else:
            self._open_db_conn()
            self.status_label.config(fg="#4caf50"); self.placeholder.config(text=READY_PH)
            threading.Thread(target=self._preload_semantic_model, daemon=True).start()
            # v3.5: if the currently selected AI model's embeddings already
            # cover every content_store row (e.g. a previous --update data run
            # finished AI too), jump straight to Blue instead of staying Green
            # until the next Update DB run notices.
            # v7.10 FIX: the dropdown always starts back at DEFAULT_SEMANTIC_
            # MODEL ("jina_v3") on every app launch -- it isn't persisted. So
            # if Jina's build had failed in a previous run but another model
            # (e.g. BGE-Gemma2) succeeded, every future launch kept checking
            # Jina, found it incomplete, and stayed stuck on Green/AI-Search-
            # greyed-out forever, even though a perfectly usable AI index
            # existed under BGE. Now: if the current selection isn't fully
            # indexed, check every other model and auto-switch to the first
            # one that is.
            def _check_ai_startup():
                global _sem_model_key
                _cur_key = self.ai_model_var.get() if self.ai_model_var.get() in SEMANTIC_MODELS else DEFAULT_SEMANTIC_MODEL
                if not self._ai_fully_indexed(_cur_key):
                    for _mk in SEMANTIC_MODELS.keys():
                        if _mk != _cur_key and self._ai_fully_indexed(_mk):
                            print(f"[Semantic] '{_cur_key}' isn't fully indexed but '{_mk}' is -- "
                                  f"switching the AI model selection to '{_mk}'")
                            _cur_key = _mk
                            _sem_model_key = _mk
                            self.ai_model_var.set(_mk)

                            def _update_combo_display(k=_mk):
                                try:
                                    _combo = getattr(self, "_ai_model_combo", None)
                                    if _combo and _combo.winfo_exists():
                                        _combo.set(SEMANTIC_MODELS[k]["label"])
                                except Exception:
                                    pass
                            self.root.after(0, _update_combo_display)
                            break
                if self._ai_fully_indexed(_cur_key):
                    self.root.after(0, lambda: (self.status_label.config(fg="#2196f3"), self._sync_ai_adv_lock()))
            threading.Thread(target=_check_ai_startup, daemon=True).start()
        self._sync_ai_adv_lock()
        self.root.after(1000, self._ramp_watchdog_tick)
    
    def _add_files_dialog(self):
        """'+' button handler — pick one or more files and add them to the
        index right away, without a full Update DB rescan."""
        paths = filedialog.askopenfilenames(title="Add files to index")
        if not paths:
            return
        try:
            self.add_file_btn.config(state="disabled", fg="#888888")
        except Exception:
            pass
        threading.Thread(target=self._add_files_worker, args=(list(paths),), daemon=True).start()

    def _add_files_worker(self, paths):
        """Runs off the UI thread. Inserts each file into files/
        content_store/content_index (same tables Update DB uses), and —
        best-effort only — also embeds it into any AI model table that
        already has data, so it shows up in AI Search immediately too,
        not just BM25/File Name. If a model was never built (table empty/
        missing), we skip AI embedding for it here rather than force a
        model load just for one file; it'll get picked up on the next full
        Update DB run for that model like any other file."""
        added, failed = 0, []
        try:
            conn = self.db_conn
            if conn is None:
                raise RuntimeError("Database connection not ready (indexing in progress?)")
            c = conn.cursor()
            active_sem_tables = []
            for mk, info in SEMANTIC_MODELS.items():
                try:
                    c.execute(f"SELECT COUNT(*) FROM {info['table']}")
                    if c.fetchone()[0] > 0:
                        active_sem_tables.append((mk, info["table"]))
                except Exception:
                    pass   # that model's table doesn't exist yet -- never built, skip
            for p in paths:
                try:
                    if not os.path.isfile(p):
                        failed.append(p); continue
                    size = os.path.getsize(p)
                    name = os.path.basename(p)
                    mtime = os.path.getmtime(p)
                    c.execute("INSERT OR REPLACE INTO files (type, name, path, size) VALUES (?,?,?,?)",
                              ("File", name, p, size))
                    content = self.get_file_content(p) or ""
                    c.execute("INSERT OR REPLACE INTO content_store (path, content, mtime) VALUES (?,?,?)",
                              (p, content, mtime))
                    c.execute("DELETE FROM content_index WHERE path=?", (p,))
                    if content:
                        c.execute("INSERT INTO content_index (path, content) VALUES (?,?)", (p, content))
                    conn.commit()
                    # Best-effort AI embedding into already-built models —
                    # same "filename + content" text the bulk embedder uses.
                    embed_text = (name + " " + content[:300]).replace("\n", " ")
                    for mk, table in active_sem_tables:
                        try:
                            if _load_semantic_model(mk):
                                vec = _encode_passages([embed_text], mk, batch_size=1)[0]
                                c.execute(f"INSERT OR REPLACE INTO {table} VALUES (?,?,?)",
                                          (p, embed_text, vec.astype("float32").tobytes()))
                                conn.commit()
                        except Exception as _e:
                            print(f"[Add file] embed failed for {mk} on {p}: {_e}")
                    added += 1
                except Exception as _e:
                    print(f"[Add file] failed for {p}: {_e}")
                    failed.append(p)
            # NOTE: do NOT close conn here — it's self.db_conn, the shared
            # persistent connection _smart_search_realtime() reads from on
            # every keystroke, not a connection we opened ourselves.
        except Exception as _e:
            print(f"[Add file] DB error: {_e}")

        def _done():
            try:
                self.add_file_btn.config(state="normal", fg="#7ec8e3")
            except Exception:
                pass
            msg = f"Added {added} file(s) to the index."
            if failed:
                msg += f"\n{len(failed)} failed (see console log for details)."
            messagebox.showinfo("Add files", msg)
        self.root.after(0, _done)

    def get_file_content(self, filepath):
        if not filepath: return ""
        ext = os.path.splitext(filepath)[1].lower()
        text = ""
        try:
            print(f"[Extract] {ext}  {filepath}", flush=True)
            if is_text_file(filepath):
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read(MAX_CHARS_TO_INDEX)
            elif ext == '.pdf':
                import pypdf
                with open(filepath, 'rb') as f:
                    reader = pypdf.PdfReader(f)
                    for page in reader.pages[:10]:
                        text += (page.extract_text() or "") + " "
            elif ext == '.docx':
                doc = docx.Document(filepath)
                text = " ".join([p.text for p in doc.paragraphs])
            elif ext == '.doc':
                import win32com.client
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(filepath, ReadOnly=True, Visible=False)
                text = doc.Range().Text
                doc.Close()
                word.Quit()
            elif ext in ['.xlsx', '.xls', '.csv']:
                if ext == '.xls':
                    import pandas as pd
                    df = pd.read_excel(filepath, engine='xlrd')
                    text = df.head(100).to_string(index=False)
                else:
                    df = pd.read_excel(filepath) if ext != '.csv' else pd.read_csv(filepath)
                    text = df.head(100).to_string(index=False)
            elif ext == '.pptx':
                prs = Presentation(filepath)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + " "
            elif ext == '.ppt':
                import win32com.client
                app = win32com.client.Dispatch("PowerPoint.Application")
                pres = app.Presentations.Open(filepath, ReadOnly=True, Untitled=False, WithWindow=False)
                for slide in pres.Slides:
                    for shape in slide.Shapes:
                        if hasattr(shape, "TextFrame") and shape.TextFrame.HasText:
                            text += shape.TextFrame.TextRange.Text + " "
                pres.Close()
                app.Quit()
            elif ext == '.one':
                # Direct binary extraction — works without OneNote running.
                # .one files store text as UTF-16-LE strings interspersed in binary data.
                try:
                    with open(filepath, 'rb') as f:
                        raw = f.read(MAX_CHARS_TO_INDEX * 8)
                    import re as _re
                    # Extract UTF-16-LE runs of printable chars (min 3 chars = 6 bytes)
                    utf16_chunks = _re.findall(rb'(?:[\x20-\x7e\x00][\x00]){3,}', raw)
                    parts = []
                    for chunk in utf16_chunks:
                        try:
                            s = chunk.decode('utf-16-le', errors='ignore').strip()
                            if len(s) >= 3 and not all(c in ' \t\n\r' for c in s):
                                parts.append(s)
                        except:
                            pass
                    # Also extract plain ASCII runs (page titles, tags often ASCII)
                    ascii_chunks = _re.findall(b'[\x20-\x7e]{4,}', raw)
                    for chunk in ascii_chunks:
                        try:
                            s = chunk.decode('ascii', errors='ignore').strip()
                            if len(s) >= 4:
                                parts.append(s)
                        except:
                            pass
                    text = " ".join(parts)
                except:
                    pass
            elif ext in _OCR_IMAGE_EXTS:
                if OCR_ENABLED:
                    text = _run_ocr(filepath)
            elif ext == '.msg':
                try:
                    import win32com.client
                    outlook = win32com.client.Dispatch("Outlook.Application")
                    ns = outlook.GetNamespace("MAPI")
                    msg_obj = ns.OpenSharedItem(os.path.abspath(filepath))
                    text = f"{msg_obj.Subject or ''} {msg_obj.Body or ''} {msg_obj.SenderName or ''} {msg_obj.SenderEmailAddress or ''}"
                    msg_obj.Close(0)
                except:
                    try:
                        # Fallback: read raw bytes and extract printable strings
                        with open(filepath, 'rb') as f:
                            raw = f.read(MAX_CHARS_TO_INDEX * 4)
                        import re as _re
                        # Extract UTF-16-LE strings (common in .msg) — catches full-width chars like ：
                        utf16_strings = []
                        for s in _re.findall(b'(?:[\x20-\x7e\x00-\xff]\x00){4,}', raw):
                            try:
                                decoded = s.decode('utf-16-le', errors='ignore').strip()
                                if decoded: utf16_strings.append(decoded)
                            except: pass
                        # Extract UTF-8 chunks — catches CJK, full-width punctuation ：＜＞ etc.
                        utf8_strings = []
                        try:
                            utf8_text = raw.decode('utf-8', errors='ignore')
                            import re as _re2
                            for chunk in _re2.findall(r'[\x20-\x7e\u3000-\u9fff\uff00-\uffef\u4e00-\u9fff]{3,}', utf8_text):
                                utf8_strings.append(chunk)
                        except: pass
                        ascii_strings = [s.decode('ascii', errors='ignore')
                                         for s in _re.findall(b'[\x20-\x7e]{4,}', raw)]
                        text = " ".join(utf16_strings + utf8_strings + ascii_strings)
                    except:
                        pass
        except: pass
        return text.strip()[:MAX_CHARS_TO_INDEX]

    def _preload_semantic_model(self):
        """Load AI model in background at startup — model ready immediately when needed."""
        table = _semantic_table_for()
        try:
            conn = sqlite3.connect(DB_FILE, timeout=5)
            c = conn.cursor()
            c.execute("SELECT count(*) FROM sqlite_master WHERE type='table' AND name=?", (table,))
            has_table = c.fetchone()[0] > 0
            count = 0
            if has_table:
                c.execute(f"SELECT count(*) FROM {table}")
                count = c.fetchone()[0]
            conn.close()
        except Exception as _e:
            print(f"[Semantic] Startup check failed: {_e}")
            return
        if count == 0:
            print(f"[Semantic] No {table} data — run --update data to build (model: {_sem_model_key})")
            return
        print(f"[Semantic] Loading model at startup ({count} indexed docs, {_sem_model_key})...")
        if _load_semantic_model():
            print("[Semantic] Model ready!")
            # Update hybrid status label in UI — just mark AI as ready, not active
            def _update_lbl():
                if hasattr(self, '_hybrid_status_lbl'):
                    try:
                        self._hybrid_status_lbl.config(
                            text="📊 BM25  (AI ready)", fg="#7ec8e3")
                    except Exception: pass
            self.root.after(0, _update_lbl)
        else:
            print("[Semantic] Model load FAILED — check path and packages")

    def _ai_fully_indexed(self, model_key=None):
        """True when the given model's (or, if omitted, the currently
        selected model's) embedding table covers every row currently in
        content_store — i.e. AI Search data is completely up to date, not
        just BM25/content. Used to decide whether the ramp light should read
        Green (BM25 ready, AI pending) or Blue (BM25 + AI both ready). Safe
        to call from a background thread — opens its own short-lived
        connection rather than touching self.db_conn."""
        try:
            conn = sqlite3.connect(DB_FILE, timeout=5)
            c = conn.cursor()
            c.execute("SELECT count(*) FROM content_store")
            total = c.fetchone()[0]
            if total == 0:
                conn.close()
                return False
            table = _semantic_table_for(model_key)
            c.execute("SELECT count(*) FROM sqlite_master WHERE type='table' AND name=?", (table,))
            if not c.fetchone()[0]:
                conn.close()
                return False
            c.execute(f"SELECT count(*) FROM {table}")
            sem_count = c.fetchone()[0]
            conn.close()
            return sem_count >= total
        except Exception:
            return False

    _RAMP_BLINK_COLORS = {
        "yellow": ("#ffcc00", "#6b5200"),  # BM25/content indexing stage
        "green":  ("#4caf50", "#1f4623"),  # AI/semantic embedding stage
    }

    def _ramp_blink_start(self, color="yellow"):
        """Start blinking the ramp-light while an Update DB run is actively
        working. color="yellow" for the BM25/content stage, color="green"
        for the AI/semantic-embedding stage that follows it — lets you tell
        the two stages apart at a glance instead of both looking the same.
        Call _ramp_blink_stop() once all work finishes — success snaps to a
        solid Green/Blue (see indexing_worker), failure snaps to solid Red.

        Calling this again with a DIFFERENT color while already blinking
        (e.g. switching from the yellow BM25 stage straight into the green
        AI stage) just re-colors the ongoing blink in place, without
        stopping/restarting the loop — this is what keeps the light
        continuously blinking across the stage transition instead of
        pausing solid in between.

        Idempotent by design: this gets called from more than one place for
        the same run (the button/command handler AND indexing_worker's own
        startup), and without the guard below each call would spin up its
        own independent self-rescheduling after() loop — two loops toggling
        the same shared _ramp_blink_on flag in close succession effectively
        cancel each other out, which is why the dot used to look static
        instead of blinking."""
        self._ramp_blink_on_color, self._ramp_blink_off_color = \
            self._RAMP_BLINK_COLORS.get(color, self._RAMP_BLINK_COLORS["yellow"])
        if getattr(self, "_ramp_blinking", False):
            return  # already blinking -- color above takes effect on the next tick
        self._ramp_blinking = True
        self._ramp_blink_on = True
        self._ramp_blink_tick()

    def _ramp_blink_tick(self):
        if not getattr(self, "_ramp_blinking", False):
            return
        try:
            on_c = getattr(self, "_ramp_blink_on_color", "#ffcc00")
            off_c = getattr(self, "_ramp_blink_off_color", "#6b5200")
            self.status_label.config(fg=on_c if self._ramp_blink_on else off_c)
        except Exception:
            pass
        self._ramp_blink_on = not self._ramp_blink_on
        try:
            self.root.after(500, self._ramp_blink_tick)
        except Exception:
            pass

    def _ramp_blink_stop(self, final_fg=None):
        """Stop blinking. If final_fg is given, snap the ramp to that solid
        color; otherwise leave whatever color the last tick painted."""
        self._ramp_blinking = False
        if final_fg is not None:
            try:
                self.status_label.config(fg=final_fg)
            except Exception:
                pass

    def _begin_update_ramp(self, text):
        """Called once at the very start of indexing_worker (from the main
        thread via root.after) — sets the Update DB button/status text and
        starts the Yellow blink that runs through the whole BM25 content
        stage."""
        self._set_index_status(text, "#ffcc00")
        self._ramp_blink_start()

    def _sync_ai_adv_lock(self):
        """Grey-out (disable) the Advanced button while the ramp light is Red
        or Yellow (DB content not fully indexed yet); re-enabled the moment
        it turns Green or Blue.

        The AI Search button and the Jina/BGE model dropdown are stricter
        still: they require the ramp to be solid Blue (BM25 content AND the
        currently selected AI model's embeddings both fully up to date).
        Green is NOT enough for AI Search — Green can mean BM25 is ready but
        AI embeddings are still being built in the background (e.g. "AI
        2/2: 20%..."), and AI Search/the model dropdown must stay greyed out
        through all of that, only unlocking once the ramp actually reaches
        Blue. Safe to call any time -- widgets may not exist yet (idle mode,
        before results are shown) so every lookup is guarded."""
        try:
            _fg = self.status_label.cget("fg")
            is_ready = _fg in ("#4caf50", "#2196f3")  # Green or Blue = BM25 content ready
            is_blue = _fg == "#2196f3"                # Blue = BM25 + AI both fully ready
        except Exception:
            is_ready = False
            is_blue = False
        updating = bool(getattr(self, "_update_db_running", False))
        adv_state = "normal" if is_ready else "disabled"
        ai_ready = is_blue and not updating
        ai_state = "normal" if ai_ready else "disabled"
        combo_state = "readonly" if ai_ready else "disabled"
        try:
            _adv = getattr(self, "_adv_search_btn", None)
            if _adv and _adv.winfo_exists():
                _adv.config(state=adv_state)
        except Exception:
            pass
        try:
            _ai = getattr(self, "_ai_search_btn", None)
            if _ai and _ai.winfo_exists():
                _ai.config(state=ai_state)
        except Exception:
            pass
        try:
            _combo = getattr(self, "_ai_model_combo", None)
            if _combo and _combo.winfo_exists():
                _combo.config(state=combo_state)
        except Exception:
            pass

    def _set_index_status(self, text, fg, done=False, error=False):
        """Central place indexing_worker() reports progress through — updates
        the small status dot on the root searchbox AND the Update DB button
        (if the Results window happens to be open right now), so whichever
        one the user is currently looking at shows live progress.
        done/error also resets self._update_db_running and re-enables the button.

        The ramp dot itself only ever shows the "●" glyph and changes color
        (grey/red/yellow/green) — it never swaps to the progress text (that
        would duplicate what the Update DB button/tooltip already say and
        made the dot visually jump between a dot and a sentence). The
        Update DB button, on the other hand, does show the live text.

        v9.11 fix: also remember `text` in self._last_index_status_text.
        The search box (and its Update DB button) gets torn down and
        recreated when minimized/reopened — without remembering the last
        real progress string here, the freshly-recreated button had no way
        to know we were mid-run at "AI 2/2: 28%" and fell back to a bare
        "Updating..." with no percentage, discarding progress info that
        was still perfectly valid, just not visible anymore."""
        if not done and not error:
            self._last_index_status_text = text
        try:
            self.status_label.config(text="●", fg=fg)
        except Exception:
            pass
        self._sync_ai_adv_lock()
        if done or error:
            self._update_db_running = False
            self._last_index_status_text = None
        try:
            btn = self._update_db_btn
            if btn and btn.winfo_exists():
                if done:
                    btn.config(state="normal", text="Update DB", fg="#33363c", bg="#e6e8ec")
                elif error:
                    btn.config(state="normal", text="Update DB (error)", fg="#ff5555", bg="#3a1a1a")
                else:
                    btn.config(state="disabled", text=text, fg="#ffcc00", bg="#2a2a1a")
        except Exception:
            pass

    def indexing_worker(self, selected_tiers=None, selected_models=None, ocr_enabled=False):
        """selected_tiers: None = all tiers (1-4, default/backward-compatible).
        Otherwise a set of 0-indexed tier numbers (0=Tier1 ... 3=Tier4) — only
        files in these tiers get their CONTENT extracted/embedded. Filename
        metadata (files table, used by filename search) is still built for
        every file regardless of tier filter — the tier filter only limits
        the (slow) content-extraction step.
        selected_models: None = build AI embeddings for every model in
        SEMANTIC_MODELS (default/backward-compatible). An empty list/set =
        build NO AI embeddings this run (Tier-only update — BM25/content
        only, AI stage entirely skipped). Otherwise a list/set of model_key
        strings — only these get embeddings built this run, skipping the
        others entirely (saves a lot of time if the user only cares about
        one model, e.g. just Jina-v3).
        ocr_enabled: (v9.13) False by default — when True, image files
        (.jpg/.png/...) also get their content extracted via OCR
        (EasyOCR), same as Office/PDF/text files. Off by default because
        OCR is noticeably slower than the other extraction methods and
        downloads its own model weights on first use — an explicit opt-in
        via the Update DB dialog's "OCR images" checkbox."""
        global OCR_ENABLED
        OCR_ENABLED = bool(ocr_enabled)
        try:
            self._update_db_running = True
            self.root.after(0, lambda: self._begin_update_ramp("Updating..."))
            # Close any existing persistent connection before rebuilding DB
            if self.db_conn:
                try: self.db_conn.close()
                except: pass
                self.db_conn = None
            conn = sqlite3.connect(DB_FILE)
            c = conn.cursor()
            c.execute("CREATE TABLE IF NOT EXISTS history (id INTEGER PRIMARY KEY, query TEXT, date TEXT)")
            c.execute("CREATE TABLE IF NOT EXISTS files (type TEXT, name TEXT, path TEXT, size INTEGER)")
            c.execute("CREATE VIRTUAL TABLE IF NOT EXISTS content_index USING fts5(path, content, tokenize='trigram case_sensitive 0')")
            c.execute("CREATE TABLE IF NOT EXISTS content_store (path TEXT PRIMARY KEY, content TEXT, mtime REAL)")
            # v1.5: checkpoint table — tracks resume state across interrupted runs.
            # stage: 'scan' | 'content' | 'embed'
            # model_key: NULL for scan/content stages, model name for embed stage
            # last_path: last successfully processed path (resume point)
            c.execute("""CREATE TABLE IF NOT EXISTS update_checkpoint (
                            stage     TEXT,
                            model_key TEXT,
                            last_path TEXT,
                            done_count INTEGER,
                            total_count INTEGER,
                            updated_at TEXT,
                            PRIMARY KEY (stage, model_key)
                         )""")
            conn.commit()

            def _save_checkpoint(stage, model_key, last_path, done_count, total_count):
                try:
                    c.execute("""INSERT INTO update_checkpoint
                                 (stage, model_key, last_path, done_count, total_count, updated_at)
                                 VALUES (?,?,?,?,?,?)
                                 ON CONFLICT(stage, model_key) DO UPDATE SET
                                 last_path=excluded.last_path, done_count=excluded.done_count,
                                 total_count=excluded.total_count, updated_at=excluded.updated_at""",
                              (stage, model_key or '', last_path, done_count, total_count,
                               datetime.now().isoformat()))
                    conn.commit()
                except Exception:
                    pass

            def _get_checkpoint(stage, model_key):
                try:
                    c.execute("SELECT last_path, done_count, total_count FROM update_checkpoint WHERE stage=? AND model_key=?",
                              (stage, model_key or ''))
                    row = c.fetchone()
                    return row if row else (None, 0, 0)
                except Exception:
                    return (None, 0, 0)

            def _clear_checkpoint(stage, model_key=None):
                try:
                    if model_key is None:
                        c.execute("DELETE FROM update_checkpoint WHERE stage=?", (stage,))
                    else:
                        c.execute("DELETE FROM update_checkpoint WHERE stage=? AND model_key=?", (stage, model_key))
                    conn.commit()
                except Exception:
                    pass

            # ── STAGE 1: File scan ────────────────────────────────────────────
            # File scan is fast (minutes) — always rerun fully to catch new/deleted/moved
            # files. Not the bottleneck, so no resume needed here.
            c.execute("DROP TABLE IF EXISTS files")
            c.execute("CREATE TABLE files (type TEXT, name TEXT, path TEXT, size INTEGER)")
            c.execute("DROP TABLE IF EXISTS files_temp")
            c.execute("CREATE TABLE files_temp (type TEXT, name TEXT, path TEXT, size INTEGER)")
            # content_index (FTS5) + content_store are content caches — do NOT drop them.
            # We diff against existing mtime to skip unchanged files (huge time save on resume).
            c.execute("CREATE VIRTUAL TABLE IF NOT EXISTS content_index USING fts5(path, content, tokenize='trigram case_sensitive 0')")
            c.execute("CREATE TABLE IF NOT EXISTS content_store (path TEXT PRIMARY KEY, content TEXT, mtime REAL)")
            c.execute("CREATE INDEX IF NOT EXISTS idx_content_store_path ON content_store(path)")
            # v1.5: migrate old DBs (pre-checkpoint) that lack the mtime column
            c.execute("PRAGMA table_info(content_store)")
            _cs_cols = [row[1] for row in c.fetchall()]
            if 'mtime' not in _cs_cols:
                c.execute("ALTER TABLE content_store ADD COLUMN mtime REAL")
                conn.commit()
                print("[Migration] Added mtime column to content_store (old DB detected)")

            # Load existing content_store mtimes — used to skip re-extracting unchanged files
            c.execute("SELECT path, mtime FROM content_store")
            _existing_mtime = {r[0]: r[1] for r in c.fetchall()}
            # v1.5: semantic tables are now resumed/appended, never dropped here.
            # Table creation + resume logic happens later in the embedding loop below.
            
            batch_files = []; batch_content = []
            _pending_extract = []  # v2.5: (tier, path, mtime) — extracted AFTER full walk, sorted by tier
            _scan_total = 0; _scan_skipped = 0; _scan_new = 0
            _seen_paths = set()  # v1.5: deduplicate — prevents double-extract if drives overlap

            # v1.5: known text extensions — skip is_text_file() open() call for these
            # is_text_file() reads 1024 bytes from every unknown file → huge bottleneck
            # with millions of files. Whitelist covers 99% of indexable text files.
            # v2.5: trimmed compiled-language source extensions (.c/.h/.cpp/.java/.cs/
            # .go/.rs/.php/.sh) — these are almost always installer/SDK payload noise
            # (Abaqus, CATIA, etc.) in an engineering file share, not content users
            # actually search for. Add them back below if you do want them indexed.
            _TEXT_EXTS = {
                '.txt', '.md', '.rst', '.csv', '.log', '.ini', '.cfg', '.conf',
                '.json', '.xml', '.yaml', '.yml', '.toml', '.html', '.htm',
                '.py', '.js', '.ts', '.css', '.sql', '.bat', '.ps1',
                '.tsv', '.spck', '.env',
            }
            # .env / .env.local / .env.production etc. are dotfiles — Python's
            # splitext() treats the leading dot as the filename, not an extension
            # (splitext(".env") == (".env", "")), so they'd never match _TEXT_EXTS
            # by extension. Match by basename prefix instead.
            # NOTE: .env files commonly hold API keys/credentials — indexing their
            # content means those secrets become searchable/readable via the app.
            # Keep this only if that's an acceptable tradeoff in your environment.
            _ENV_FILE_PREFIX = ".env"
            # Extensions that are definitively binary — never open, never extract
            _BINARY_EXTS = {
                '.exe', '.dll', '.lib', '.obj', '.pyc', '.bin',
                '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.ico', '.svg', '.webp',
                '.zip', '.7z', '.rar', '.tar', '.gz', '.bz2',
                '.mp3', '.mp4', '.avi', '.mov', '.mkv', '.wav', '.flac',
                '.db', '.sqlite', '.ldb', '.sst',
                '.lock', '.tmp', '.bak', '.cache',
            }

            # v1.5: Whitelist-based scan — much faster than blacklisting system folders.
            # Strategy: scan C:\Users\<current_user> + all non-C drives (D, E, F...) fully.
            # This covers 99% of user documents while skipping Windows/app install folders.
            import getpass as _getpass
            _username = _getpass.getuser()
            _scan_roots = []
            # Always include current user's home folder on C:\
            _user_home = os.path.join("C:\\", "Users", _username)
            if os.path.exists(_user_home):
                _scan_roots.append(_user_home)
            # Include all non-C drives fully
            for _dl in string.ascii_uppercase:
                if _dl == "C": continue
                _dp = f"{_dl}:\\"
                if os.path.exists(_dp):
                    _scan_roots.append(_dp)
            print(f"[Scan] Roots: {_scan_roots}")

            for d in _scan_roots:
                for root_dir, dirs, files in os.walk(d):
                    def _should_skip(name):
                        nl = name.lower()
                        if nl in SKIP_FOLDERS: return True
                        if nl.startswith("python") and len(nl) <= 12: return True
                        return False
                    dirs[:] = [dd for dd in dirs
                               if not _should_skip(dd)
                               and not _looks_like_model_repo(os.path.join(root_dir, dd))]
                    for item in (dirs + files):
                        try:
                            full_p = os.path.normpath(os.path.join(root_dir, item))
                            if full_p in _seen_paths:
                                continue
                            _seen_paths.add(full_p)
                            is_dir = item in dirs
                            
                            f_size = os.path.getsize(full_p) if not is_dir else 0
                            batch_files.append(("Folder" if is_dir else "File", item, full_p, f_size))
                            
                            if not is_dir:
                                f_ext = os.path.splitext(full_p)[1].lower()
                                # Skip sensitive files (IT security / DLP)
                                f_basename = os.path.basename(full_p)
                                if _SENSITIVE_FILE_RE.search(f_basename):
                                    continue
                                is_allowed = False
                                # Priority 1: known Office/PDF — always extract
                                if f_ext in {'.pdf', '.docx', '.doc', '.xlsx', '.xls',
                                             '.csv', '.pptx', '.ppt', '.one', '.msg'}:
                                    is_allowed = True
                                # Priority 2: known text extensions — no need to open file
                                elif f_ext in _TEXT_EXTS:
                                    if f_size < 2 * 1024 * 1024:
                                        is_allowed = True
                                # Priority 3: .env / .env.local / .env.* dotfiles
                                elif f_basename.lower().startswith(_ENV_FILE_PREFIX):
                                    if f_size < 2 * 1024 * 1024:
                                        is_allowed = True
                                # Priority 4 (v9.13): images, only when OCR is
                                # enabled for this run (Update DB "OCR images"
                                # checkbox) — off by default since OCR is much
                                # slower than the other extraction methods.
                                # Capped at 15MB: legitimate screenshots/scans
                                # are almost always well under this; anything
                                # bigger is more likely a huge raw photo/scan
                                # where OCR would be slow for little benefit.
                                elif f_ext in _OCR_IMAGE_EXTS and OCR_ENABLED:
                                    if f_size < 15 * 1024 * 1024:
                                        is_allowed = True
                                # v2.5: removed the old "unknown extension -> open file
                                # and sniff for text" fallback. It was extracting a lot
                                # of installer/SDK metadata (.catnls, .clsid, .iid,
                                # .intinfo, .tmw, ...) that just happens to be plain-text
                                # formatted, plus it opened every single unrecognized
                                # file on the drive (slow). Anything not explicitly
                                # whitelisted above is now skipped outright.
                                if is_allowed:
                                    _scan_total += 1
                                    # v1.5: skip cloud-only (not-yet-downloaded) OneDrive/SharePoint files
                                    # FILE_ATTRIBUTE_RECALL_ON_DATA_ACCESS = 0x400000 means cloud placeholder
                                    try:
                                        _attrs = os.stat(full_p).st_file_attributes
                                        if _attrs & 0x400000:  # cloud-only placeholder
                                            _scan_skipped += 1
                                            continue
                                    except Exception:
                                        pass
                                    # v1.5: skip re-extracting content if file unchanged since last run
                                    try:
                                        _cur_mtime = os.path.getmtime(full_p)
                                    except Exception:
                                        _cur_mtime = None
                                    _prev_mtime = _existing_mtime.get(full_p)
                                    if _prev_mtime is not None and _cur_mtime is not None and abs(_prev_mtime - _cur_mtime) < 1.0:
                                        _scan_skipped += 1
                                        continue  # content unchanged — keep existing content_store/content_index row
                                    # v2.5: DON'T read/extract content here — just queue it
                                    # (tier, path, mtime). Content is extracted AFTER the
                                    # full walk across ALL drives finishes, sorted by tier
                                    # (see _ext_tier) — so tier-1 (office/pdf) content is
                                    # fully indexed across the WHOLE filesystem before
                                    # tier-2 starts, then tier-3, then tier-4. If the run
                                    # gets interrupted (Ctrl+C, power loss, etc.), whatever
                                    # was already committed is always the highest-value
                                    # content first, never a random alphabetical slice.
                                    # A re-run resumes naturally: unchanged files are
                                    # already skipped above by the mtime check, so only
                                    # genuinely new/changed files get queued again.
                                    _tier = 2 if (f_ext == '' and f_basename.lower().startswith(_ENV_FILE_PREFIX)) \
                                            else self._ext_tier(f_ext)
                                    # v9.13: images don't belong to any of the
                                    # 4 Tiers (they'd fall through to an
                                    # "unlisted" tier index that no Tier
                                    # checkbox can ever select), so gate them
                                    # ONLY by the OCR_ENABLED checkbox
                                    # (already checked above via is_allowed),
                                    # not by Tier selection at all.
                                    if f_ext in _OCR_IMAGE_EXTS:
                                        pass
                                    elif selected_tiers is not None and _tier not in selected_tiers:
                                        continue  # tier not requested — skip content extraction
                                    _pending_extract.append((_tier, full_p, _cur_mtime))
                            if len(batch_files) >= 2000:
                                c.executemany("INSERT INTO files_temp VALUES (?,?,?,?)", batch_files); batch_files = []
                        except: continue
            if batch_files: c.executemany("INSERT INTO files_temp VALUES (?,?,?,?)", batch_files)
            conn.commit()

            # ── v2.5: extract content in tier order ──────────────────────────
            # Tier 1: office/pdf | Tier 2: msg/txt/log/one | Tier 3: scripts/
            # config/spck | Tier 4: markup/misc. Stable sort keeps files within
            # the same tier in their original (discovery) order.
            _pending_extract.sort(key=lambda x: x[0])
            _tier_labels = {0: "Tier1 office/pdf", 1: "Tier2 msg/txt/log/one",
                            2: "Tier3 scripts/config", 3: "Tier4 markup/misc"}
            _tier_counts = {}
            for _t, _p, _m in _pending_extract:
                _tier_counts[_t] = _tier_counts.get(_t, 0) + 1
            print(f"[Scan] Tier filter: {'ALL (1-4)' if selected_tiers is None else ','.join(str(t+1) for t in sorted(selected_tiers))}")
            print(f"[Scan] {len(_pending_extract)} files queued for extraction — "
                  + ", ".join(f"{_tier_labels.get(t, f'Tier{t+1}')}: {n}" for t, n in sorted(_tier_counts.items())))

            for _tier, full_p, _cur_mtime in _pending_extract:
                try:
                    content = self.get_file_content(full_p)
                    if content:
                        batch_content.append((full_p, content, _cur_mtime))
                        _scan_new += 1
                except Exception:
                    continue
                if len(batch_content) >= 100:
                    # Remove stale rows before insert (file changed → re-index FTS + store)
                    c.executemany("DELETE FROM content_index WHERE path=?", [(bc[0],) for bc in batch_content])
                    c.executemany("INSERT INTO content_index (path, content) VALUES (?,?)",
                                  [(bc[0], bc[1]) for bc in batch_content])
                    c.executemany("INSERT OR REPLACE INTO content_store VALUES (?,?,?)", batch_content)
                    conn.commit()
                    batch_content = []
            if batch_content:
                c.executemany("DELETE FROM content_index WHERE path=?", [(bc[0],) for bc in batch_content])
                c.executemany("INSERT INTO content_index (path, content) VALUES (?,?)",
                              [(bc[0], bc[1]) for bc in batch_content])
                c.executemany("INSERT OR REPLACE INTO content_store VALUES (?,?,?)", batch_content)
                conn.commit()
            print(f"[Scan] {_scan_total} files checked, {_scan_skipped} unchanged (skipped), {_scan_new} new/changed (re-extracted)")
            c.execute("DELETE FROM files"); c.execute("INSERT INTO files SELECT * FROM files_temp")
            c.execute("DROP TABLE IF EXISTS files_temp")
            c.execute("CREATE INDEX IF NOT EXISTS idx_files_path ON files(path)")
            c.execute("CREATE INDEX IF NOT EXISTS idx_files_name ON files(name)")
            c.execute("CREATE INDEX IF NOT EXISTS idx_files_type ON files(type)")
            conn.commit()

            # ── BM25/content indexing stage complete ─────────────────────────
            # files + content_store + content_index are all rebuilt and usable
            # right now, so Search and Advanced become usable immediately even
            # though the AI embedding pass below can still take a long time.
            # If an AI embedding pass IS about to run, keep the ramp blinking
            # but switch it from Yellow to Green so it's visually obvious
            # which stage is active. If there's no AI pass this run
            # (selected_models == [] , a deliberate "Tier-only" update), stop
            # blinking and go solid Green instead. The ramp only advances to
            # Blue once AI embeddings finish too (see the end of this
            # function / _ai_fully_indexed()).
            _will_run_ai_phase = (selected_models is None) or (len(selected_models) > 0)
            if _will_run_ai_phase:
                self.root.after(0, lambda: (self._ramp_blink_start(color="green"), self._sync_ai_adv_lock()))
            else:
                self.root.after(0, lambda: (self._ramp_blink_stop(final_fg="#4caf50"), self._sync_ai_adv_lock()))

            # v1.5: track paths whose content was just re-extracted (new/changed files).
            # Their OLD embeddings (if any, from previous run) are now stale and must
            # be purged from every semantic table so they get re-embedded below.
            _changed_paths = set()
            c.execute("SELECT path, mtime FROM content_store")
            for _p, _m in c.fetchall():
                _prev = _existing_mtime.get(_p)
                if _prev is None or (_m is not None and abs(_prev - _m) >= 1.0):
                    _changed_paths.add(_p)
            if _changed_paths:
                for _mk in SEMANTIC_MODELS.keys():
                    _tbl = _semantic_table_for(_mk)
                    try:
                        c.execute(f"SELECT name FROM sqlite_master WHERE type='table' AND name=?", (_tbl,))
                        if c.fetchone():
                            c.executemany(f"DELETE FROM {_tbl} WHERE path=?", [(p,) for p in _changed_paths])
                    except Exception:
                        pass
                conn.commit()
                print(f"[Semantic] Invalidated stale embeddings for {len(_changed_paths)} changed file(s)")

            # ── Build semantic embeddings for ALL models in one pass ────────
            import numpy as np
            c2 = conn.cursor()
            # v1.5b: only pull path + first 300 chars (SUBSTR in SQL) — avoids loading
            # full file content (can be MBs per row) into RAM for hundreds of thousands
            # of rows. This was causing RAM to fill up (86%+) and trigger Windows page
            # file swapping, which is far slower than CPU-bound embedding itself.
            c2.execute("SELECT path, SUBSTR(content, 1, 300) FROM content_store")
            rows = c2.fetchall()
            total_rows = len(rows)
            rows_by_path = {r[0]: r[1] for r in rows}  # dict for O(1) lookup, replaces rows list

            all_model_keys = list(SEMANTIC_MODELS.keys())
            # None => every model (backward-compatible default). An empty
            # list/set is a deliberate "Tier-only, no AI this run" choice —
            # must be checked with `is not None`, not truthiness, or an
            # explicit empty selection would be silently treated the same
            # as None and build every model anyway.
            if selected_models is not None:
                all_model_keys = [k for k in all_model_keys if k in selected_models]
            n_models = len(all_model_keys)
            for model_idx, model_key in enumerate(all_model_keys, 1):
                try:
                    label = SEMANTIC_MODELS[model_key]["label"]
                    print(f"[Semantic] Building index for model {model_idx}/{n_models}: {model_key} ({label})")
                    self.root.after(0, lambda mi=model_idx, mn=n_models, lbl=label:
                        self._set_index_status(f"AI {mi}/{mn}...", "#4caf50"))

                    if not _load_semantic_model(model_key):
                        print(f"[Semantic] Skipping {model_key} — model load failed")
                        continue

                    sem_table = _semantic_table_for(model_key)
                    # v1.5: DO NOT drop table — resume by skipping paths already embedded.
                    # If table doesn't exist yet, create fresh (first run for this model).
                    c2.execute(f"CREATE TABLE IF NOT EXISTS {sem_table} (path TEXT PRIMARY KEY, snippet TEXT, embedding BLOB)")
                    conn.commit()

                    # Load paths already embedded for this model — these are skipped.
                    c2.execute(f"SELECT path FROM {sem_table}")
                    _already_done = set(r[0] for r in c2.fetchall())
                    # v1.5b: iterate keys directly instead of filtering a full copy of `rows`
                    remaining_paths = [p for p in rows_by_path.keys() if p not in _already_done]
                    skipped_count = total_rows - len(remaining_paths)
                    if skipped_count > 0:
                        print(f"[Semantic] {model_key}: resuming — {skipped_count}/{total_rows} already embedded, "
                              f"{len(remaining_paths)} remaining")

                    if not remaining_paths:
                        print(f"[Semantic] {model_key}: nothing to do — already complete ({total_rows} docs)")
                        continue

                    SEM_BATCH = 8 if model_key == "bge_gemma2" else 32
                    sem_buf = []
                    total_done = skipped_count

                    for i in range(0, len(remaining_paths), SEM_BATCH):
                        paths = remaining_paths[i:i+SEM_BATCH]
                        # v9.8 fix: prepend the filename to the text that gets
                        # embedded, not just the extracted file content. Before
                        # this, a file like "Switch-model.zip" — a format we
                        # never extract text from (.zip is in _BINARY_EXTS) —
                        # got embedded from an EMPTY content string, producing
                        # a near-meaningless vector totally disconnected from
                        # its obviously-relevant filename. Even for text-bearing
                        # formats (.spck, etc.) whose content is mostly numeric/
                        # config data rather than natural language, the filename
                        # is often the single strongest topical signal available
                        # and was previously being thrown away entirely.
                        clean_snippets = [
                            (os.path.basename(p) + " " + (rows_by_path[p] or "")).replace("\n", " ")
                            for p in paths
                        ]
                        vecs = _encode_passages(clean_snippets, model_key, batch_size=SEM_BATCH)
                        for path, snippet, vec in zip(paths, clean_snippets, vecs):
                            sem_buf.append((path, snippet, vec.astype("float32").tobytes()))
                        total_done += len(paths)
                        if len(sem_buf) >= 500:
                            c2.executemany(f"INSERT OR REPLACE INTO {sem_table} VALUES (?,?,?)", sem_buf)
                            conn.commit(); sem_buf = []
                            # v1.5: checkpoint — survives interrupt; next run resumes from here
                            _save_checkpoint('embed', model_key, paths[-1], total_done, total_rows)
                        if total_done % (SEM_BATCH * 10) == 0 or total_done == total_rows:
                            _pct = int(total_done * 100 / max(total_rows, 1))
                            self.root.after(0, lambda p=_pct, mi=model_idx, mn=n_models:
                                self._set_index_status(f"AI {mi}/{mn}: {p}%", "#4caf50"))
                    if sem_buf:
                        c2.executemany(f"INSERT OR REPLACE INTO {sem_table} VALUES (?,?,?)", sem_buf)
                        conn.commit()
                        _save_checkpoint('embed', model_key, paths[-1] if paths else '', total_done, total_rows)
                    print(f"[Semantic] {model_key}: indexed {total_done}/{total_rows} docs into {sem_table}")
                    _clear_checkpoint('embed', model_key)  # model fully done — checkpoint no longer needed

                except Exception as _se:
                    print(f"[Semantic] Embedding build failed for {model_key}: {_se}")
                    print(f"[Semantic] {model_key}: progress saved — next --update data run will resume from here")
            # ────────────────────────────────────────────────────────────────

            conn.commit()
            conn.close()
            import time; time.sleep(0.5)  # let WAL checkpoint flush before reopening
            # ⚠️ VACUUM removed: 44GB DB needs ~88GB free disk + 30min → causes RED LIGHT
            # v3.5: Blue only once the currently selected AI model's embeddings
            # actually cover every content_store row (they might not, e.g. if
            # _load_semantic_model() failed above) — otherwise stay Green.
            # v7.10 FIX: previously the ramp only turned Blue (and AI Search
            # unlocked) if the model happening to be selected in the dropdown
            # was fully indexed -- if that model's build failed (e.g. Jina-v3
            # here: "'XLMRobertaLoRA' object has no attribute
            # 'all_tied_weights_keys'") but a DIFFERENT model in the same run
            # succeeded (e.g. BGE-Gemma2), the ramp stayed stuck on Green and
            # AI Search stayed greyed out even though a usable AI index
            # existed. Now: if the currently selected model isn't fully
            # indexed, check the other models that were part of this run and
            # auto-switch to the first one that IS fully indexed.
            global _sem_model_key
            _cur_key = self.ai_model_var.get() if self.ai_model_var.get() in SEMANTIC_MODELS else DEFAULT_SEMANTIC_MODEL
            if not self._ai_fully_indexed(_cur_key):
                for _mk in all_model_keys:
                    if _mk != _cur_key and self._ai_fully_indexed(_mk):
                        print(f"[Semantic] '{_cur_key}' isn't fully indexed but '{_mk}' is -- "
                              f"switching the AI model selection to '{_mk}'")
                        _cur_key = _mk
                        _sem_model_key = _mk
                        self.ai_model_var.set(_mk)

                        def _update_combo_display(k=_mk):
                            try:
                                _combo = getattr(self, "_ai_model_combo", None)
                                if _combo and _combo.winfo_exists():
                                    _combo.set(SEMANTIC_MODELS[k]["label"])
                            except Exception:
                                pass
                        self.root.after(0, _update_combo_display)
                        break
            _final_fg = "#2196f3" if self._ai_fully_indexed(_cur_key) else "#4caf50"
            print(f"[Semantic] Indexing finished — final ramp color: "
                  f"{'BLUE (AI fully indexed)' if _final_fg == '#2196f3' else 'GREEN (AI not fully indexed)'}")

            def _finish(fg=_final_fg):
                # v5.3 fix: each step is now guarded independently. Before, this
                # was a single list-expression -- if any ONE call raised, the
                # remaining calls (including _set_index_status(done=True), the
                # one that actually flips the ramp to Blue) silently never ran,
                # leaving the button stuck on "AI x/x..." / Green until restart.
                try: self._open_db_conn()   # reopen persistent search connection
                except Exception as _e1: print(f"[Semantic] _open_db_conn failed: {_e1}")
                try: self._ramp_blink_stop()
                except Exception as _e2: print(f"[Semantic] _ramp_blink_stop failed: {_e2}")
                try: self._set_index_status("●", fg, done=True)
                except Exception as _e3: print(f"[Semantic] _set_index_status failed: {_e3}")
                try: self.placeholder.config(text=READY_PH)
                except Exception as _e4: print(f"[Semantic] placeholder update failed: {_e4}")
                try: self._sync_ai_adv_lock()  # belt-and-suspenders re-sync
                except Exception as _e5: print(f"[Semantic] _sync_ai_adv_lock failed: {_e5}")

            self.root.after(0, _finish)
        except Exception as e:
            import traceback
            err = traceback.format_exc()
            print(f"Indexing Error:\n{err}")
            try:
                log_path = os.path.join(os.path.dirname(os.path.abspath(DB_FILE)), "search_error.log")
                with open(log_path, "a", encoding="utf-8") as lf:
                    from datetime import datetime as _dt
                    lf.write(f"\n[{_dt.now()}] INDEXING ERROR:\n{err}\n")
            except: pass
            self.root.after(0, lambda: (self._ramp_blink_stop(), self._set_index_status("●", "#ff0000", error=True)))

    def on_key_release(self, event):
        if event.keysym in ["Up", "Down", "Return", "Escape", "Shift_L", "Shift_R", "Control_L", "Control_R", "Alt_L", "Alt_R"]: return
        q = self.entry_var.get().strip()

        # v5.8: any keystroke invalidates a pending idle-save-to-History
        # check from a previous pause — reschedule fresh below (once we know
        # this isn't an empty box / a "-" command).
        if getattr(self, '_hist_idle_timer', None):
            try: self.root.after_cancel(self._hist_idle_timer)
            except Exception: pass
            self._hist_idle_timer = None

        if len(q) < 1:
            # v2.7: user request — clearing the box to retype shouldn't shrink the
            # results window back down. Just cancel any pending debounced search and
            # leave the last results on screen; only the ✕ close button (or a special
            # command below) should shrink the window back to the small search box.
            if self.search_timer:
                self.root.after_cancel(self.search_timer)
                self.search_timer = None
            return

        _CMD_LIST = ["--update data", "--exit", "--quit"]
        ql = q.lower()

        # Cancel any pending "close-on-command" check from a previous keystroke --
        # if we're here, the user just typed/changed something, so an earlier
        # scheduled check is stale and must not fire on top of the new text.
        if getattr(self, '_cmd_close_timer', None):
            try: self.root.after_cancel(self._cmd_close_timer)
            except Exception: pass
            self._cmd_close_timer = None

        # v3.1 FIX: the box starting with "-" (e.g. "-", "--", "--h", "--u"...)
        # no longer closes the results window / clears the box on the spot --
        # that made "--" itself vanish and forced retyping "--help" in the
        # shrunken box. Instead treat it as "possibly typing a command": don't
        # search, and only AFTER the user pauses (400ms with no further
        # keystroke) do we check whether the finished text is a real command
        # and close the results window then. If the user keeps typing, this
        # check keeps getting cancelled/rescheduled above, so it never fires
        # mid-typing.
        if ql.startswith("-"):
            if self.search_timer:
                self.root.after_cancel(self.search_timer)
                self.search_timer = None
            self._cmd_close_timer = self.root.after(
                400, lambda snapshot=q: self._maybe_close_for_cmd(snapshot))
            return

        self.current_search_id += 1
        search_version = self.current_search_id
        if self.search_timer: self.root.after_cancel(self.search_timer)
        self.root.update_idletasks()
        box_x, box_y, box_h = self.root.winfo_x(), self.root.winfo_y(), self.root.winfo_height()

        # v5.8: save to Search History automatically once the user pauses on
        # this text for 10s without pressing Enter (see _maybe_save_hist_idle).
        self._hist_idle_timer = self.root.after(
            10000, lambda snapshot=q: self._maybe_save_hist_idle(snapshot))

        # v2.3: Always launch MFT scan (realtime, no DB needed) for File/Folder Name tabs
        threading.Thread(
            target=self._mft_scan_search,
            args=(q, search_version, box_x, box_y, box_h), daemon=True
        ).start()

        # v2.3: Also launch BM25 (File Content tab) if DB is available — blend results
        if self.db_conn is not None:
            self.search_timer = self.root.after(350, lambda: threading.Thread(
                target=self._smart_search_realtime,
                args=(q, search_version, box_x, box_y, box_h, False), daemon=True
            ).start())

    def _rerun_current_search(self):
        """v7.10: re-run the search currently in the box with the SAME text
        -- used when a toggle that changes matching behavior (e.g. "Whole
        word") flips, so results refresh immediately without the user having
        to retype/re-trigger the query themselves."""
        q = self.entry_var.get().strip()
        if not q or q.startswith("-"):
            return
        self.current_search_id += 1
        search_version = self.current_search_id
        if self.search_timer:
            try: self.root.after_cancel(self.search_timer)
            except Exception: pass
        self.root.update_idletasks()
        box_x, box_y, box_h = self.root.winfo_x(), self.root.winfo_y(), self.root.winfo_height()
        threading.Thread(
            target=self._mft_scan_search,
            args=(q, search_version, box_x, box_y, box_h), daemon=True
        ).start()
        if self.db_conn is not None:
            threading.Thread(
                target=self._smart_search_realtime,
                args=(q, search_version, box_x, box_y, box_h, False), daemon=True
            ).start()

    def _maybe_close_for_cmd(self, snapshot):
        """Fired ~400ms after the user stops typing a string starting with '-'.
        Only closes the results window if the box still holds exactly the same
        text (i.e. the user has genuinely paused / finished) AND that text is
        one of the recognized commands. Actually running the command (help
        popup, --update data, exit, etc.) still happens on Enter via
        handle_action() as before -- this only tidies up the results window
        so it isn't sitting open behind a command about to run."""
        self._cmd_close_timer = None
        if self.entry_var.get().strip() != snapshot:
            return  # user kept typing / edited since — stale check, ignore
        ql = snapshot.lower()
        _CMD_LIST = ["--update data", "--exit", "--quit"]
        if ql in _CMD_LIST or ql.startswith("--update data"):
            if self.active_result_win and tk.Toplevel.winfo_exists(self.active_result_win):
                self._close_results_window()
                self.entry_var.set(snapshot)  # restore the command text the user typed
                self.entry.icursor("end")

    # ── v2.3: MFT / os.walk realtime scan ────────────────────────────────────
    def _mft_scan_search(self, q, sid, box_x, box_y, box_h):
        """Realtime filename scan — no DB needed. Streams results into File Name
        and Folder Name tabs as they are found. If DB is also available, BM25
        runs in parallel and populates File Content tab independently."""
        try:
            # Open result window if needed — guarded by a lock so a burst of
            # KeyRelease events (e.g. Japanese/CJK IME composition) can't race
            # each other into creating two separate result windows.
            with self._win_create_lock:
                win_missing = not self.active_result_win or not tk.Toplevel.winfo_exists(self.active_result_win)
                already_opening = self._opening_result_win
                if win_missing and not already_opening:
                    self._opening_result_win = True

            if win_missing and not already_opening:
                try:
                    self.root.after(0, lambda: self.show_results(
                        [], [], q, False, box_x, box_y, box_h, sid))
                    time.sleep(0.15)
                finally:
                    self._opening_result_win = False
            elif win_missing and already_opening:
                # Another thread is already opening the window — wait for it
                # to appear instead of opening a second one.
                for _ in range(30):
                    if self.active_result_win and tk.Toplevel.winfo_exists(self.active_result_win):
                        break
                    time.sleep(0.02)
            else:
                # v2.3 fix: window already open from a previous query — MFT streams
                # results into it but never touched the title, so it stayed stale.
                def _set_title():
                    try:
                        if sid != self.current_search_id: return
                        if self.active_result_win and tk.Toplevel.winfo_exists(self.active_result_win):
                            self.active_result_win.title(f"Results: {q}")
                    except Exception: pass
                self.root.after(0, _set_title)

            # Clear File Name + Folder Name trees (not Content — BM25 handles that)
            def _clear():
                try:
                    if hasattr(self, 'tree_f')   and self.tree_f.winfo_exists():
                        self.tree_f.delete(*self.tree_f.get_children())
                    if hasattr(self, 'tree_fol') and self.tree_fol.winfo_exists():
                        self.tree_fol.delete(*self.tree_fol.get_children())
                except Exception: pass
            self.root.after(0, _clear)
            time.sleep(0.05)

            import getpass
            _user_home = os.path.join("C:\\", "Users", getpass.getuser())
            scan_roots = [_user_home] if os.path.exists(_user_home) else []
            for letter in string.ascii_uppercase:
                if letter == "C": continue
                dp = f"{letter}:\\"
                if os.path.exists(dp): scan_roots.append(dp)

            # AND match: each whitespace-separated keyword must appear
            # SOMEWHERE in the name (not required to be contiguous) — this is
            # what the app's own help text promises ("space = AND search").
            # We also track an OR match (any keyword, but not full AND) in
            # parallel. AND results always render first; OR-only results are
            # appended BELOW them at the end of the scan (row[5] = 0 for AND,
            # 1 for OR — see _sort_priority) — applied independently for File
            # Name and Folder Name, so one tab having enough AND matches
            # doesn't suppress the other tab's OR fallback.
            keywords = [k for k in q.lower().split() if k]
            if not keywords:
                return
            # v5.8: drop filler words ("to", "relevant", ...) before AND/OR
            # matching -- see STOPWORDS comment above for why this matters.
            keywords = _strip_stopwords(keywords)
            multi_kw = len(keywords) > 1
            ww = self.whole_word_var.get()  # v7.10: "Whole word" toggle
            file_idx    = [0]
            folder_idx  = [0]
            or_file_idx   = [0]
            or_folder_idx = [0]
            MAX = 500
            # Hard time budget: for a genuinely rare query (or on huge drives)
            # neither MAX may ever be reached, so without a cap we'd walk the
            # entire filesystem before showing anything. Bound the worst case
            # latency instead of scanning forever.
            #
            # v2.5 fix: drives are scanned sequentially (C-home, then D, E,
            # F...). A single slow/large drive (esp. a network share) could
            # eat the ENTIRE global budget, so later drives (e.g. F:) never
            # got scanned at all — even though they exist and the BM25/DB
            # index (built offline, no live scanning needed) still has them.
            # Give each root a fair per-root time slice too, so a slow drive
            # gets cut off and moves on instead of starving the rest.
            # v7.10: bumped from 3.0s -- a query whose matches sit deep inside
            # a large drive (e.g. "adas" living several folders down on a big
            # D:/E: drive full of engineering data) could get cut off by
            # os.walk's per-root time slice before ever reaching those paths,
            # while a query whose matches happen to sit shallower (or on a
            # smaller/faster drive) found plenty within the same budget. This
            # isn't about the keyword itself, just where on disk it happens to
            # live vs. how much of the tree got walked before the clock ran
            # out. 10s is still a soft cap (not exhaustive on huge drives) --
            # for guaranteed complete results regardless of location/depth,
            # run --update data once so search uses the full DB index instead
            # of this live best-effort scan.
            TIME_BUDGET_SEC = 12.0
            PER_ROOT_BUDGET_SEC = max(2.0, TIME_BUDGET_SEC / max(1, len(scan_roots)))
            t_start = time.time()
            self._mft_file_res   = []
            self._mft_folder_res = []
            or_file_res   = []   # OR-only matches, appended below AND at the end
            or_folder_res = []
            timed_out = False

            for root_dir in scan_roots:
                if sid != self.current_search_id: return
                if timed_out: break
                t_root_start = time.time()
                for root, dirs, files in os.walk(root_dir):
                    if sid != self.current_search_id: return
                    if time.time() - t_start > TIME_BUDGET_SEC:
                        timed_out = True
                        break
                    if time.time() - t_root_start > PER_ROOT_BUDGET_SEC:
                        break   # move on to the next drive, don't starve it
                    dirs[:] = [d for d in dirs
                                if d.lower() not in SKIP_FOLDERS
                                and not (d.lower().startswith("python") and len(d) <= 12)]

                    # ── Folders ───────────────────────────────────────────
                    for d in dirs:
                        dl = d.lower()
                        is_and = all(_kw_matches(kw, dl, ww) for kw in keywords)
                        if is_and and folder_idx[0] < MAX:
                            full_p = os.path.normpath(os.path.join(root, d))
                            mt = get_live_mtime(full_p)
                            folder_idx[0] += 1
                            row = ("Folder", d, full_p, 0, mt, 0)
                            self._mft_folder_res.append(row)
                            self._schedule_mft_render("fol", sid)
                        elif multi_kw and or_folder_idx[0] < MAX:
                            # v5.8: track how many keywords matched (not just
                            # any/none) so results matching MORE keywords can
                            # be ranked above single-keyword-only matches
                            # instead of sitting at the same priority.
                            m_count = sum(1 for kw in keywords if _kw_matches(kw, dl, ww))
                            if m_count > 0:
                                full_p = os.path.normpath(os.path.join(root, d))
                                mt = get_live_mtime(full_p)
                                or_folder_idx[0] += 1
                                or_folder_res.append(("Folder", d, full_p, 0, mt, 1, m_count))

                    # ── Files ─────────────────────────────────────────────
                    for f in files:
                        fl = f.lower()
                        is_and = all(_kw_matches(kw, fl, ww) for kw in keywords)
                        if is_and and file_idx[0] < MAX:
                            full_p = os.path.normpath(os.path.join(root, f))
                            try: sz = os.path.getsize(full_p)
                            except: sz = 0
                            mt = get_live_mtime(full_p)
                            file_idx[0] += 1
                            row = ("File", f, full_p, sz, mt, 0)
                            self._mft_file_res.append(row)
                            self._schedule_mft_render("f", sid)
                        elif multi_kw and or_file_idx[0] < MAX:
                            m_count = sum(1 for kw in keywords if _kw_matches(kw, fl, ww))
                            if m_count > 0:
                                full_p = os.path.normpath(os.path.join(root, f))
                                try: sz = os.path.getsize(full_p)
                                except: sz = 0
                                mt = get_live_mtime(full_p)
                                or_file_idx[0] += 1
                                or_file_res.append(("File", f, full_p, sz, mt, 1, m_count))

                    # Stop early once we've got plenty of AND results AND
                    # plenty of OR fallback data too (no point walking the
                    # whole disk just to keep collecting more of either pool).
                    #
                    # v7.4 FIX: the old condition put "not multi_kw" at the
                    # top level of the OR, e.g. (file_idx>=MAX or not multi_kw
                    # or or_file_idx>=MAX). For a single-keyword query
                    # (multi_kw=False), "not multi_kw" is always True, which
                    # made the WHOLE clause True regardless of file_idx --
                    # the scan broke out after the very FIRST os.walk
                    # directory, before ever descending into subfolders.
                    # That's why a single generic keyword (e.g. "adas") could
                    # return 0 results while a 2-keyword query (multi_kw=True,
                    # where this bug doesn't trigger) correctly walked the
                    # whole tree and found plenty. The AND-count requirement
                    # must always hold; only the OR-count requirement should
                    # be skipped when there's just one keyword.
                    if (file_idx[0] >= MAX and (or_file_idx[0] >= MAX or not multi_kw)) and \
                       (folder_idx[0] >= MAX and (or_folder_idx[0] >= MAX or not multi_kw)):
                        break

            # ── Always blend: AND results on top, OR-only results appended
            # below — independently for File Name and Folder Name tabs.
            # (Previously this only kicked in when AND was completely empty,
            # so "simpack realtime" with exactly 1 AND hit never got the OR
            # results appended below it.)
            # v5.8: OR-only rows are sorted by match_count (desc) before
            # being appended, so e.g. a file matching 2 of 3 keywords shows
            # above one matching only 1 of 3 -- then the temporary 7th
            # "match_count" field is stripped back off (rows are 6-tuples
            # everywhere else, incl. _sort_priority/_fill_files/_fill_folders).
            if multi_kw:
                or_file_res.sort(key=lambda r: r[6], reverse=True)
                or_folder_res.sort(key=lambda r: r[6], reverse=True)
                self._mft_file_res   = self._mft_file_res   + [r[:6] for r in or_file_res]
                self._mft_folder_res = self._mft_folder_res + [r[:6] for r in or_folder_res]

            # Final render once the scan finishes, so the very last batch of
            # rows (which might not have hit the coalescing window) is shown.
            self._schedule_mft_render("f", sid, force=True)
            self._schedule_mft_render("fol", sid, force=True)

        except Exception as e:
            print(f"[MFT] Error: {e}", flush=True)

    # ── v2.4: coalesced MFT render — fixes 4 related bugs at once ───────────
    # 1) stale "#" numbering from old-query callbacks still landing after a
    #    new search started (no sid guard before)
    # 2) MFT rows ignored priority sort (office/pdf/msg first, log/code last)
    # 3) horizontal scrollbar / column width reset on every single insert
    # 4) active ext/size/name filter being bypassed by raw streaming inserts
    # Instead of inserting each found row immediately, we append to the
    # result cache and schedule one debounced full re-render (sorted +
    # filtered) a little later. Multiple finds within the debounce window
    # collapse into a single re-render.
    def _schedule_mft_render(self, which, sid, force=False):
        pending_attr = "_mft_render_pending_f" if which == "f" else "_mft_render_pending_fol"
        if getattr(self, pending_attr) and not force:
            return
        setattr(self, pending_attr, True)
        delay = 0 if force else 120

        def _do():
            setattr(self, pending_attr, False)
            if sid != self.current_search_id:
                return
            if which == "f":
                self._render_mft_file_tree()
            else:
                self._render_mft_folder_tree()
        self.root.after(delay, _do)

    def _merge_mft_with_db(self, mft_rows, is_folder):
        """v7.7 FIX: merge the live MFT-scan rows with the last DB-backed
        result snapshot (_db_rendered_file_res), if that snapshot is for the
        SAME search id, instead of letting whichever thread finishes last
        (DB search vs. live disk scan) blindly overwrite the other's results.

        The DB-backed search scans the full indexed corpus and is treated as
        the authoritative/more complete set; the live scan only walks the
        user's home folder + non-C: drives (see _mft_scan_search), so it can
        legitimately have FEWER matches than the DB even though it finishes
        later. Rather than let that narrower set stomp the fuller one, we
        union both by path (DB rows first, then any MFT-only paths the DB
        doesn't know about -- e.g. very recent files not yet indexed)."""
        if self._db_rendered_sid != self.current_search_id:
            # DB hasn't painted anything for this query (yet, or at all) --
            # nothing to merge with, just show what the live scan found.
            return list(mft_rows)
        db_rows = [r for r in self._db_rendered_file_res
                   if (str(r[0]).lower() == 'folder') == is_folder]
        seen_paths = set()
        merged = []
        for r in db_rows:
            p = r[2]
            if p in seen_paths:
                continue
            seen_paths.add(p)
            # Pad DB rows (4-tuple: type,name,path,size) to the 6-tuple shape
            # MFT rows use (type,name,path,size,mtime,and_or_flag) so the
            # renderers below can treat both uniformly.
            merged.append(r if len(r) >= 6 else (r[0], r[1], r[2], r[3], None, 0))
        for r in mft_rows:
            p = r[2]
            if p in seen_paths:
                continue
            seen_paths.add(p)
            merged.append(r)
        return merged

    def _render_mft_file_tree(self):
        """Full (but cheap) re-render of tree_f from self._mft_file_res,
        merged with any DB-backed results already shown for this same query:
        priority-sorted, current filter applied, scroll/column width kept."""
        try:
            if not hasattr(self, 'tree_f') or not self.tree_f.winfo_exists():
                return
        except Exception:
            return
        tree = self.tree_f
        merged_rows = self._merge_mft_with_db(self._mft_file_res or [], is_folder=False)
        rows = self._sort_priority(merged_rows, 1)
        rows = self._filter_file_rows(rows)

        x0 = self._save_xview(tree)
        widths = self._save_col_widths(tree)
        for item in tree.get_children():
            tree.delete(item)
        for rn, item in enumerate(rows, start=1):
            try:
                path = item[2]
                img = get_tree_icon_image(path, is_folder=False)
                name_txt = item[1] if img else get_file_icon(path) + item[1]
                mt = item[4] if len(item) > 4 and item[4] is not None else get_live_mtime(path)
                tree.insert("", "end", text="", **({"image": img} if img else {}), values=(
                    name_txt, format_size(item[3]), mt,
                    get_file_type(path), os.path.dirname(path), path))
            except Exception:
                pass
        self._restore_col_widths(tree, widths)
        self._restore_xview(tree, x0)

    def _render_mft_folder_tree(self):
        """Full (but cheap) re-render of tree_fol from self._mft_folder_res,
        merged with any DB-backed results already shown for this same query:
        current filter applied, scroll/column width kept."""
        try:
            if not hasattr(self, 'tree_fol') or not self.tree_fol.winfo_exists():
                return
        except Exception:
            return
        tree = self.tree_fol
        merged_rows = self._merge_mft_with_db(self._mft_folder_res or [], is_folder=True)
        rows = self._sort_priority(merged_rows, 1)
        name_needle = _norm_txt(self.name_filter_var.get().strip().lower())
        if name_needle:
            rows = [r for r in rows if name_needle in _norm_txt(os.path.basename(r[2]).lower())]

        x0 = self._save_xview(tree)
        widths = self._save_col_widths(tree)
        for item in tree.get_children():
            tree.delete(item)
        for rn, item in enumerate(rows, start=1):
            try:
                path = item[2]
                img = get_tree_icon_image(path, is_folder=True)
                name_txt = item[1] if img else "📁 " + item[1]
                mt = item[4] if len(item) > 4 and item[4] is not None else get_live_mtime(path)
                tree.insert("", "end", text="", **({"image": img} if img else {}), values=(name_txt, mt, path, path))
            except Exception:
                pass
        self._restore_col_widths(tree, widths)
        self._restore_xview(tree, x0)

    def _filter_file_rows(self, rows):
        """Apply the File Name tab's active size/ext/name filters to a row list.
        Shared by the MFT live renderer and the manual filter re-apply so both
        paths always agree — this is what fixes the 'MFT stream ignores the
        active filter' inconsistency."""
        op, size_bytes = self._get_size_filter_bytes()
        ext_filter  = self._get_ext_filter()
        name_needle = _norm_txt(self.name_filter_var.get().strip().lower())
        if op is None and ext_filter is None and not name_needle:
            return rows
        cmp_fn = {'>':  lambda a, b: a >  b, '>=': lambda a, b: a >= b,
                  '<':  lambda a, b: a <  b, '<=': lambda a, b: a <= b,
                  '=':  lambda a, b: a == b}.get(op) if op else None
        result = []
        for item in rows:
            if cmp_fn is not None:
                if not cmp_fn(get_live_size(item[2], item[3]), size_bytes):
                    continue
            if ext_filter is not None:
                _, ext = os.path.splitext(item[2])
                if ext.lower() not in ext_filter:
                    continue
            if name_needle:
                if name_needle not in _norm_txt(os.path.basename(item[2]).lower()):
                    continue
            result.append(item)
        return result

    # ── shared helpers: preserve scroll position + column widths across a
    #    full clear+rebuild of a Treeview (used by MFT render and filters) ──
    def _save_xview(self, tree):
        try:
            return tree.xview()[0]
        except Exception:
            return 0.0

    def _restore_xview(self, tree, x0):
        if x0 and x0 > 0.0:
            try:
                tree.xview_moveto(x0)
            except Exception:
                pass

    def _save_col_widths(self, tree):
        widths = {}
        try:
            for col in tree["columns"]:
                widths[col] = tree.column(col, "width")
        except Exception:
            pass
        return widths

    def _restore_col_widths(self, tree, widths):
        for col, w in widths.items():
            try:
                tree.column(col, width=w)
            except Exception:
                pass
    # ─────────────────────────────────────────────────────────────────────────

    def _semantic_search(self, query, top_k=SEMANTIC_TOP_K, threshold=None):
        """Return list of (path, score_pct) using cosine similarity against the embedding
        table of the currently selected model. The 2 models (jina_v3/bge_gemma2) have
        different embedding dimensions/vector spaces so each uses a separate table,
        and (v9.3) each can have its own similarity threshold -- see
        SEMANTIC_THRESHOLD_BY_MODEL near the top of the file."""
        try:
            if not _load_semantic_model():
                return []
            if threshold is None:
                threshold = SEMANTIC_THRESHOLD_BY_MODEL.get(_sem_model_key, SEMANTIC_THRESHOLD_DEFAULT)
            query = _maybe_restore_diacritics(query)
            import numpy as np
            q_vec = _encode_query(query, _sem_model_key).astype("float32")
            conn = self.db_conn
            if conn is None:
                return []
            c = conn.cursor()
            table = _semantic_table_for()
            try:
                c.execute(f"SELECT path, embedding FROM {table}")
                rows = c.fetchall()
            except sqlite3.OperationalError:
                # This model's table has not been built yet (--update data not run for this model)
                rows = []
            if not rows:
                return []
            paths = [r[0] for r in rows]
            mat = np.frombuffer(b"".join(r[1] for r in rows), dtype="float32").reshape(len(rows), -1)
            scores = mat @ q_vec   # cosine similarity (vectors normalized)
            top_idx = np.argsort(scores)[::-1][:top_k]
            results = [(paths[i], float(scores[i])) for i in top_idx if scores[i] >= threshold]
            if not results and len(top_idx) > 0:
                # v9.3: everything got filtered out by `threshold`. This is
                # worth logging rather than silently returning [] -- the
                # single SEMANTIC_THRESHOLD (0.25) is shared by both models,
                # but jina_v3 and bge_gemma2 are different architectures
                # with different cosine-similarity score distributions, so
                # the same cutoff can behave very differently between them
                # (e.g. one model still clears 0.25 for a bad/ambiguous
                # query, like non-diacritic Vietnamese, while the other's
                # best score sits just under it and gets filtered to
                # nothing). Print the actual best score reached so this is
                # easy to diagnose instead of just looking "broken".
                _best = float(scores[top_idx[0]])
                print(f"[Semantic] {_sem_model_key}: 0 results after threshold filter "
                      f"(best raw score was {_best:.3f}, threshold is {threshold}). "
                      f"If this best score is consistently just under the threshold for "
                      f"this model, consider lowering SEMANTIC_THRESHOLD for it specifically.")
            return results
        except Exception as _e:
            print(f"[Semantic] Search error: {_e}")
            return []

    def _ai_search_and_update(self, query):
        """Run semantic search and merge with cached BM25 results for ALL tabs.
        Calling again when already in Hybrid mode restores BM25-only (toggle)."""
        # Strategy: BM25 results always on top (sorted by BM25 score),
        # AI-only results appended below (sorted by semantic score).
        # This respects BM25 precision while still surfacing AI extras.
        BM25_WEIGHT = 0.7
        SEM_WEIGHT  = 0.3

        def _set_btn(text, fg, bg, state="normal"):
            def _do():
                try:
                    if self._ai_search_btn and self._ai_search_btn.winfo_exists():
                        self._ai_search_btn.config(text=text, fg=fg, bg=bg, state=state)
                except Exception: pass
            self.root.after(0, _do)

        def _set_status(text, color):
            def _do():
                try:
                    if hasattr(self, '_hybrid_status_lbl') and self._hybrid_status_lbl.winfo_exists():
                        self._hybrid_status_lbl.config(text=text, fg=color)
                except Exception: pass
            self.root.after(0, _do)

        # ── TOGGLE: se dang o Hybrid thi restore BM25 ─────────────────────
        if self._ai_mode_active:
            def _restore_bm25():
                try:
                    if not self.active_result_win or not tk.Toplevel.winfo_exists(self.active_result_win):
                        return
                    bm25_cont  = self._last_bm25_cont_res
                    bm25_files = self._last_bm25_file_res
                    only_files   = [r for r in bm25_files if str(r[0]).lower() != "folder"]
                    only_folders = [r for r in bm25_files if str(r[0]).lower() == "folder"]
                    self._all_content_data = bm25_cont
                    self._all_files_data   = only_files

                    for item in self.tree_c.get_children():   self.tree_c.delete(item)
                    for item in self.tree_f.get_children():   self.tree_f.delete(item)
                    for item in self.tree_fol.get_children(): self.tree_fol.delete(item)
                    for rn, item in enumerate(bm25_cont,  start=1):
                        self._insert_row(self.tree_c,   item, rn, 0)
                    for rn, item in enumerate(only_files,   start=1):
                        self._insert_row(self.tree_f,   item, rn, 1)
                    for rn, item in enumerate(only_folders, start=1):
                        self._insert_row(self.tree_fol, item, rn, 2)
                    if self.content_filter_count_label:
                        self.content_filter_count_label.config(text=f"{len(bm25_cont)} files")
                    if self.filter_count_label:
                        self.filter_count_label.config(text=f"{len(only_files)} files")
                    self.nb.tab(0, text=" File Name ")
                    self.nb.tab(1, text=" Folder Name ")
                except Exception as _e:
                    print(f"[AI Toggle] Restore BM25 error: {_e}")

            self._ai_mode_active = False
            self.root.after(0, _restore_bm25)
            self.root.after(0, self._close_ai_split)
            _set_btn("🤖 AI Search", "#7ec8e3", "#1e3a5f")
            _set_status("📊 BM25", "#aaaaaa")
            return
        # ──────────────────────────────────────────────────────────────────

        _set_btn("⏳ Running...", "#ffcc00", "#2a2a1a", "disabled")
        _set_status("⏳ AI searching...", "#ffcc00")

        try:
            if not _load_semantic_model():
                _set_btn("🤖 AI Search", "#7ec8e3", "#1e3a5f")
                _set_status("⚠️ AI model not available", "#ff6666")
                return

            sem_pairs = self._semantic_search(query)
            if not sem_pairs:
                _set_btn("🤖 AI Search", "#7ec8e3", "#1e3a5f")
                _set_status("ℹ️ AI: no additional results", "#888888")
                return

            sem_scores = {p: s for p, s in sem_pairs}   # path → 0.0-1.0

            # ── A: Merge File Content ─────────────────────────────────────
            # Group 1: BM25 paths → on top, sorted by BM25 score (small AI boost)
            # Group 2: AI-only paths → appended below, sorted by semantic score
            bm25_cont = self._last_bm25_cont_res
            bm25_scores_c = {}
            size_map_c = {}
            for item in bm25_cont:
                path = item[0]; sz = item[1] if len(item) > 1 else 0
                size_map_c[path] = sz
                bm25_scores_c[path] = (item[2] / 99.0) if len(item) > 2 and item[2] else 0.0

            new_paths = [p for p in sem_scores if p not in size_map_c]
            if new_paths and self.db_conn:
                try:
                    ph2 = ",".join("?" * len(new_paths))
                    c2 = self.db_conn.cursor()
                    c2.execute(f"SELECT path, size FROM files WHERE path IN ({ph2})", new_paths)
                    for row in c2.fetchall():
                        size_map_c[row[0]] = row[1]
                except Exception: pass

            bm25_group = []
            for path in bm25_scores_c:
                b = bm25_scores_c[path]
                s = sem_scores.get(path, 0.0)
                score = BM25_WEIGHT * b + SEM_WEIGHT * s
                bm25_group.append((path, size_map_c.get(path, 0), score))
            bm25_group.sort(key=lambda x: x[2], reverse=True)

            ai_only_group = []
            for path in sem_scores:
                if path not in bm25_scores_c:
                    ai_only_group.append((path, size_map_c.get(path, 0), sem_scores[path]))
            ai_only_group.sort(key=lambda x: x[2], reverse=True)

            combined = bm25_group + ai_only_group
            max_c = combined[0][2] if combined else 1.0
            cont_res = [(p, sz, int((sc / max_c) * 99) if max_c > 0 else 0)
                        for p, sz, sc in combined]

            # ── B: File Name + Folder Name — AI-augmented ────────────────
            # Problem: BM25 file_res may be empty when query has special chars/short tokens
            # (e.g. "Python file 3.2_e5-base" → "file" is too short, "3.2_e5-base" has dots)
            # Solution: pull file records for ALL paths that semantic returned, then
            # merge with BM25 file_res (BM25 on top, AI-only extras below).
            bm25_files = self._last_bm25_file_res  # list of (type, name, path, size)

            # Collect AI-suggested paths from content semantic results
            sem_paths_all = list(sem_scores.keys())

            # Also look up parent folders and sibling files for semantic paths
            ai_file_extras = []
            if sem_paths_all and self.db_conn:
                try:
                    import re as _re_ai

                    # ── Format keyword → extension mapping ───────────────
                    _FORMAT_MAP = {
                        "excel": [".xlsx", ".xls", ".csv"],
                        "xlsx":  [".xlsx"],
                        "xls":   [".xls"],
                        "csv":   [".csv"],
                        "word":  [".doc", ".docx"],
                        "doc":   [".doc", ".docx"],
                        "docx":  [".docx"],
                        "pdf":   [".pdf"],
                        "ppt":   [".ppt", ".pptx"],
                        "pptx":  [".pptx"],
                        "powerpoint": [".ppt", ".pptx"],
                        "python": [".py"],
                        "py":    [".py"],
                        "text":  [".txt", ".log"],
                        "txt":   [".txt"],
                        "log":   [".log"],
                        "image": [".png", ".jpg", ".jpeg", ".bmp", ".gif"],
                        "png":   [".png"],
                        "jpg":   [".jpg", ".jpeg"],
                        "zip":   [".zip", ".7z", ".rar"],
                        "onenote": [".one"],
                        "one":   [".one"],
                        "msg":   [".msg"],
                        "outlook": [".msg"],
                        "email": [".msg"],
                    }

                    # Split query into: format keywords (→ ext filter) + name tokens
                    query_lower = query.lower()
                    raw_parts = _re_ai.split(r'[\s]+', query_lower)
                    ext_filter_set = set()   # extensions to filter by
                    name_toks = []           # tokens to LIKE-search in filenames
                    for part in raw_parts:
                        clean = _re_ai.sub(r'[\-_\.\(\)]', '', part)  # strip punctuation
                        if clean in _FORMAT_MAP:
                            ext_filter_set.update(_FORMAT_MAP[clean])
                        else:
                            # Use sub-tokens split by non-alphanum as name search tokens
                            # v5.8: drop filler words here too, otherwise
                            # e.g. "to" LIKE-matches every ".toc" file.
                            subtoks = [t for t in _re_ai.split(r'[\s\-_\.]+', part)
                                       if len(t) >= 2 and t not in STOPWORDS]
                            name_toks.extend(subtoks)

                    # Remove duplicate/empty tokens, deduplicate case-insensitively
                    seen_t = set()
                    name_toks_clean = []
                    for t in name_toks:
                        tl = t.lower()
                        if tl not in seen_t and not tl.isdigit():
                            seen_t.add(tl)
                            name_toks_clean.append(tl)
                    name_toks = name_toks_clean[:8]  # cap at 8 tokens

                    # Get file records whose path appears in semantic results
                    ph_sem = ",".join("?" * len(sem_paths_all))
                    c3 = self.db_conn.cursor()
                    c3.execute(
                        f"SELECT type, name, path, size FROM files "
                        f"WHERE path IN ({ph_sem}) AND type != 'Folder' LIMIT 500",
                        sem_paths_all)
                    sem_file_rows = c3.fetchall()

                    # ── Name token search with OR logic ──────────────────
                    # Each token searched independently → union of results
                    # Score = how many tokens match the filename (used for ranking)
                    token_score_map = {}   # path → match_count
                    token_row_map   = {}   # path → row tuple

                    for tok in name_toks:
                        c3.execute(
                            "SELECT type, name, path, size FROM files "
                            "WHERE (name LIKE ? OR path LIKE ?) AND type != 'Folder' LIMIT 500",
                            (f"%{tok}%", f"%{tok}%"))
                        for row in c3.fetchall():
                            p = row[2]
                            token_score_map[p] = token_score_map.get(p, 0) + 1
                            if p not in token_row_map:
                                token_row_map[p] = row

                    # Apply ext filter if any format keywords were detected
                    if ext_filter_set:
                        token_row_map = {
                            p: r for p, r in token_row_map.items()
                            if os.path.splitext(p)[1].lower() in ext_filter_set
                        }
                        token_score_map = {p: s for p, s in token_score_map.items() if p in token_row_map}
                        # Also filter sem_file_rows by ext
                        sem_file_rows = [r for r in sem_file_rows
                                         if os.path.splitext(r[2])[1].lower() in ext_filter_set]

                    # Sort token results by match count (descending) — files matching more tokens rank higher
                    token_rows_sorted = sorted(token_row_map.items(),
                                               key=lambda x: token_score_map.get(x[0], 0),
                                               reverse=True)
                    token_rows = [row for _, row in token_rows_sorted]

                    # ── Folder search ─────────────────────────────────────
                    folder_rows = []
                    for tok in name_toks[:5]:
                        c3.execute(
                            "SELECT type, name, path, size FROM files "
                            "WHERE name LIKE ? AND type = 'Folder' LIMIT 200",
                            (f"%{tok}%",))
                        folder_rows.extend(c3.fetchall())

                    # Deduplicate by path; BM25 file_res paths excluded (they're already on top)
                    seen_paths = {r[2] for r in bm25_files}
                    for row in (sem_file_rows + token_rows):
                        if row[2] not in seen_paths:
                            seen_paths.add(row[2])
                            ai_file_extras.append(row)
                    for row in folder_rows:
                        if row[2] not in seen_paths:
                            seen_paths.add(row[2])
                            ai_file_extras.append(row)
                except Exception as _dbe:
                    print(f"[AI Search] file name DB query error: {_dbe}")

            # Re-rank using BM25 on filename — higher match-count files rise naturally
            import re as _re_ai2
            def _tok(t):
                return [x.lower() for x in _re_ai2.split(r'[\s\W]+', str(t)) if len(x) >= 2]

            def _bm25_rank(rows, q_tok):
                if not rows or not q_tok:
                    return rows
                try:
                    from rank_bm25 import BM25Okapi
                    corpus = [_tok(os.path.basename(r[2]) + " " + r[2]) for r in rows]
                    bm25_fn = BM25Okapi(corpus)
                    raw = bm25_fn.get_scores(q_tok)
                    ranked = sorted(zip(rows, raw), key=lambda x: x[1], reverse=True)
                    return [r for r, _ in ranked]
                except Exception:
                    return rows

            q_tok_fn = _tok(query)
            bm25_files_ranked = _bm25_rank(bm25_files, q_tok_fn)  # BM25 group re-ranked
            ai_extras_ranked  = _bm25_rank(ai_file_extras, q_tok_fn)  # AI extras ranked by name relevance

            # Combine: BM25 on top, AI extras below
            file_res_new = bm25_files_ranked + ai_extras_ranked

            only_files_new   = [r for r in file_res_new if str(r[0]).lower() != "folder"]
            only_folders_new = [r for r in file_res_new if str(r[0]).lower() == "folder"]

            # Priority sort: office/pdf/msg first, txt/md middle, log/html/code/... last
            cont_res         = self._sort_priority(cont_res, 0)
            only_files_new   = self._sort_priority(only_files_new, 1)

            # ── C: Update all 3 trees on main thread ─────────────────────
            def _update_ui():
                try:
                    if not self.active_result_win or not tk.Toplevel.winfo_exists(self.active_result_win):
                        return
                    self._all_content_data = cont_res
                    self._all_files_data   = only_files_new
                    # FIX 3: cache AI-specific results for split pane display
                    self._ai_cont_res  = cont_res
                    self._ai_file_res  = only_files_new + only_folders_new

                    for item in self.tree_c.get_children():   self.tree_c.delete(item)
                    for item in self.tree_f.get_children():   self.tree_f.delete(item)
                    for item in self.tree_fol.get_children(): self.tree_fol.delete(item)

                    for rn, item in enumerate(cont_res, start=1):
                        self._insert_row(self.tree_c, item, rn, 0)
                    for rn, item in enumerate(only_files_new, start=1):
                        self._insert_row(self.tree_f, item, rn, 1)
                    for rn, item in enumerate(only_folders_new, start=1):
                        self._insert_row(self.tree_fol, item, rn, 2)

                    if self.content_filter_count_label:
                        self.content_filter_count_label.config(text=f"{len(cont_res)} files")
                    if self.filter_count_label:
                        self.filter_count_label.config(text=f"{len(only_files_new)} files")
                    self.nb.tab(0, text=" File Name ")
                    self.nb.tab(1, text=" Folder Name ")

                    self._ai_mode_active = True
                    self._ai_active_query = query  # v5.4: remember which query AI mode belongs to
                    _set_btn("↩ BM25", "#aaaaaa", "#2a3a2a")
                    _set_status("🔀 Hybrid: BM25 + AI", "#7ec8e3")
                    # Open AI split window (left/right 60-40)
                    self.root.after(50, self._open_ai_split_win)
                except Exception as _ue:
                    print(f"[AI Search] UI update error: {_ue}")

            self.root.after(0, _update_ui)

        except Exception as _e:
            print(f"[AI Search] Error: {_e}")
            _set_btn("🤖 AI Search", "#7ec8e3", "#1e3a5f")
            _set_status("⚠️ AI search failed", "#ff6666")

    def _smart_search_realtime(self, q, version, box_x, box_y, box_h, adv_mode=False):
        try:
            if len(q) < 2: return
            cleaned_q, op, size_val = parse_size_filter(q)
            kw = [k.lower() for k in cleaned_q.split() if k]
            # v5.8: same stopword strip as the live MFT scan -- this is a
            # SEPARATE search path (DB-backed) that feeds _last_bm25_file_res,
            # which the "Default" pane switches to as soon as AI Search /
            # Advanced mode is toggled on. Without this, "to" (from "relevant
            # to HILS") still slipped into the OR-fallback LIKE clause below
            # and matched every ".toc" file the moment AI Search was clicked,
            # even though the live-scan path (used before AI Search) was
            # already fixed.
            kw = _strip_stopwords(kw)

            # Use persistent connection — no open/close overhead per keystroke
            conn = self.db_conn
            if conn is None:
                return   # DB not ready yet (indexing in progress)
            try:
                conn.execute("SELECT 1")  # quick liveness check
            except Exception:
                self.db_conn = None
                return   # connection closed (indexing just started)
            c = conn.cursor()
            
            if not kw and op is not None:
                file_sql = f"SELECT type, name, path, size FROM files WHERE size {op} ? LIMIT 1000"
                c.execute(file_sql, (size_val,))
                file_res = c.fetchall()
                cont_res = []
            elif not kw:
                conn.close(); return
            else:
                # name-only AND: all keywords must appear in the file/folder name
                # Only use keywords >= 2 chars for name matching to avoid false positives like 'a'
                size_clause = f" AND size {op} ?" if op is not None else ""
                size_extra  = [size_val] if op is not None else []
                # Name/folder search: require at least one keyword >= 3 chars (anchor).
                # Short keywords (1-2 chars) are included in LIKE only when anchored by a long keyword.
                # e.g. "simpack a" → ok (simpack anchors), "a b" → skip (no anchor).
                has_anchor = any(_is_anchor_kw(k) for k in kw)
                if has_anchor:
                    name_conds  = ["name LIKE ?" for k in kw]
                    name_params = [f"%{k}%" for k in kw]
                    file_name_sql = ("SELECT type, name, path, size FROM files WHERE type != 'Folder' AND "
                                + " AND ".join(name_conds) + size_clause + " LIMIT 1000")
                    c.execute(file_name_sql, name_params + size_extra)
                    file_name_res = c.fetchall()

                    # ── OR fallback: if AND returns nothing, try OR so partial matches surface ──
                    if not file_name_res and len(kw) > 1:
                        or_conds  = ["name LIKE ?" for k in kw]
                        or_params = [f"%{k}%" for k in kw]
                        file_name_sql_or = ("SELECT type, name, path, size FROM files WHERE type != 'Folder' AND ("
                                    + " OR ".join(or_conds) + ")" + size_clause + " LIMIT 1000")
                        c.execute(file_name_sql_or, or_params + size_extra)
                        file_name_res = c.fetchall()
                    # ─────────────────────────────────────────────────────────────────────────────

                    folder_name_sql = ("SELECT type, name, path, size FROM files WHERE type = 'Folder' AND "
                                + " AND ".join(name_conds) + size_clause + " LIMIT 1000")
                    c.execute(folder_name_sql, name_params + size_extra)
                    folder_name_res = c.fetchall()

                    # v7.10: "Whole word" toggle -- SQL LIKE has no notion of
                    # word boundaries, so fetch the (broader) substring hits
                    # above as usual, then post-filter in Python when the
                    # toggle is on. This drops buried-substring hits like
                    # "readasync.xml" for "adas" while keeping "ADAS_systems.pdf".
                    if self.whole_word_var.get():
                        file_name_res   = [r for r in file_name_res
                                            if all(_kw_matches(k, r[1].lower(), True) for k in kw)]
                        folder_name_res = [r for r in folder_name_res
                                            if all(_kw_matches(k, r[1].lower(), True) for k in kw)]

                    file_res = file_name_res + folder_name_res
                else:
                    file_res = []
                
                # ── Content search (v10.63) ──────────────────────────────────
                # SPEED STRATEGY for large DB (44GB):
                #
                # LIKE '%phrase%' on content_store = full table scan = slow on 44GB
                # FTS5 trigram MATCH = uses index = fast even on large DB
                #
                # New approach:
                #   1. Extract alphanum tokens from phrase (FTS-safe parts)
                #   2. FTS MATCH → get small candidate path set (fast, indexed)
                #   3. LIKE verify full phrase on candidate set only (fast, small set)
                #
                # This avoids scanning all 44GB for every keystroke.
                # LIKE is only run on the ~few hundred paths FTS returns, not millions.

                import re as _re

                _STOP_WORDS = {
                    'a','an','the','is','in','on','at','to','of','or','and','as',
                    'be','by','do','for','has','had','he','her','him','his',
                    'how','i','if','it','its','me','my','no','not','off',
                    'our','out','own','so','than','that','them','then',
                    'they','this','us','was','we','who','why','will','with',
                    'you','your','also','been','but','can','did','does','from',
                    'get','got','have','into','just','may','new','now','one',
                    'see','set','she','time','what','when','which','would',
                }

                def _is_fts_safe(k):
                    # NOTE: content_index uses SQLite's fts5 'trigram' tokenizer, which
                    # can only match terms >= 3 characters — SQLite silently returns ZERO
                    # rows (not an error) for shorter terms, even a valid 2-char CJK word
                    # like 解析. So this length check must stay at 3 regardless of script;
                    # short CJK anchors are instead routed to the LIKE fallback below via
                    # _is_anchor_kw(), which has no such length limit.
                    FTS5_OPS = set('+-*:^"()：、。・<>@[]{}|\\/?!#$%&=~`\'')
                    if len(k) < 3: return False
                    if any(ch in FTS5_OPS for ch in k): return False
                    if '-' in k: return False   # hyphen is NOT operator in FTS5
                    if '.' in k and not k.replace('.','').isdigit(): return False  # dots cause issues
                    if _re.search(r'[^\w\u3000-\u9fff\uff00-\uffef\u4e00-\u9fff]', k): return False
                    return True

                # ── Tunable display limits (adv_mode = Advanced button pressed) ──
                if adv_mode:
                    FTS_LIMIT      = 5000  # Advanced: full result set
                    DISPLAY_LIMIT  = 2000  # Advanced: show all
                    BM25_THRESHOLD = 0.0   # Advanced: no threshold, show everything
                else:
                    FTS_LIMIT      = 300   # Realtime: fast, top candidates only
                    DISPLAY_LIMIT  = 100   # Realtime: max 100 rows for speed
                    BM25_THRESHOLD = 0.05  # Realtime: drop noise < 5% of top
                # ─────────────────────────────────────────────────────────────

                def _fts_candidate_paths(cursor, tokens):
                    """Use FTS index to get candidate path set — fast even on 44GB."""
                    fts_tokens = [k for k in tokens
                                  if len(k) >= 3 and k.lower() not in _STOP_WORDS and _is_fts_safe(k)]
                    if not fts_tokens:
                        return None   # no FTS tokens available
                    fts_terms = " AND ".join(f'"{k}"' for k in fts_tokens)
                    try:
                        cursor.execute(
                            "SELECT path FROM content_index WHERE content MATCH ? LIMIT ?",
                            (fts_terms, FTS_LIMIT))
                        return set(r[0] for r in cursor.fetchall())
                    except Exception as _e:
                        print(f"FTS error [{fts_terms}]: {_e}")
                        return None

                def _like_on_paths(cursor, phrase, paths_set):
                    """Run LIKE verify only on known candidate paths — avoids full scan."""
                    if not paths_set:
                        return []
                    ph = ",".join("?" * len(paths_set))
                    cursor.execute(
                        f"SELECT path FROM content_store WHERE path IN ({ph}) AND content LIKE ?",
                        list(paths_set) + [f"%{phrase}%"])
                    return set(r[0] for r in cursor.fetchall())

                def _paths_to_rows(cursor, paths_set):
                    if not paths_set: return []
                    ph = ",".join("?" * len(paths_set))
                    cursor.execute(
                        f"SELECT f.path, f.size FROM files f WHERE f.path IN ({ph})",
                        list(paths_set))
                    return cursor.fetchall()

                phrase = cleaned_q.strip()
                cont_res = []

                if _is_anchor_kw(phrase):
                    # Extract alphanum tokens from phrase for FTS narrowing
                    phrase_tokens = [k for k in _re.split(r'[\s:+\-<>@"()#.\[\]{}|\\/?!&=~`]+', phrase) if k]

                    # Step 1: FTS → candidate paths (indexed, fast)
                    candidates = _fts_candidate_paths(c, phrase_tokens)

                    if candidates is not None and len(candidates) > 0:
                        if len(phrase_tokens) == 1 and _is_fts_safe(phrase_tokens[0]):
                            # Single clean token → FTS result is already exact, no LIKE needed
                            cont_res = _paths_to_rows(c, candidates)
                        else:
                            # Multi-token or special chars → LIKE verify on candidate subset
                            verified = _like_on_paths(c, phrase, candidates)
                            if verified:
                                cont_res = _paths_to_rows(c, verified)
                            else:
                                # LIKE verify missed → FTS results are good enough
                                cont_res = _paths_to_rows(c, candidates)
                    elif candidates is not None and len(candidates) == 0:
                        # FTS returned nothing → no results
                        cont_res = []
                    else:
                        # No FTS tokens available (all special chars) → LIKE only on full table
                        # This is the slow path, only hits for queries like "+81" alone
                        useful_like = [k for k in phrase_tokens
                                       if _is_anchor_kw(k) and k.lower() not in _STOP_WORDS]
                        if useful_like:
                            best = max(useful_like, key=len)  # use longest token only
                            try:
                                c.execute(
                                    "SELECT path FROM content_store WHERE content LIKE ? LIMIT 3000",
                                    (f"%{best}%",))
                                fallback_paths = set(r[0] for r in c.fetchall())
                                if len(phrase_tokens) > 1:
                                    fallback_paths = _like_on_paths(c, phrase, fallback_paths)
                                cont_res = _paths_to_rows(c, fallback_paths)
                            except: pass
                
            # Do NOT close conn — it's the persistent db_conn, reused every search

            # ── BM25 ranking for File Name + Folder Name ────────────────────
            if file_res and len(cleaned_q.split()) >= 1:
                try:
                    from rank_bm25 import BM25Okapi
                    import re as _re_bm25fn
                    import time as _time_fn
                    def _tok_fn(text):
                        return [t.lower() for t in _re_bm25fn.split(r'[\s\W]+', str(text)) if len(t) >= 2]
                    q_tok_fn = _tok_fn(cleaned_q)
                    if q_tok_fn:
                        # v1.5: filename weighted 3x (more important than full path)
                        def _weighted_fn(path):
                            fname = os.path.basename(path)
                            return _tok_fn(f"{fname} {fname} {fname} {path}")
                        corpus_fn = [_weighted_fn(r[2]) for r in file_res]
                        bm25_fn = BM25Okapi(corpus_fn)
                        raw_fn = bm25_fn.get_scores(q_tok_fn)
                        max_fn = max(raw_fn) if max(raw_fn) > 0 else 1.0

                        # v1.5: extension boost for File Name tab
                        _EXT_BOOST_FN = {
                            '.pdf': 1.30, '.docx': 1.30, '.doc': 1.25,
                            '.xlsx': 1.20, '.xls': 1.15, '.csv': 1.10,
                            '.pptx': 1.20, '.ppt': 1.15, '.msg': 1.15,
                            '.one': 1.10, '.txt': 1.00, '.md': 1.00,
                            '.py': 0.90,  '.log': 0.75, '.ini': 0.70,
                        }
                        # v1.5: recency boost — newer files rank higher
                        _now_fn = _time_fn.time()
                        def _recency_fn(path):
                            try:
                                age_days = (_now_fn - os.path.getmtime(path)) / 86400
                                return 1.0 / (1.0 + age_days / 365)
                            except Exception:
                                return 0.5

                        def _final_fn(i, path):
                            bm25_norm = raw_fn[i] / max_fn
                            ext       = os.path.splitext(path)[1].lower()
                            ext_b     = _EXT_BOOST_FN.get(ext, 1.0)
                            rec_b     = _recency_fn(path)
                            # v9.7 fix: relevance must dominate ranking — file-type
                            # preference is now a small additive nudge (±0.03 max),
                            # not a multiplier directly on the relevance term. The
                            # old formula (bm25_norm * ext_b * 0.75) let a 30%
                            # extension boost outrank a genuinely more relevant
                            # result in a different format (e.g. a highly relevant
                            # .zip/.spck file losing to a weakly-relevant .pdf just
                            # because .pdf carries ext_b=1.30 vs .zip's default 1.0).
                            # ext_b's 0.70–1.30 range maps to a ±0.03 nudge here —
                            # relevant only as a tie-breaker between near-equal
                            # matches, never enough to flip a real relevance gap.
                            return bm25_norm * 0.90 + (ext_b - 1.0) * 0.10 + rec_b * 0.05

                        file_res = [file_res[i] for i in sorted(
                            range(len(file_res)),
                            key=lambda i: _final_fn(i, file_res[i][2]),
                            reverse=True)]
                except Exception:
                    pass  # keep original order if BM25 fails
            # Cap file_res to DISPLAY_LIMIT for realtime (adv_mode keeps all)
            _file_display_limit = 2000 if adv_mode else 100
            file_res = file_res[:_file_display_limit]
            # ────────────────────────────────────────────────────────────────

            # ── BM25 ranking for File Content ────────────────────────────────
            # Default: BM25 only (fast). AI Search button triggers hybrid later.
            sem_res = []
            bm25_scores = {}
            size_map = {}

            # ── Display limits — controlled by adv_mode ─────────────────────────────
            if adv_mode:
                DISPLAY_LIMIT  = 2000  # Advanced: all results
                BM25_THRESHOLD = 0.0   # Advanced: no threshold filtering
            else:
                DISPLAY_LIMIT  = 100   # Realtime: top 100 for fast display
                BM25_THRESHOLD = 0.05  # Realtime: discard noise below 5% of top score
            # ─────────────────────────────────────────────────────────────────────────────

            if cont_res:
                try:
                    from rank_bm25 import BM25Okapi
                    import re as _re_bm25

                    def _tokenize(text):
                        return [t.lower() for t in _re_bm25.split(r"[\s\W]+", str(text)) if len(t) >= 2]

                    q_tokens = _tokenize(cleaned_q)
                    if q_tokens:
                        # ⚡ SPEED: rank by filename + 2 parent dirs — skip content DB fetch.
                        # FTS already guarantees file contains keyword → BM25 only re-ranks.
                        # ~20x faster than content-fetch, results not truncated.
                        paths_in_cont = [r[0] for r in cont_res]
                        size_map = {r[0]: r[1] for r in cont_res}

                        path_order = []
                        corpus = []
                        for path in paths_in_cont:
                            fname  = os.path.basename(path)
                            parent = os.path.basename(os.path.dirname(path))
                            gp     = os.path.basename(os.path.dirname(os.path.dirname(path)))
                            # v1.5: filename weighted 3x — most relevant signal for local search
                            text = f"{fname} {fname} {fname} {parent} {parent} {gp}"
                            path_order.append(path)
                            corpus.append(_tokenize(text))

                        bm25 = BM25Okapi(corpus)
                        raw_scores = bm25.get_scores(q_tokens)
                        max_s = max(raw_scores) if max(raw_scores) > 0 else 1.0

                        # v1.5: extension boost + recency boost for File Content tab
                        _EXT_BOOST_C = {
                            '.pdf': 1.30, '.docx': 1.30, '.doc': 1.25,
                            '.xlsx': 1.20, '.xls': 1.15, '.csv': 1.10,
                            '.pptx': 1.20, '.ppt': 1.15, '.msg': 1.15,
                            '.one': 1.10, '.txt': 1.00, '.md': 1.00,
                            '.py': 0.90,  '.log': 0.75, '.ini': 0.70,
                        }
                        import time as _time_c
                        _now_c = _time_c.time()
                        def _recency_c(path):
                            try:
                                age_days = (_now_c - os.path.getmtime(path)) / 86400
                                return 1.0 / (1.0 + age_days / 365)
                            except Exception:
                                return 0.5

                        bm25_scores = {}
                        for i, path in enumerate(path_order):
                            bm25_norm = raw_scores[i] / max_s
                            ext       = os.path.splitext(path)[1].lower()
                            ext_b     = _EXT_BOOST_C.get(ext, 1.0)
                            rec_b     = _recency_c(path)
                            # v9.7 fix: same as the File Name tab above — relevance
                            # must dominate, extension type is now a small additive
                            # nudge (±0.03 max) instead of a multiplier that could
                            # outrank a genuinely more relevant result in a
                            # different file format.
                            bm25_scores[path] = bm25_norm * 0.90 + (ext_b - 1.0) * 0.10 + rec_b * 0.05
                    else:
                        size_map = {r[0]: r[1] for r in cont_res}
                except Exception as _bm25_e:
                    print(f"[BM25] Error: {_bm25_e}")
                    size_map = {r[0]: r[1] for r in cont_res}

            # Sort by BM25 score → threshold → display cap
            all_paths = list(bm25_scores.keys()) if bm25_scores else [r[0] for r in cont_res]
            hybrid = []
            for path in all_paths:
                score = bm25_scores.get(path, 0.0)
                hybrid.append((path, size_map.get(path, 0), score))
            hybrid.sort(key=lambda x: x[2], reverse=True)
            max_h = hybrid[0][2] if hybrid else 1.0
            # Noise filter: drop results scoring below threshold relative to top result
            if max_h > 0 and bm25_scores:
                cutoff = max_h * BM25_THRESHOLD
                hybrid = [x for x in hybrid if x[2] >= cutoff]
            # Limit displayed rows to keep UI responsive
            hybrid = hybrid[:DISPLAY_LIMIT]
            cont_res = [(p, sz, int((sc / max_h) * 99) if max_h > 0 else 0)
                        for p, sz, sc in hybrid]
            # ────────────────────────────────────────────────────────────────

            if version == self.current_search_id:
                # Priority sort at source: office/pdf/msg first, txt/md middle,
                # log/html/code/... last. All places reusing this cache (restore,
                # Advanced reset, AI merge) receive pre-sorted data.
                cont_res  = self._sort_priority(cont_res, 0)
                file_res  = self._sort_priority(file_res, 1)
                # Cache BM25 results + query for later AI Search merge
                self._last_query = q
                self._last_bm25_cont_res = cont_res
                self._last_bm25_file_res = file_res
                # v7.7 FIX: remember that the DB-backed (comprehensive) result
                # set has been painted for this search id, and snapshot it, so
                # the live MFT scan (narrower scope, may still be running/
                # finish later) merges with it in _render_mft_file_tree /
                # _render_mft_folder_tree instead of overwriting it outright.
                self._db_rendered_sid = version
                self._db_rendered_file_res = list(file_res)
                # v5.4: only clear AI mode if this refresh is for a genuinely
                # DIFFERENT query than the one AI results are currently shown
                # for. Some UI interactions (e.g. selecting a model in the AI
                # dropdown) can end up re-triggering this realtime BM25 refresh
                # for the SAME query in the background -- that used to always
                # silently kick the view back to BM25 and reset the AI Search
                # button, even though the user never asked to leave AI mode.
                # v9.12 fix: once AI Search has been turned on, keep it
                # "sticky" across new keywords instead of silently falling
                # back to BM25 and making the user re-click "AI Search"
                # every time — the semantic model is already loaded/cached
                # in memory (_load_semantic_model() is a no-op if the right
                # model is already loaded), so re-running AI search for a
                # NEW query here is cheap, not a full reload.
                # _ai_search_and_update() TOGGLES: calling it while
                # _ai_mode_active is already True restores BM25 instead of
                # refreshing, so clear the flag first here to make it
                # re-engage for the new query instead of turning itself off.
                if self._ai_mode_active and q != getattr(self, "_ai_active_query", None):
                    self._ai_mode_active = False
                    threading.Thread(target=self._ai_search_and_update, args=(q,), daemon=True).start()
                elif not (self._ai_mode_active and q == getattr(self, "_ai_active_query", None)):
                    self._ai_mode_active = False
                self._adv_mode_active = adv_mode
                # Update Advanced button label to reflect current mode
                def _update_adv_btn():
                    try:
                        if self._adv_search_btn and self._adv_search_btn.winfo_exists():
                            if adv_mode:
                                self._adv_search_btn.config(text="Simple ↩", fg="#90ee90", bg="#1a3a1a", state="normal")
                            else:
                                self._adv_search_btn.config(text="Advanced", fg="#33363c", bg="#e6e8ec", state="normal")
                    except Exception: pass
                self.root.after(0, _update_adv_btn)
                self.root.after(0, lambda: self.update_or_show_results(file_res, cont_res, q, box_x, box_y, box_h, version, sem_res))
        except Exception as e:
            if "closed database" in str(e).lower() or "cannot operate" in str(e).lower():
                self.db_conn = None  # mark as closed, will reopen after indexing
            else:
                print(f"Search Error: {e}")

    def update_or_show_results(self, file_res, cont_res, query, box_x, box_y, box_h, version, sem_res=None):
        if sem_res is None: sem_res = []
        if not self.active_result_win or not tk.Toplevel.winfo_exists(self.active_result_win):
            self.show_results(file_res, cont_res, query, self.has_args, box_x, box_y, box_h, version, sem_res)
            return
            
        if version != self.current_search_id: return
        
        self.active_result_win.title(f"Results: {query}")
        
        self.tree_c.search_id = version
        self.tree_f.search_id = version
        self.tree_fol.search_id = version
        
        for item in self.tree_c.get_children(): self.tree_c.delete(item)
        for item in self.tree_f.get_children(): self.tree_f.delete(item)
        for item in self.tree_fol.get_children(): self.tree_fol.delete(item)
        
        # cont_res / file_res are already priority-sorted at source (search thread),
        # filter preserves relative order so only file/folder split is needed here.
        only_files = [r for r in file_res if str(r[0]).lower() != "folder"]
        only_folders = [r for r in file_res if str(r[0]).lower() == "folder"]
        # v7.7 FIX: the DB-backed search is comprehensive but only knows about
        # already-indexed files. Fold in any live MFT-scan matches (same
        # search id) for very recent files/folders the index doesn't have
        # yet, instead of silently dropping them when this DB render lands.
        if version == self.current_search_id:
            only_files   = self._merge_mft_with_db(self._mft_file_res or [], is_folder=False) \
                           if only_files or self._mft_file_res else only_files
            only_folders = self._merge_mft_with_db(self._mft_folder_res or [], is_folder=True) \
                           if only_folders or self._mft_folder_res else only_folders
        self._all_files_data = only_files
        self._all_content_data = cont_res
        self.size_op_var.set(">"); self.size_num_var.set(""); self.size_unit_var.set("MB")
        self.ext_filter_var.set("")
        self.c_size_op_var.set(">"); self.c_size_num_var.set(""); self.c_size_unit_var.set("MB")
        self.c_ext_filter_var.set("")
        self.name_filter_var.set("")
        self.c_name_filter_var.set("")
        self.nb.tab(0, text=" File Name ")
        self.nb.tab(1, text=" Folder Name ")

        # v4.4: "Search again" from the History tab (Help) used to leave the
        # notebook sitting on whatever tab it was triggered from, so results
        # were invisible until the user manually clicked "File Name". Jump
        # there automatically -- but only for that flow (flag set in
        # use_query_from_hist), not on every live keystroke re-search, which
        # would otherwise yank the user off a tab they're actively reading.
        if getattr(self, "_force_file_tab", False):
            self._force_file_tab = False
            try:
                self.nb.select(0)
            except Exception:
                pass

        # Reset AI Search + Advanced button for new query — but not if this
        # is just a background refresh for the SAME query AI mode is already
        # showing results for (see matching guard in _smart_search_realtime).
        _same_ai_query = self._ai_mode_active and query == getattr(self, "_ai_active_query", None)
        if not _same_ai_query:
            self._ai_mode_active = False
            self._adv_mode_active = False
            self._adv_page = 0
            self._adv_all_cont  = []
            self._adv_all_files = []
            self._ai_cont_res   = []
            self._ai_file_res   = []
            # Close any open split windows
            self._close_adv_split()
            self._close_ai_split()
            try:
                if self._ai_search_btn and self._ai_search_btn.winfo_exists():
                    self._ai_search_btn.config(text="🤖 AI Search", fg="#7ec8e3",
                                               bg="#1e3a5f", state="normal")
                if self._adv_search_btn and self._adv_search_btn.winfo_exists():
                    self._adv_search_btn.config(text="Advanced", fg="#33363c",
                                                bg="#e6e8ec", state="normal")
            except Exception: pass
            try:
                if hasattr(self, '_hybrid_status_lbl') and self._hybrid_status_lbl.winfo_exists():
                    self._hybrid_status_lbl.config(text="📊 BM25", fg="#aaaaaa")
            except Exception: pass
        
        def bg_load_ui(tree, data, mode, current_vid):
            if not data or getattr(tree, 'search_id', 0) != current_vid: return
            CHUNK = 500
            first_chunk, remaining_chunk = data[:CHUNK], data[CHUNK:]
            offset = [0]
            for item in first_chunk:
                if getattr(tree, 'search_id', 0) != current_vid: return
                try:
                    offset[0] += 1
                    self._insert_row(tree, item, offset[0], mode)
                except: pass

            def _update_count2():
                if mode == 1 and self.filter_count_label:
                    self.filter_count_label.config(text=f"{len(self._last_bm25_file_res or [])} files")
                elif mode == 0 and self.content_filter_count_label:
                    self.content_filter_count_label.config(text=f"{len(self._last_bm25_cont_res or [])} files")

            def load_rest(step=0):
                if not self.active_result_win or not tk.Toplevel.winfo_exists(self.active_result_win): return
                if getattr(tree, 'search_id', 0) != current_vid: return
                sub_chunk = remaining_chunk[step: step + CHUNK]
                if sub_chunk:
                    for item in sub_chunk:
                        if getattr(tree, 'search_id', 0) != current_vid: return
                        try:
                            offset[0] += 1
                            self._insert_row(tree, item, offset[0], mode)
                        except: pass
                    self.root.after(5, lambda: load_rest(step + CHUNK))
                else:
                    _update_count2()
            if remaining_chunk:
                self.root.after(10, lambda: load_rest(0))
            else:
                self.root.after(20, _update_count2)

        bg_load_ui(self.tree_c, cont_res, 0, version)
        bg_load_ui(self.tree_f, only_files, 1, version)
        bg_load_ui(self.tree_fol, only_folders, 2, version)

        # Reset pane-tree dicts to single-pane state for new search
        self._c_pane_trees   = {"main": self.tree_c}
        self._f_pane_trees   = {"main": self.tree_f}
        self._fol_pane_trees = {"main": self.tree_fol}

    def _parse_tier_filter(self, raw_text):
        """Parse a *standalone* tier expression (no '--update data' prefix) —
        used by the Update DB button so the user can just type 'tier 1,2' or
        even '1,2' into the Searchbox and click the button, instead of typing
        the full '--update data tier 1,2' command and pressing Enter.
        Accepts: '', 'tier 1', 'tier 1,2', '1,2', '1 2 3', 'tier 1 tier 2', ...
        Returns (selected_tiers, is_tier_expr):
          - ('', True)         -> box empty -> (None, True)  = no filter, all tiers
          - valid tier text    -> ({0-indexed tiers}, True)
          - anything else      -> (None, False) = not a tier expression at all
                                   (e.g. a normal search query like 'SR01403940'),
                                   caller should fall back to a full scan.
        """
        text = (raw_text or "").strip()
        if not text:
            return None, True
        rest = re.sub(r'(?i)\btier\b', ' ', text)
        if not re.fullmatch(r'[\d,\s]+', rest.strip()):
            return None, False
        nums = re.findall(r'\d+', rest)
        parsed = {int(n) - 1 for n in nums if 1 <= int(n) <= 4}
        if not parsed:
            return None, False
        return parsed, True

    def handle_action(self):
        raw_query = self.entry_var.get().strip()
        if not raw_query: return
        
        if getattr(self, '_hist_idle_timer', None):
            try: self.root.after_cancel(self._hist_idle_timer)
            except Exception: pass
            self._hist_idle_timer = None
        self._save_hist(raw_query)
        q_norm = raw_query
        if q_norm.lower() in ["--exit", "--quit"]: (self.root.destroy() or None); return
            
        q_stripped = q_norm.strip()
        q_stripped_lower = q_stripped.lower()
        if q_stripped_lower == "--update data" or q_stripped_lower.startswith("--update data "):
            # v2.5: optional tier filter — lets you index just the tiers you
            # care about instead of the whole drive. Accepts any of:
            #   --update data                    -> all tiers (1-4)
            #   --update data tier 1              -> tier 1 only
            #   --update data tier 1,2            -> tiers 1 and 2
            #   --update data tier 1 tier 2       -> tiers 1 and 2
            #   --update data tier 1 2 3          -> tiers 1, 2, 3
            _tier_arg = q_stripped[len("--update data"):].strip()
            _selected_tiers = None  # None = all tiers (default, unchanged behavior)
            if _tier_arg:
                _rest = _tier_arg.lower().replace("tier", " ")
                _nums = re.findall(r'\d+', _rest)
                _parsed = set()
                for n in _nums:
                    tn = int(n)
                    if 1 <= tn <= 4:
                        _parsed.add(tn - 1)  # store 0-indexed to match _ext_tier()
                if _parsed:
                    _selected_tiers = _parsed
                else:
                    self.status_label.config(text="Bad tier filter — use e.g. 'tier 1,2'", fg="#ff5555")
                    self.entry_var.set(""); return
            _tier_desc = ("ALL (1-4)" if _selected_tiers is None
                          else ",".join(str(t + 1) for t in sorted(_selected_tiers)))
            self._update_db_running = True  # lock AI Search/model combo immediately, don't wait for the thread
            self._set_index_status(f"Updating (tier {_tier_desc})...", "#ffcc00"); self._ramp_blink_start(); self.root.update_idletasks()
            threading.Thread(target=self.indexing_worker, args=(_selected_tiers,), daemon=True).start()
            self.entry_var.set(""); return

        q_norm = unicodedata.normalize('NFKC', q_norm)
        url = self.resolve_id(q_norm)
        if url: (webbrowser.open(url) or None); return
        
        if re.match(r"^(abaqus|abq\d{4})\s+cae$", q_norm, re.I):
            try:
                # cmd /c start launches detached, finds abaqus in PATH, opens GUI
                subprocess.Popen(f'cmd /c start "" {q_norm}', shell=True, close_fds=True)
            except Exception as e:
                messagebox.showerror("Launch Error", f"Cannot launch: {q_norm}\n{e}")
        else:
            # v3.4 FIX: this used to just .lift() the window and do nothing
            # else whenever active_result_win already existed — harmless when
            # Enter is pressed after typing (on_key_release's debounce had
            # already fired the real search on keystrokes), but silently
            # broken for any caller that sets entry_var programmatically and
            # calls handle_action() directly without going through
            # <KeyRelease> first — e.g. "Search again" from History, which
            # left the keyword sitting in the box with no results ever
            # fetched. Now this branch always launches a fresh search itself
            # (same calls on_key_release makes), and additionally lifts the
            # window to the front if one was already open.
            if self.active_result_win and tk.Toplevel.winfo_exists(self.active_result_win):
                self.active_result_win.lift()
            self.current_search_id += 1; self.root.update_idletasks()
            box_x, box_y, box_h = self.root.winfo_x(), self.root.winfo_y(), self.root.winfo_height()
            threading.Thread(target=self._mft_scan_search, args=(
                q_norm, self.current_search_id, box_x, box_y, box_h), daemon=True).start()
            if self.db_conn is not None:
                threading.Thread(target=self._smart_search_realtime, args=(
                    q_norm, self.current_search_id, box_x, box_y, box_h, False), daemon=True).start()

    # ── Priority sort helpers ──────────────────────────────────────────────────
    # v2.5: user-defined 4-tier system — used for BOTH (a) the order files are
    # extracted in during --update data (tier 1 finishes across the whole
    # filesystem before tier 2 starts, etc. — so an interrupted/partial run
    # always has the highest-value content indexed first) and (b) the "#"
    # sort order shown in File Content search results. Single source of
    # truth so the two never drift apart.
    #
    # Tier 1: primary office documents — the main things people search for
    _PRIORITY_EXTS_TIER0 = {
        '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.pdf', '.csv',
    }
    # Tier 2: everyday readable content — email, notes, plain text, logs
    _PRIORITY_EXTS_TIER1 = {
        '.msg', '.eml', '.txt', '.one', '.log',
    }
    # Tier 3: scripts / config / domain-specific engineering files
    #   .spck = Simpack model files — domain-specific to this environment
    _PRIORITY_EXTS_TIER2 = {
        '.spck', '.py', '.env', '.ini', '.cfg', '.conf',
        '.js', '.ts', '.bat', '.ps1', '.json', '.yaml', '.yml', '.toml',
    }
    # Tier 4: markup / query / misc structured text
    _PRIORITY_EXTS_TIER3 = {
        '.htm', '.html', '.sql', '.xml', '.css', '.md', '.rst', '.tsv',
    }
    _ALL_TIERS = [_PRIORITY_EXTS_TIER0, _PRIORITY_EXTS_TIER1, _PRIORITY_EXTS_TIER2, _PRIORITY_EXTS_TIER3]

    @classmethod
    def _ext_tier(cls, ext):
        for i, tier_exts in enumerate(cls._ALL_TIERS):
            if ext in tier_exts:
                return i
        return len(cls._ALL_TIERS)  # unlisted extension -> put at the very end

    # v9.10: tier only breaks ties WITHIN a band of this many
    # originally-adjacent (by relevance) results, instead of being an
    # absolute override of relevance order across the whole result set.
    # Tune this if results still feel too format-biased (raise it) or too
    # relevance-only / not enough office-doc preference (lower it).
    _TIER_SORT_BAND_SIZE = 8

    def _sort_priority(self, rows, mode):
        """Sort theo tier (xem _ext_tier ở trên) — used as a gentle
           tie-breaker among near-equally-relevant results, NOT an absolute
           override of the relevance order `rows` already arrived in.

           v9.10 fix: the old version sorted PURELY by (grp, tier,
           drive_score), throwing away the incoming relevance order
           entirely except as a same-tier tie-breaker. That meant ANY
           .pdf/.docx (tier 0) always sorted above EVERY other file format,
           regardless of how relevant it actually was — a barely-relevant
           PDF could rank above a highly relevant .zip/.spck simply because
           archives/simulation files aren't in the tier list (fall through
           to the lowest tier). Now the incoming rank is banded into groups
           of _TIER_SORT_BAND_SIZE and used as the PRIMARY key -- tier can
           only reorder results that were already close in relevance, never
           flip a real relevance gap.

           Outermost key unchanged: match group — rows tagged as an
           OR-fallback match (row[5] == 1, used by MFT multi-keyword
           search) always sort BELOW rows that matched the full AND query
           (row[5] == 0 or untagged).
        """
        def _key(indexed_row):
            i, r = indexed_row
            path = r[0] if mode == 0 else r[2]
            drive = path[0].upper() if path and len(path) >= 2 and path[1] == ':' else ''
            drive_score = 99 if drive == 'C' else 0
            ext = os.path.splitext(path)[1].lower()
            tier = self._ext_tier(ext)
            grp = r[5] if len(r) > 5 else 0
            band = i // self._TIER_SORT_BAND_SIZE
            return (grp, band, tier, drive_score, i)
        return [r for _, r in sorted(enumerate(rows), key=_key)]

    # ── Sortable columns (class methods — shared across ALL trees in ALL panes) ──
    def _sort_tree(self, tree, col, reverse):
        """Sort treeview by column; toggle direction on repeated click.
        Works with any Treeview (main/adv/ai, across all 3 tabs)."""
        try:
            items = [(tree.set(k, col), k) for k in tree.get_children("")]
            # Try numeric sort for size (contains digits + unit)
            def _sort_key(val):
                v = val[0]
                # File Size: convert "1.2 MB" → bytes for proper numeric sort
                _units = {'B': 1, 'KB': 1024, 'MB': 1024**2, 'GB': 1024**3, 'TB': 1024**4}
                parts = v.strip().split()
                if len(parts) == 2 and parts[1] in _units:
                    try: return (0, float(parts[0]) * _units[parts[1]])
                    except: pass
                # Date: already ISO-ish "YYYY-MM-DD HH:MM" → sorts correctly as string
                # Fallback: case-insensitive string
                return (1, v.lower())
            items.sort(key=lambda x: _sort_key(x), reverse=reverse)
            for idx, (_, k) in enumerate(items):
                tree.move(k, "", idx)
            # v5.8: #0 is icon-only now (no row numbers to renumber after sort)
            # Update heading with arrow indicator
            for c in tree["columns"]:
                tree.heading(c, text=tree.heading(c)["text"].replace(" ▲","").replace(" ▼",""))
            arrow = " ▼" if reverse else " ▲"
            tree.heading(col, text=tree.heading(col)["text"] + arrow,
                         command=lambda: self._sort_tree(tree, col, not reverse))
        except Exception as _se:
            print(f"Sort error: {_se}")

    def _make_sortable(self, tree, sortable_cols):
        """Bind click-to-sort for the specified columns. Call this immediately when a tree
        is created (including Advanced/AI panes born from a split) so Size/
        Modified/Type columns are always sortable, regardless of which pane the tree is in."""
        for col in sortable_cols:
            tree.heading(col, command=lambda c=col: self._sort_tree(tree, c, False))

    # ══════════════════════════════════════════════════════════════════════════
    # SPLIT LAYOUT ENGINE
    # Each Tab uses the same 3-phase layout strategy:
    #
    #   Default  (1 pane):  [  Main  ]
    #   Advanced (2 panes): [  Main  ]   ← top 50%
    #                       [ Adv.   ]   ← bottom 50%
    #   +AI      (3 panes): [Main|AI ]   ← top-left 60% | right 40%
    #                       [Adv.|AI ]   ← bottom-left  | right shared
    #
    # Layout as required:
    #   Advanced only  → top/bottom 50%/50%  (fixed, no sash)
    #   AI only        → left/right  50%/50% (fixed, no sash)
    #   Advanced + AI  → outer H-split 60/40; left side continues with V-split 50/50
    # ══════════════════════════════════════════════════════════════════════════

    # ── Treeview factory helpers ──────────────────────────────────────────────
    def _build_tree_content(self, parent, label_text, label_bg, label_fg, show_label=True):
        """Build a content-style tree (path, size, mtime, type) in parent."""
        lbl_widget = None
        if show_label:
            lbl_widget = tk.Label(parent, text=label_text, bg=label_bg, fg=label_fg,
                     font=("Segoe UI", 9, "bold"))
            lbl_widget.pack(fill="x")
        tree = ttk.Treeview(parent,
                            columns=("path","size","mtime","ftype","fp"),
                            show="tree headings")
        tree.heading("#0", text="")
        tree.column("#0", width=40, minwidth=40, stretch=False, anchor="center")
        for col, w, txt in [("path",750,"File Location"),
                             ("size",90,"Size"), ("mtime",130,"Modified"), ("ftype",55,"Type")]:
            tree.heading(col, text=txt)
            tree.column(col, width=w, stretch=False, minwidth=30)
        tree["displaycolumns"] = ("path","size","mtime","ftype")
        sb  = ttk.Scrollbar(parent, orient="vertical",   command=tree.yview)
        sbx = ttk.Scrollbar(parent, orient="horizontal", command=tree.xview)
        tree.configure(yscrollcommand=sb.set, xscrollcommand=sbx.set)
        sb.pack(side="right", fill="y"); sbx.pack(side="bottom", fill="x")
        tree.pack(fill="both", expand=True)
        tree._header_label = lbl_widget  # store ref for dynamic update
        tree._base_label   = label_text  # base text (e.g. '📄 Default (33)')
        self._make_sortable(tree, ["size", "mtime", "ftype"])  # wire sort for all panes
        return tree

    def _build_tree_file(self, parent, label_text, label_bg, label_fg, show_label=True):
        """Build a file-name-style tree (name, size, mtime, type, loc) in parent."""
        lbl_widget = None
        if show_label:
            lbl_widget = tk.Label(parent, text=label_text, bg=label_bg, fg=label_fg,
                     font=("Segoe UI", 9, "bold"))
            lbl_widget.pack(fill="x")
        tree = ttk.Treeview(parent,
                            columns=("name","size","mtime","ftype","loc","fp"),
                            show="tree headings")
        tree.heading("#0", text="")
        tree.column("#0", width=40, minwidth=40, stretch=False, anchor="center")
        for col, w, txt in [("name",360,"File Name"),
                             ("size",90,"Size"), ("mtime",130,"Modified"),
                             ("ftype",55,"Type"), ("loc",520,"Location")]:
            tree.heading(col, text=txt)
            tree.column(col, width=w, stretch=False, minwidth=30)
        tree["displaycolumns"] = ("name","size","mtime","ftype","loc")
        sb  = ttk.Scrollbar(parent, orient="vertical",   command=tree.yview)
        sbx = ttk.Scrollbar(parent, orient="horizontal", command=tree.xview)
        tree.configure(yscrollcommand=sb.set, xscrollcommand=sbx.set)
        sb.pack(side="right", fill="y"); sbx.pack(side="bottom", fill="x")
        tree.pack(fill="both", expand=True)
        tree._header_label = lbl_widget
        tree._base_label   = label_text
        self._make_sortable(tree, ["size", "mtime", "ftype"])  # wire sort for all panes
        return tree

    def _build_tree_folder(self, parent, label_text, label_bg, label_fg, show_label=True):
        """Build a folder-name-style tree (name, mtime, loc) in parent."""
        lbl_widget = None
        if show_label:
            lbl_widget = tk.Label(parent, text=label_text, bg=label_bg, fg=label_fg,
                     font=("Segoe UI", 9, "bold"))
            lbl_widget.pack(fill="x")
        tree = ttk.Treeview(parent,
                            columns=("name","mtime","loc","fp"),
                            show="tree headings")
        tree.heading("#0", text="")
        tree.column("#0", width=40, minwidth=40, stretch=False, anchor="center")
        for col, w, txt in [("name",420,"Folder Name"),
                             ("mtime",130,"Modified"), ("loc",650,"Location")]:
            tree.heading(col, text=txt)
            tree.column(col, width=w, stretch=False, minwidth=30)
        tree["displaycolumns"] = ("name","mtime","loc")
        sb  = ttk.Scrollbar(parent, orient="vertical",   command=tree.yview)
        sbx = ttk.Scrollbar(parent, orient="horizontal", command=tree.xview)
        tree.configure(yscrollcommand=sb.set, xscrollcommand=sbx.set)
        sb.pack(side="right", fill="y"); sbx.pack(side="bottom", fill="x")
        tree.pack(fill="both", expand=True)
        tree._header_label = lbl_widget
        tree._base_label   = label_text
        self._make_sortable(tree, ["mtime"])  # wire sort for all panes
        return tree

    # ── Data populators ───────────────────────────────────────────────────────
    def _fill_content(self, tree, data):
        """Populate a content tree from _all_content_data items."""
        for row in tree.get_children(): tree.delete(row)
        for i, item in enumerate(self._sort_priority(data, 0), 1):
            try:
                img = get_tree_icon_image(item[0], False)
                name_txt = item[0] if img else get_file_icon(item[0], False) + item[0]
                tree.insert("", tk.END, text="", **({"image": img} if img else {}), values=(
                    name_txt,
                    format_size(item[1]) if item[1] else "",
                    get_live_mtime(item[0]), get_file_type(item[0]), item[0]))
            except: pass

    def _fill_files(self, tree, data):
        """Populate a file tree from _all_files_data items."""
        for row in tree.get_children(): tree.delete(row)
        for i, item in enumerate(self._sort_priority(data, 1), 1):
            try:
                img = get_tree_icon_image(item[2], False)
                name_txt = item[1] if img else get_file_icon(item[2], False) + item[1]
                tree.insert("", tk.END, text="", **({"image": img} if img else {}), values=(
                    name_txt,
                    format_size(item[3]) if item[3] else "",
                    get_live_mtime(item[2]), get_file_type(item[2]),
                    os.path.dirname(item[2]), item[2]))
            except: pass

    def _fill_folders(self, tree, data):
        """Populate a folder tree from folder items."""
        for row in tree.get_children(): tree.delete(row)
        for i, item in enumerate(data, 1):
            try:
                img = get_tree_icon_image(item[2], True)
                name_txt = item[1] if img else get_file_icon(item[2], True) + item[1]
                tree.insert("", tk.END, text="", **({"image": img} if img else {}), values=(name_txt, get_live_mtime(item[2]), item[2], item[2]))
            except: pass

    # ── Layout builder — called for EACH tab_frame independently ─────────────
    def _rebuild_tab_layout(self, tab_frame, tab_mode,
                            main_data, adv_data, ai_data,
                            main_tree_attr):
        """
        Destroy and rebuild the layout inside tab_frame.

        tab_mode : "c"   → File Content tab  (main = content tree)
                   "f"   → File Name tab      (main = file tree)
                   "fol" → Folder Name tab    (main = folder tree)

        Layout states (same for every tab):
          adv=False, ai=False  →  1 pane  : main only
          adv=True,  ai=False  →  2 panes : top(main) / bottom(adv)  50/50
          adv=False, ai=True   →  2 panes : left(main) / right(ai)   60/40
          adv=True,  ai=True   →  3 panes : outer H 60/40;
                                            left V-split top(main)/bot(adv) 50/50;
                                            right = ai
        adv_mode / ai_mode are stored in self._adv_mode_active / self._ai_mode_active.
        """
        adv = self._adv_mode_active
        ai  = self._ai_mode_active

        # Clear everything in the container
        for w in tab_frame.winfo_children(): w.destroy()

        # ── helpers per tab_mode ─────────────────────────────────────────────
        # Show ">5000" only when result count hits the 5000 hard cap
        _HARD_CAP = 5000
        def _fmt(n):
            return f">{_HARD_CAP}" if n >= _HARD_CAP else str(n)
        n_main = len(main_data)
        n_adv  = len(adv_data)
        n_ai   = len(ai_data)
        s_main = _fmt(n_main)
        s_adv  = _fmt(n_adv)
        s_ai   = _fmt(n_ai)

        def make_main(parent, show_label=True):
            lbl = f"📄 Default ({s_main})" if tab_mode == "c" else f"📁 Default ({s_main})"
            if tab_mode == "c":
                return self._build_tree_content(parent, lbl, BG_COLOR, TEXT_COLOR, show_label=show_label)
            elif tab_mode == "f":
                return self._build_tree_file(parent, lbl, BG_COLOR, TEXT_COLOR, show_label=show_label)
            else:
                return self._build_tree_folder(parent, lbl, BG_COLOR, TEXT_COLOR, show_label=show_label)

        def make_adv(parent):
            lbl = f"📄 Advanced ({s_adv})" if tab_mode == "c" else f"📁 Advanced ({s_adv})"
            if tab_mode == "c":
                return self._build_tree_content(parent, lbl, "#1a3a2a", "#aaffaa")
            elif tab_mode == "f":
                return self._build_tree_file(parent, lbl, "#1a2a3a", "#7ec8e3")
            else:
                return self._build_tree_folder(parent, lbl, "#1a1a2a", "#aaaaff")

        def make_ai(parent):
            _model_tags = {"jina_v3": "Jina", "bge_gemma2": "Gemma2"}
            _size_tag = _model_tags.get(_sem_model_key, _sem_model_key)
            lbl = f"🤖 AI Search [{_size_tag}] ({s_ai})"
            if tab_mode == "c":
                return self._build_tree_content(parent, lbl, "#1a1a3a", "#a0a0ff")
            elif tab_mode == "f":
                return self._build_tree_file(parent, lbl, "#2a1a2a", "#d090f0")
            else:
                return self._build_tree_folder(parent, lbl, "#2a1a1a", "#ffaa88")

        def fill_main(tree):
            if tab_mode == "c":   self._fill_content(tree, main_data)
            elif tab_mode == "f": self._fill_files(tree, main_data)
            else:                 self._fill_folders(tree, main_data)

        def fill_adv(tree):
            if tab_mode == "c":   self._fill_content(tree, adv_data)
            elif tab_mode == "f": self._fill_files(tree, adv_data)
            else:                 self._fill_folders(tree, adv_data)

        def fill_ai(tree):
            if tab_mode == "c":   self._fill_content(tree, ai_data)
            elif tab_mode == "f": self._fill_files(tree, ai_data)
            else:                 self._fill_folders(tree, ai_data)

        # ── Determine which pane-tree dict to update ─────────────────────────
        if tab_mode == "c":   pane_dict = self._c_pane_trees
        elif tab_mode == "f": pane_dict = self._f_pane_trees
        else:                 pane_dict = self._fol_pane_trees
        pane_dict.clear()

        # ── Build layout — strict 50/50 using place geometry ─────────────────
        SEP = 3  # separator thickness in px

        if not adv and not ai:
            # ── 1 pane ───────────────────────────────────────────────────────
            pane = tk.Frame(tab_frame, bg=BG_COLOR)
            pane.place(relx=0, rely=0, relwidth=1.0, relheight=1.0)
            main_tree = make_main(pane, show_label=False)
            fill_main(main_tree)
            setattr(self, main_tree_attr, main_tree)
            pane_dict["main"] = main_tree
            self._wire_tree_clicks(main_tree, tab_mode)

        elif adv and not ai:
            # ── 2 panes top/bottom strict 50/50 ──────────────────────────────
            sep_h = tk.Frame(tab_frame, bg="#555", height=SEP)
            sep_h.place(relx=0, rely=0.5, relwidth=1.0, height=SEP, anchor="w")

            top = tk.Frame(tab_frame, bg=BG_COLOR)
            top.place(relx=0, rely=0, relwidth=1.0, relheight=0.5)
            bot = tk.Frame(tab_frame, bg=BG_COLOR)
            bot.place(relx=0, rely=0.5, relwidth=1.0, relheight=0.5)

            main_tree = make_main(top); fill_main(main_tree)
            adv_tree  = make_adv(bot);  fill_adv(adv_tree)
            setattr(self, main_tree_attr, main_tree)
            pane_dict["main"] = main_tree
            pane_dict["adv"]  = adv_tree
            self._wire_tree_clicks(main_tree, tab_mode)
            self._wire_tree_clicks(adv_tree,  tab_mode)

        elif not adv and ai:
            # ── 2 panes left/right strict 50/50 ──────────────────────────────
            sep_v = tk.Frame(tab_frame, bg="#555", width=SEP)
            sep_v.place(relx=0.5, rely=0, width=SEP, relheight=1.0, anchor="n")

            left  = tk.Frame(tab_frame, bg=BG_COLOR)
            left.place(relx=0, rely=0, relwidth=0.5, relheight=1.0)
            right = tk.Frame(tab_frame, bg=BG_COLOR)
            right.place(relx=0.5, rely=0, relwidth=0.5, relheight=1.0)

            main_tree = make_main(left);  fill_main(main_tree)
            ai_tree   = make_ai(right);   fill_ai(ai_tree)
            setattr(self, main_tree_attr, main_tree)
            pane_dict["main"] = main_tree
            pane_dict["ai"]   = ai_tree
            self._wire_tree_clicks(main_tree, tab_mode)
            self._wire_tree_clicks(ai_tree,   tab_mode)

        else:
            # ── 3 panes: left 50% | right 50%; left split top/bottom 50/50 ──
            sep_v = tk.Frame(tab_frame, bg="#555", width=SEP)
            sep_v.place(relx=0.5, rely=0, width=SEP, relheight=1.0, anchor="n")

            left_frame = tk.Frame(tab_frame, bg=BG_COLOR)
            left_frame.place(relx=0, rely=0, relwidth=0.5, relheight=1.0)
            right = tk.Frame(tab_frame, bg=BG_COLOR)
            right.place(relx=0.5, rely=0, relwidth=0.5, relheight=1.0)

            sep_h = tk.Frame(left_frame, bg="#555", height=SEP)
            sep_h.place(relx=0, rely=0.5, relwidth=1.0, height=SEP, anchor="w")

            top = tk.Frame(left_frame, bg=BG_COLOR)
            top.place(relx=0, rely=0, relwidth=1.0, relheight=0.5)
            bot = tk.Frame(left_frame, bg=BG_COLOR)
            bot.place(relx=0, rely=0.5, relwidth=1.0, relheight=0.5)

            main_tree = make_main(top); fill_main(main_tree)
            adv_tree  = make_adv(bot);  fill_adv(adv_tree)
            ai_tree   = make_ai(right); fill_ai(ai_tree)
            setattr(self, main_tree_attr, main_tree)
            pane_dict["main"] = main_tree
            pane_dict["adv"]  = adv_tree
            pane_dict["ai"]   = ai_tree
            self._wire_tree_clicks(main_tree, tab_mode)
            self._wire_tree_clicks(adv_tree,  tab_mode)
            self._wire_tree_clicks(ai_tree,   tab_mode)

    # ── Public split triggers (called by Advanced / AI Search buttons) ────────

    def _wire_tree_clicks(self, tree, tab_mode):
        """Wire double-click, single-click (cell-edit), and right-click context menu
        to any Treeview — including panes created after the initial layout by split."""
        open_file     = getattr(self, "_fn_open_file",     None)
        open_explorer = getattr(self, "_fn_open_explorer", None)
        show_ctx      = getattr(self, "_fn_show_ctx_menu", None)
        on_click      = getattr(self, "_fn_on_tree_click", None)
        dbl_select    = getattr(self, "_fn_dbl_select",    None)

        if not all([open_file, open_explorer, show_ctx, on_click, dbl_select]):
            return  # handlers not ready yet

        is_folder  = (tab_mode == "fol")
        is_content = not is_folder  # both "c" and "f" tabs can open files

        # Use default-argument capture to avoid late-binding closure trap
        if is_folder:
            tree.bind("<Double-1>", lambda e, t=tree: (dbl_select(t), open_explorer(t)))
        else:
            tree.bind("<Double-1>", lambda e, t=tree: (dbl_select(t), open_file(t)))

        tree.bind("<Button-1>", lambda e, t=tree: on_click(e, t))
        tree.bind("<Button-3>", lambda e, t=tree, f=is_folder, c=is_content:
                  show_ctx(e, t, is_folder=f, is_content=c))

    def _open_adv_split_win(self):
        """Advanced activated → rebuild ALL tabs with adv=True layout."""
        if not self.active_result_win or not tk.Toplevel.winfo_exists(self.active_result_win):
            return
        # FIX 3: main=BM25 realtime, adv=full Advanced results, ai=AI results
        bm25_cont    = self._last_bm25_cont_res or []
        bm25_files   = [r for r in (self._last_bm25_file_res or []) if str(r[0]).lower() != "folder"]
        bm25_folders = [r for r in (self._last_bm25_file_res or []) if str(r[0]).lower() == "folder"]
        adv_cont     = self._adv_all_cont  or []
        adv_files    = [r for r in (self._adv_all_files or []) if str(r[0]).lower() != "folder"]
        adv_folders  = [r for r in (self._adv_all_files or []) if str(r[0]).lower() == "folder"]
        ai_cont      = self._ai_cont_res  or []
        ai_files     = [r for r in (self._ai_file_res  or []) if str(r[0]).lower() != "folder"]
        ai_folders   = [r for r in (self._ai_file_res  or []) if str(r[0]).lower() == "folder"]

        self._rebuild_tab_layout(self.c_tree_frame,   "c",
            bm25_cont,    adv_cont,    ai_cont,    "tree_c")
        self._rebuild_tab_layout(self.f_tree_frame,   "f",
            bm25_files,   adv_files,   ai_files,   "tree_f")
        self._rebuild_tab_layout(self.fol_tree_frame, "fol",
            bm25_folders, adv_folders, ai_folders, "tree_fol")

    def _close_adv_split(self):
        """Advanced deactivated → rebuild ALL tabs without adv layout."""
        if not (hasattr(self, 'c_tree_frame') and self.c_tree_frame.winfo_exists()):
            return
        bm25_cont    = self._last_bm25_cont_res or []
        bm25_files   = [r for r in (self._last_bm25_file_res or []) if str(r[0]).lower() != "folder"]
        bm25_folders = [r for r in (self._last_bm25_file_res or []) if str(r[0]).lower() == "folder"]
        ai_cont      = self._ai_cont_res  or []
        ai_files     = [r for r in (self._ai_file_res  or []) if str(r[0]).lower() != "folder"]
        ai_folders   = [r for r in (self._ai_file_res  or []) if str(r[0]).lower() == "folder"]
        self._rebuild_tab_layout(self.c_tree_frame,   "c",
            bm25_cont,    [], ai_cont,    "tree_c")
        self._rebuild_tab_layout(self.f_tree_frame,   "f",
            bm25_files,   [], ai_files,   "tree_f")
        self._rebuild_tab_layout(self.fol_tree_frame, "fol",
            bm25_folders, [], ai_folders, "tree_fol")

    def _open_ai_split_win(self):
        """AI Search activated → rebuild ALL tabs with ai=True layout."""
        if not self.active_result_win or not tk.Toplevel.winfo_exists(self.active_result_win):
            return
        # FIX 3: main=BM25 realtime, adv=Advanced full (if active), ai=AI results
        bm25_cont    = self._last_bm25_cont_res or []
        bm25_files   = [r for r in (self._last_bm25_file_res or []) if str(r[0]).lower() != "folder"]
        bm25_folders = [r for r in (self._last_bm25_file_res or []) if str(r[0]).lower() == "folder"]
        adv_cont     = self._adv_all_cont  or []
        adv_files    = [r for r in (self._adv_all_files or []) if str(r[0]).lower() != "folder"]
        adv_folders  = [r for r in (self._adv_all_files or []) if str(r[0]).lower() == "folder"]
        ai_cont      = self._ai_cont_res  or []
        ai_files     = [r for r in (self._ai_file_res  or []) if str(r[0]).lower() != "folder"]
        ai_folders   = [r for r in (self._ai_file_res  or []) if str(r[0]).lower() == "folder"]

        self._rebuild_tab_layout(self.c_tree_frame,   "c",
            bm25_cont,    adv_cont,    ai_cont,    "tree_c")
        self._rebuild_tab_layout(self.f_tree_frame,   "f",
            bm25_files,   adv_files,   ai_files,   "tree_f")
        self._rebuild_tab_layout(self.fol_tree_frame, "fol",
            bm25_folders, adv_folders, ai_folders, "tree_fol")

    def _close_ai_split(self):
        """AI Search deactivated → rebuild ALL tabs without ai layout."""
        if not (hasattr(self, 'c_tree_frame') and self.c_tree_frame.winfo_exists()):
            return
        bm25_cont    = self._last_bm25_cont_res or []
        bm25_files   = [r for r in (self._last_bm25_file_res or []) if str(r[0]).lower() != "folder"]
        bm25_folders = [r for r in (self._last_bm25_file_res or []) if str(r[0]).lower() == "folder"]
        adv_cont     = self._adv_all_cont  or []
        adv_files    = [r for r in (self._adv_all_files or []) if str(r[0]).lower() != "folder"]
        adv_folders  = [r for r in (self._adv_all_files or []) if str(r[0]).lower() == "folder"]
        self._rebuild_tab_layout(self.c_tree_frame,   "c",
            bm25_cont,    adv_cont,    [], "tree_c")
        self._rebuild_tab_layout(self.f_tree_frame,   "f",
            bm25_files,   adv_files,   [], "tree_f")
        self._rebuild_tab_layout(self.fol_tree_frame, "fol",
            bm25_folders, adv_folders, [], "tree_fol")

    def _open_triple_split_win(self):
        """Both Advanced + AI active → same as calling open_adv/ai (flags already set)."""
        self._open_adv_split_win()

    def _insert_row(self, tree, item, rn, mode):
        """Insert one row into a result tree — uses DB size, no disk I/O.
        v5.8: column #0 is icon-ONLY now (no row number). ttk.Treeview only
        supports image= on the tree column (#0), which is always the
        leftmost column and can't be reordered relative to the data
        columns, so the layout is [icon] → File Name → Size → ... .
        Emoji fallback if the icon backend isn't available or extraction
        fails for that row (emoji gets prefixed onto the Name text itself,
        since it can't be a real image=)."""
        if mode == 0:
            # item = (path, size, score) from BM25 content search
            path = item[0]
            size = item[1] if len(item) > 1 else None
            img = get_tree_icon_image(path, is_folder=False)
            name_txt = path if img else get_file_icon(path, is_folder=False) + path
            readable_sz = format_size(size) if size else ""
            mtime = get_live_mtime(path)
            ftype = get_file_type(path)
            kwargs = dict(text="", values=(name_txt, readable_sz, mtime, ftype, path))
            if img:
                kwargs["image"] = img
            tree.insert("", tk.END, **kwargs)
        elif mode == 1:
            # File Name tab: (type, name, path, size)
            img = get_tree_icon_image(item[2], is_folder=False)
            name_txt = item[1] if img else get_file_icon(item[2], is_folder=False) + item[1]
            readable_sz = format_size(item[3]) if item[3] else ""
            mtime = get_live_mtime(item[2])
            ftype = get_file_type(item[2])
            kwargs = dict(text="", values=(name_txt, readable_sz, mtime, ftype, os.path.dirname(item[2]), item[2]))
            if img:
                kwargs["image"] = img
            tree.insert("", tk.END, **kwargs)
        else:
            # Folder Name tab: (type, name, path, size)
            img = get_tree_icon_image(item[2], is_folder=True)
            name_txt = item[1] if img else get_file_icon(item[2], is_folder=True) + item[1]
            mtime = get_live_mtime(item[2])
            kwargs = dict(text="", values=(name_txt, mtime, item[2], item[2]))
            if img:
                kwargs["image"] = img
            tree.insert("", tk.END, **kwargs)

    def _draw_search_icon(self, parent, size=18, color=None):
        """v4.9: flat vector magnifying-glass icon drawn on a Canvas — looks
        consistent across OSes/fonts and matches the theme, unlike the
        🔍 emoji glyph which renders as a colorful, font-dependent picture
        that clashes with a flat UI. Thicker strokes + longer handle + a
        blue tint by default, closer to the Windows 11 / Everything look."""
        color = color or "#0078d7"
        cv = tk.Canvas(parent, width=size, height=size, bg=BG_COLOR, highlightthickness=0)
        r = size * 0.30
        cx, cy = size * 0.40, size * 0.40
        lw = max(2, round(size * 0.13))  # stroke width scales with icon size
        cv.create_oval(cx - r, cy - r, cx + r, cy + r, outline=color, width=lw)
        ang = math.radians(45)
        x1 = cx + (r + lw * 0.3) * math.cos(ang)
        y1 = cy + (r + lw * 0.3) * math.sin(ang)
        x2 = x1 + size * 0.36 * math.cos(ang)
        y2 = y1 + size * 0.36 * math.sin(ang)
        cv.create_line(x1, y1, x2, y2, fill=color, width=lw, capstyle=tk.ROUND)
        return cv

    def _close_results_window(self):
        """v2.8: results now render inside self.root instead of a separate
        Toplevel, so "closing" them means tearing down the results-only widgets
        (notebook, filter bars, AI/Advanced/Update DB buttons, resize grip) and
        shrinking the same window back down to the small floating search box --
        there is no second window to destroy/deiconify anymore."""
        self._in_results_mode = False
        try:
            if self.results_frame and self.results_frame.winfo_exists():
                self.results_frame.destroy()
        except Exception:
            pass
        self.results_frame = None
        try:
            if self._results_extra_bar and self._results_extra_bar.winfo_exists():
                self._results_extra_bar.destroy()
        except Exception:
            pass
        self._results_extra_bar = None
        try:
            self.close_btn.pack_forget()
        except Exception:
            pass
        # v4.4 fix: while results were showing, self._r_p (the ramp-light
        # frame) was repacked with after=self._results_extra_bar so it would
        # sit just left of the button group (see show_results()). Now that
        # _results_extra_bar has just been destroyed above, that "after="
        # anchor no longer exists -- simply forgetting/destroying the anchor
        # widget left r_p with a dangling packing reference and it stopped
        # being drawn at all (the ramp light vanishing after closing results).
        # Explicitly re-pack it back to its normal idle-mode spot (far right
        # edge of the search bar) so it's guaranteed visible again.
        try:
            self._r_p.pack_forget()
            self._r_p.pack(side="right", fill="y", padx=(1, 5), pady=5)
        except Exception:
            pass
        self.entry.bind("<Escape>", lambda e: self.root.destroy())
        try:
            self.root.protocol("WM_DELETE_WINDOW", lambda: self.root.destroy())
        except Exception:
            pass
        self.active_result_win = None
        self._result_win = None
        self.entry_var.set("")
        try:
            self.root.geometry(f"{SMALL_SIZE}+{self.x_pos}+5")
            self.root.after(50, self.entry.focus_set)
            # v5.9: same DWM repaint nudge as _real_shrink() -- closing results
            # shrinks the window the same way, and could leave the ramp light
            # (correct color, just unpainted) looking like it had vanished.
            self.root.after(60, self._force_repaint)
        except Exception:
            pass

    def show_results(self, file_res, cont_res, query, should_exit, box_x, box_y, box_h, version, sem_res=None):
        if sem_res is None: sem_res = []
        # v2.8: results now render inside THIS SAME window (self.root) instead of
        # a separate Toplevel that popped up while root was withdrawn. That old
        # approach felt like a stutter/flash the instant the first character was
        # typed -- one whole OS window vanishing and a different one appearing.
        # Now the top search bar (self.bg_f, built once in __init__) never moves
        # or gets recreated; we just grow the window and add a results frame
        # below it. `win` is kept as a local alias so the rest of this (long)
        # function -- originally written against a standalone Toplevel -- needs
        # no further changes below this point.
        win = self.root
        self._in_results_mode = True
        self.res_edit = None
        # New Notebook instance below always defaults to its first-added tab
        # (File Name) anyway, so just clear the flag here for consistency.
        self._force_file_tab = False

        if should_exit:
            win.protocol("WM_DELETE_WINDOW", lambda: self.root.destroy())
        else:
            win.protocol("WM_DELETE_WINDOW", self._close_results_window)
        win.geometry(f"{RESULT_SIZE}+{box_x}+{box_y}")

        self.active_result_win = win; self._result_win = win

        # Reveal the "✕" close button on the persistent search bar (hidden while
        # idle) and point it at the right close behavior for this session.
        _close_cmd = (lambda: self.root.destroy()) if should_exit else self._close_results_window
        self.close_btn.configure(command=_close_cmd)
        self.close_btn.pack(side="right", padx=(0, 12), pady=5, before=self._r_p)

        self.entry.bind("<Escape>", (lambda e: self.root.destroy()) if should_exit
                                      else (lambda e: self._close_results_window()))

        # v2.5 fix: focus_set() alone only sets *Tk-internal* focus — if the
        # window itself doesn't have OS-level window focus yet (e.g. this runs
        # from a background thread's after(0, ...) callback), keystrokes kept
        # landing nowhere until the user clicked the box manually. focus_force()
        # grabs real OS input focus for the window + widget, so typing continues
        # without a click.
        win.lift()
        win.focus_force()
        self.entry.focus_force()
        self.entry.icursor("end")

        # The results-only content (notebook, AI/Advanced/Update DB buttons,
        # resize grip) lives in its own frame below the persistent search bar,
        # so closing results is just "destroy this one frame" -- see
        # _close_results_window().
        self.results_frame = tk.Frame(win, bg=BG_COLOR)
        self.results_frame.pack(fill="both", expand=True, side="top")

        # AI Search / Advanced / Update DB buttons still live on the search bar
        # row (same as before), but grouped in their own sub-frame so they can
        # all be torn down together when results close instead of accumulating
        # on the persistent bar across searches.
        self._results_extra_bar = tk.Frame(self.bg_f, bg=BG_COLOR)
        self._results_extra_bar.pack(side="right")
        search_bar = self._results_extra_bar  # local alias used further below

        # v3.4: while results are showing, move the ramp light so it sits
        # between the Searchbox and the button group (Advanced/AI Search/
        # Update DB) instead of staying pinned at the far-right edge past all
        # the buttons. `after=` tells pack() to treat r_p as if it had been
        # packed right after search_bar, so it lands just to search_bar's
        # left (i.e. between Entry and the buttons) instead of at the very
        # end. When results close and search_bar is destroyed, r_p simply
        # settles back to the far-right edge on its own (idle-mode look,
        # unchanged from before).
        try:
            self._r_p.pack(side="right", padx=(1, 5), pady=5, after=self._results_extra_bar)
        except Exception:
            pass

        # v2.6: overrideredirect(True) removes the native titlebar/border, which
        # also removes the OS resize grip — the window used to be permanently
        # stuck at RESULT_SIZE (1350x700). Draw a small draggable "◢" handle in
        # the bottom-right corner that lets the user resize by hand; default
        # size/position stay exactly as before if the user never touches it.
        # Parented to results_frame so it's cleaned up automatically on close.
        _MIN_W, _MIN_H = 900, 500
        resize_grip = tk.Label(self.results_frame, text="◢", bg=BG_COLOR, fg="#666666",
                                font=("Segoe UI", 11, "bold"), cursor="size_nw_se")
        resize_grip.place(relx=1.0, rely=1.0, anchor="se", x=-2, y=-2)

        def _start_resize(e):
            self._resize_startx, self._resize_starty = e.x_root, e.y_root
            self._resize_startw = win.winfo_width()
            self._resize_starth = win.winfo_height()

        def _do_resize(e):
            if not (self.active_result_win and tk.Toplevel.winfo_exists(self.active_result_win)):
                return
            new_w = max(_MIN_W, self._resize_startw + (e.x_root - self._resize_startx))
            new_h = max(_MIN_H, self._resize_starth + (e.y_root - self._resize_starty))
            self.active_result_win.geometry(f"{new_w}x{new_h}")

        resize_grip.bind("<Button-1>", _start_resize)
        resize_grip.bind("<B1-Motion>", _do_resize)
        resize_grip.lift()  # keep the handle clickable above the Notebook

        self.nb = ttk.Notebook(self.results_frame); self.nb.pack(fill="both", expand=True)
        only_files = [r for r in file_res if str(r[0]).lower() != "folder"]
        only_folders = [r for r in file_res if str(r[0]).lower() == "folder"]
        self._all_files_data = only_files
        self._all_content_data = cont_res
        self._ai_cont_res  = []
        self._ai_file_res  = []
        self.size_op_var.set(">"); self.size_num_var.set(""); self.size_unit_var.set("MB")
        self.ext_filter_var.set("")
        self.c_size_op_var.set(">"); self.c_size_num_var.set(""); self.c_size_unit_var.set("MB")
        self.c_ext_filter_var.set("")
        self.name_filter_var.set("")
        self.c_name_filter_var.set("")

        # v2.3: Tab order — File Name (1st), Folder Name (2nd), File Content (last/right)
        # Tabs are added to notebook in correct order below via nb.add
        # File Name + Folder Name: MFT realtime (no DB needed)
        # File Content: BM25/AI results, populated after --update data

        # 1. TAB CONTENT (frame created here, added to nb LAST after f/fol frames below)
        c_frame = tk.Frame(self.nb)

        # ── Content Filter bar ─────────────────────────────────────────────
        c_filter_bar = tk.Frame(c_frame, bg=BG_COLOR, pady=2)
        c_filter_bar.pack(fill="x", padx=6, pady=(4, 2))

        tk.Label(c_filter_bar, text="Size:", bg=BG_COLOR, fg=TEXT_COLOR,
                 font=("Segoe UI", 9)).pack(side="left", padx=(4, 2))

        c_op_menu = ttk.OptionMenu(c_filter_bar, self.c_size_op_var, self.c_size_op_var.get(),
                                    "Any", ">", ">=", "<", "<=", "=",
                                    command=lambda *_: self._apply_content_filter_to_tree())
        c_op_menu.config(width=4); c_op_menu.pack(side="left", padx=2)

        c_size_entry = tk.Entry(c_filter_bar, textvariable=self.c_size_num_var, width=7,
                                 font=("Segoe UI", 9), bg=ENTRY_BG, fg=TEXT_COLOR,
                                 insertbackground=TEXT_COLOR, bd=1, relief="flat")
        c_size_entry.pack(side="left", padx=2)
        self.c_size_num_var.trace_add("write", self._apply_content_filter_to_tree)

        c_unit_menu = ttk.OptionMenu(c_filter_bar, self.c_size_unit_var, self.c_size_unit_var.get(),
                                      "B", "KB", "MB", "GB",
                                      command=lambda *_: self._apply_content_filter_to_tree())
        c_unit_menu.config(width=4); c_unit_menu.pack(side="left", padx=2)

        # ── Content Ext filter ─────────────────────────────────────────────
        tk.Label(c_filter_bar, text="│", bg=BG_COLOR, fg="#555",
                 font=("Segoe UI", 11)).pack(side="left", padx=4)
        tk.Label(c_filter_bar, text="Ext:", bg=BG_COLOR, fg=TEXT_COLOR,
                 font=("Segoe UI", 9)).pack(side="left", padx=(0, 3))

        _C_EXT_PRESETS = [
            ("All",   ""),
            ("PDF",   "pdf"),
            ("Word",  "doc,docx"),
            ("Excel", "xls,xlsx,csv"),
            ("PPT",   "ppt,pptx"),
            ("Txt",   "txt,log"),
        ]
        def _set_c_ext(val):
            self.c_ext_filter_var.set(val)
        for _lbl, _val in _C_EXT_PRESETS:
            tk.Button(c_filter_bar, text=_lbl, font=("Segoe UI", 8),
                      bg="#e6e8ec", fg="#33363c", bd=0, padx=5, pady=1,
                      activebackground="#d3d6dc", cursor="hand2",
                      command=lambda v=_val: _set_c_ext(v)).pack(side="left", padx=1)

        tk.Label(c_filter_bar, text="│", bg=BG_COLOR, fg="#555",
                 font=("Segoe UI", 9)).pack(side="left", padx=3)

        c_ext_entry = tk.Entry(c_filter_bar, textvariable=self.c_ext_filter_var, width=14,
                                font=("Segoe UI", 9), bg=ENTRY_BG, fg="#888",
                                insertbackground=TEXT_COLOR, bd=1, relief="flat")
        c_ext_entry.pack(side="left", padx=(0, 2))
        def _c_ext_focus_in(e):
            if not self.c_ext_filter_var.get(): c_ext_entry.config(fg=TEXT_COLOR)
        def _c_ext_focus_out(e):
            if not self.c_ext_filter_var.get(): c_ext_entry.config(fg="#888")
        c_ext_entry.bind("<FocusIn>",  _c_ext_focus_in)
        c_ext_entry.bind("<FocusOut>", _c_ext_focus_out)
        self.c_ext_filter_var.trace_add("write", self._apply_content_filter_to_tree)

        tk.Button(c_filter_bar, text="✕", command=lambda: self.c_ext_filter_var.set(""),
                  bg="#c9ccd2", fg="#222222", font=("Segoe UI", 8), bd=0, padx=4,
                  activebackground="#b0b3ba", cursor="hand2").pack(side="left", padx=(0, 4))

        # ── Content Name filter ────────────────────────────────────────────
        tk.Label(c_filter_bar, text="│", bg=BG_COLOR, fg="#555",
                 font=("Segoe UI", 9)).pack(side="left", padx=3)
        tk.Label(c_filter_bar, text="Name:", bg=BG_COLOR, fg=TEXT_COLOR,
                 font=("Segoe UI", 9)).pack(side="left", padx=(0, 2))
        c_name_entry = tk.Entry(c_filter_bar, textvariable=self.c_name_filter_var, width=18,
                                 font=("Segoe UI", 9), bg=ENTRY_BG, fg="#888",
                                 insertbackground=TEXT_COLOR, bd=1, relief="flat")
        c_name_entry.pack(side="left", padx=(0, 2))
        def _c_name_focus_in(e):
            if not self.c_name_filter_var.get(): c_name_entry.config(fg=TEXT_COLOR)
        def _c_name_focus_out(e):
            if not self.c_name_filter_var.get(): c_name_entry.config(fg="#888")
        c_name_entry.bind("<FocusIn>",  _c_name_focus_in)
        c_name_entry.bind("<FocusOut>", _c_name_focus_out)
        self.c_name_filter_var.trace_add("write", self._apply_content_filter_to_tree)
        tk.Button(c_filter_bar, text="✕", command=lambda: self.c_name_filter_var.set(""),
                  bg="#c9ccd2", fg="#222222", font=("Segoe UI", 8), bd=0, padx=4,
                  activebackground="#b0b3ba", cursor="hand2").pack(side="left", padx=(0, 4))
        # ──────────────────────────────────────────────────────────────────

        self.content_filter_count_label = None  # removed

        # ── AI Search button ───────────────────────────────────────────────────
        _ai_btn_state = "normal"
        self._hybrid_status_lbl = None  # removed from UI

        def _on_ai_search():
            q_now = self._last_query or query
            if not q_now: return
            # Toggle: if currently in AI mode, restore BM25 (handled inside _ai_search_and_update)
            # Disable button immediately to prevent double-click
            try:
                if not self._ai_mode_active:
                    self._ai_search_btn.config(state="disabled", text="⏳ Running...",
                                                fg="#ffcc00", bg="#2a2a1a")
            except Exception: pass
            threading.Thread(target=self._ai_search_and_update, args=(q_now,), daemon=True).start()

        # ── AI Model selector: dropdown to pick one of 3 models ────────────
        # jina_v3 / bge_gemma2 — each model has different embedding dimension/
        # vector space so semantic data is stored in separate tables in DB
        # (semantic_index / semantic_index_jina_v3 / semantic_index_bge_gemma2).
        # Switching does not delete other model data — just run --update data
        # ONCE per model, then switch freely without rebuilding.
        _model_keys_order = ["jina_v3", "bge_gemma2"]
        _model_display = [SEMANTIC_MODELS[k]["label"] for k in _model_keys_order]
        _display_to_key = {SEMANTIC_MODELS[k]["label"]: k for k in _model_keys_order}

        def _check_table_has_data(model_key):
            try:
                table = SEMANTIC_MODELS.get(model_key, {}).get("table", "semantic_index")
                conn = sqlite3.connect(DB_FILE, timeout=5)
                c = conn.cursor()
                c.execute("SELECT count(*) FROM sqlite_master WHERE type='table' AND name=?", (table,))
                has_table = c.fetchone()[0] > 0
                count = 0
                if has_table:
                    c.execute(f"SELECT count(*) FROM {table}")
                    count = c.fetchone()[0]
                conn.close()
                return count
            except Exception:
                return 0

        def _refresh_ai_status_lbl(model_key, ok=None):
            label = SEMANTIC_MODELS.get(model_key, {}).get("label", model_key)
            data_count = _check_table_has_data(model_key)
            try:
                if ok is False:
                    self._ai_model_status_lbl.config(text=f"✗ {label}: load error", fg="#ff6666")
                else:
                    # v5.8: dropped "✓ <model>: ready" — the ramp light going
                    # Blue already communicates "ready" for the current
                    # model, so this label stayed silent otherwise (and for
                    # a genuinely-not-ready model, the dropdown/AI-Search/
                    # Advanced buttons already grey out via _sync_ai_adv_lock).
                    self._ai_model_status_lbl.config(text="", fg="#888888")
            except Exception:
                pass
            return data_count

        def _on_ai_model_change(event=None):
            global _sem_model_key
            new_key = _display_to_key.get(self._ai_model_combo.get(), DEFAULT_SEMANTIC_MODEL)
            self.ai_model_var.set(new_key)
            _sem_model_key = new_key

            data_count = _check_table_has_data(new_key)
            label = SEMANTIC_MODELS[new_key]["label"]
            if data_count == 0:
                # No embedding for this model yet — inform user to run --update data once
                messagebox.showinfo(
                    "AI Model Switch",
                    f"Switched to model: {label}\n\n"
                    f"This model has NO semantic data in search_data.db yet.\n"
                    f"Run \"--update data\" ONCE to build the AI index for this model.\n\n"
                    f"After that, you can switch freely between all 3 models "
                    f"without rebuilding (unless you re-scan all drives)."
                )
            # Load new model in background immediately so the next AI search has no delay.
            # jina-v3/bge-gemma2 on CPU can take a long time (especially bge-gemma2, ~9B params)
            # so show loading status clearly so the user knows it's loading, not frozen.
            def _bg_load():
                ok = _load_semantic_model(new_key)
                self.root.after(0, lambda: _refresh_ai_status_lbl(new_key, ok))
                # v5.8: if AI Search results are currently on screen, re-run
                # the semantic search with the newly selected model instead
                # of silently leaving the OLD model's results displayed.
                # Without this, switching Jina v3 -> BGE Gemma2 kept showing
                # Jina's hits/count until the user manually clicked BM25
                # then AI Search again -- calling _ai_search_and_update
                # directly here would just TOGGLE IT OFF (that method treats
                # a second call while _ai_mode_active as "restore BM25"), so
                # _ai_mode_active is cleared first to force a fresh run.
                if ok and data_count and self._ai_mode_active:
                    q_now = self._last_query or query
                    if q_now:
                        def _rerun_btn_state():
                            try:
                                self._ai_search_btn.config(state="disabled", text="⏳ Running...",
                                                            fg="#ffcc00", bg="#2a2a1a")
                            except Exception: pass
                        self.root.after(0, _rerun_btn_state)
                        self._ai_mode_active = False
                        self._ai_search_and_update(q_now)
            self._ai_model_status_lbl.config(text=f"⏳ loading {label}...", fg="#ffcc00")
            threading.Thread(target=_bg_load, daemon=True).start()

        # v2.6: moved onto the top search_bar (same row as the Searchbox) instead of
        # the per-tab c_filter_bar, per user request — stays visible on every tab and
        # no longer eats vertical space above the Content tab's results tree.
        ai_model_frame = tk.Frame(search_bar, bg=BG_COLOR)
        ai_model_frame.pack(side="right", padx=(2, 4))
        self._ai_model_combo = ttk.Combobox(
            ai_model_frame, values=_model_display, state="disabled",
            width=16, font=("Segoe UI", 8))
        _cur_key = self.ai_model_var.get() if self.ai_model_var.get() in SEMANTIC_MODELS else DEFAULT_SEMANTIC_MODEL
        self._ai_model_combo.set(SEMANTIC_MODELS[_cur_key]["label"])
        self._ai_model_combo.pack(side="left")
        self._ai_model_combo.bind("<<ComboboxSelected>>", _on_ai_model_change)

        self._ai_model_status_lbl = tk.Label(
            ai_model_frame, text="", font=("Segoe UI", 7, "italic"),
            bg=BG_COLOR, fg="#888888")
        self._ai_model_status_lbl.pack(side="left", padx=(4, 0))
        # Show data availability status for the current model when the result window opens
        _refresh_ai_status_lbl(_cur_key)
        # ──────────────────────────────────────────────────────────────────

        self._ai_search_btn = tk.Button(
            search_bar, text="🤖 AI Search",
            font=("Segoe UI", 8, "bold"),
            bg="#1e3a5f", fg="#7ec8e3",
            relief="raised", bd=2,
            padx=8, pady=2,
            activebackground="#2a5080", activeforeground="#a8d8f0",
            cursor="hand2",
            state="disabled",
            command=_on_ai_search)
        self._ai_search_btn.pack(side="right", padx=(4, 2), pady=4)
        def _ai_search_btn_tooltip_text():
            try:
                if self._ai_search_btn and self._ai_search_btn.winfo_exists() \
                        and self._ai_search_btn.cget("text").strip().startswith("↩"):
                    return "Currently showing AI results — click to switch back to BM25 (faster, keyword-based)"
            except Exception:
                pass
            return "AI-powered semantic search (slower, smarter)"
        add_tooltip(self._ai_search_btn, _ai_search_btn_tooltip_text)

        # ── Advanced button (pagination) ───────────────────────────────────
        # Page 0 = realtime top 100. Each click appends next 500. Last page → reset to page 0.
        PAGE_SIZE = 500

        def _adv_label():
            page = self._adv_page
            total = len(self._adv_all_cont) if hasattr(self, '_adv_all_cont') else 0
            total_f = len(self._adv_all_files) if hasattr(self, '_adv_all_files') else 0
            if page == 0:
                return "Advanced"
            return "Simple ↩"

        def _on_advanced():
            q_now = self._last_query or query
            if not q_now: return

            # First click on page 0: run full search to get all results
            if self._adv_page == 0:
                try:
                    self._adv_search_btn.config(state="disabled", text="⏳ Loading...", fg="#ffcc00", bg="#2a2a1a")
                except Exception: pass

                def _run_full_search():
                    # Run search with adv_mode=True to get full result set
                    import sqlite3 as _sq3
                    try:
                        conn = self.db_conn
                        if conn is None: return
                        c2 = conn.cursor()
                        from rank_bm25 import BM25Okapi
                        import re as _re2

                        cleaned_q2, op2, size_val2 = parse_size_filter(q_now)
                        kw2 = [k.lower() for k in cleaned_q2.split() if k]
                        kw2 = _strip_stopwords(kw2)  # v5.8: see _smart_search_realtime comment
                        size_clause2 = f" AND size {op2} ?" if op2 else ""
                        size_extra2  = [size_val2] if op2 else []
                        has_anchor2  = any(_is_anchor_kw(k) for k in kw2)

                        # File Name full search
                        file_res2 = []
                        if has_anchor2:
                            nc = ["name LIKE ?" for k in kw2]
                            np2 = [f"%{k}%" for k in kw2]
                            c2.execute("SELECT type, name, path, size FROM files WHERE type != 'Folder' AND " + " AND ".join(nc) + size_clause2 + " LIMIT 5000", np2 + size_extra2)
                            fn2 = c2.fetchall()
                            if not fn2 and len(kw2) > 1:
                                c2.execute("SELECT type, name, path, size FROM files WHERE type != 'Folder' AND (" + " OR ".join(nc) + ")" + size_clause2 + " LIMIT 5000", np2 + size_extra2)
                                fn2 = c2.fetchall()
                            c2.execute("SELECT type, name, path, size FROM files WHERE type = 'Folder' AND " + " AND ".join(nc) + size_clause2 + " LIMIT 5000", np2 + size_extra2)
                            file_res2 = fn2 + c2.fetchall()
                            # v7.10: same "Whole word" post-filter as the main
                            # realtime DB path -- SQL LIKE can't express word
                            # boundaries, so filter the substring hits in Python.
                            if self.whole_word_var.get():
                                file_res2 = [r for r in file_res2
                                              if all(_kw_matches(k, r[1].lower(), True) for k in kw2)]

                        # Content full search via FTS
                        import re as _re3
                        _STOP2 = {'a','an','the','is','in','on','at','to','of','or','and','as','be','by','do','for','has','had','he','her','him','his','how','i','if','it','its','me','my','no','not','off','our','out','own','so','than','that','them','then','they','this','us','was','we','who','why','will','with','you','your','also','been','but','can','did','does','from','get','got','have','into','just','may','new','now','one','see','set','she','time','what','when','which','would'}
                        FTS5_OPS = set('+-*:^"()：、。・<>@[]{}|\\/?!#$%&=~`\'')
                        # NOTE: same trigram-tokenizer limitation as the realtime search —
                        # FTS terms under 3 chars silently match zero rows, so this stays
                        # at len>=3; short CJK anchors fall through to the LIKE path below.
                        def _safe(k): return len(k)>=3 and not any(ch in FTS5_OPS for ch in k) and '-' not in k and not _re3.search(r'[^\w　-鿿＀-￯一-鿿]', k)
                        phrase2 = cleaned_q2.strip()
                        cont_res2 = []
                        if _is_anchor_kw(phrase2):
                            pt2 = [k for k in _re3.split(r'[\s:+\-<>@"()#.\[\]{}|\/?!&=~`]+', phrase2) if k]
                            fts_tok2 = [k for k in pt2 if len(k)>=3 and k.lower() not in _STOP2 and _safe(k)]
                            if fts_tok2:
                                fts_q2 = " AND ".join(f'"{k}"' for k in fts_tok2)
                                try:
                                    c2.execute("SELECT path FROM content_index WHERE content MATCH ? LIMIT 5000", (fts_q2,))
                                    cands2 = set(r[0] for r in c2.fetchall())
                                    if cands2:
                                        ph2 = ",".join("?"*len(cands2))
                                        c2.execute(f"SELECT f.path, f.size FROM files f WHERE f.path IN ({ph2})", list(cands2))
                                        cont_res2 = c2.fetchall()
                                except Exception: pass
                            else:
                                # No token was usable by the FTS trigram index (e.g. a
                                # short CJK anchor like 解析 — trigram needs >= 3 chars
                                # and would otherwise silently return zero results).
                                # Fall back to a direct LIKE scan on content_store —
                                # slower (full table scan) but has no length floor.
                                try:
                                    c2.execute("SELECT path FROM content_store WHERE content LIKE ? LIMIT 5000", (f"%{phrase2}%",))
                                    cands2 = set(r[0] for r in c2.fetchall())
                                    if cands2:
                                        ph2 = ",".join("?"*len(cands2))
                                        c2.execute(f"SELECT f.path, f.size FROM files f WHERE f.path IN ({ph2})", list(cands2))
                                        cont_res2 = c2.fetchall()
                                except Exception: pass

                        # BM25 rank file results
                        def _tok2(t): return [x.lower() for x in _re2.split(r'[\s\W]+', str(t)) if len(x)>=2]
                        qt2 = _tok2(cleaned_q2)
                        if file_res2 and qt2:
                            try:
                                corp_f = [_tok2(os.path.basename(r[2])+" "+r[2]) for r in file_res2]
                                b2f = BM25Okapi(corp_f)
                                sc_f = b2f.get_scores(qt2)
                                file_res2 = [file_res2[i] for i in sorted(range(len(file_res2)), key=lambda i: sc_f[i], reverse=True)]
                            except Exception: pass

                        # BM25 rank content results by filename
                        bm25_c2 = {}
                        sz_map2 = {r[0]: r[1] for r in cont_res2}
                        if cont_res2 and qt2:
                            try:
                                po2, corp_c2 = [], []
                                for p in [r[0] for r in cont_res2]:
                                    fn = os.path.basename(p); par = os.path.basename(os.path.dirname(p)); gp = os.path.basename(os.path.dirname(os.path.dirname(p)))
                                    po2.append(p); corp_c2.append(_tok2(f"{fn} {par} {gp}"))
                                b2c = BM25Okapi(corp_c2)
                                sc_c = b2c.get_scores(qt2)
                                mx2 = max(sc_c) if max(sc_c) > 0 else 1.0
                                bm25_c2 = {po2[i]: sc_c[i]/mx2 for i in range(len(po2))}
                            except Exception: pass

                        hyb2 = sorted([(p, sz_map2.get(p,0), bm25_c2.get(p,0.0)) for p in sz_map2], key=lambda x: x[2], reverse=True)
                        mx_h2 = hyb2[0][2] if hyb2 else 1.0
                        cont_res2_scored = [(p, sz, int((sc/mx_h2)*99) if mx_h2>0 else 0) for p,sz,sc in hyb2]

                        # Priority sort: office/pdf/msg first, txt/md middle, log/html/code/... last
                        cont_res2_scored = self._sort_priority(cont_res2_scored, 0)
                        file_res2        = self._sort_priority(file_res2, 1)

                        # Store full results
                        self._adv_all_cont  = cont_res2_scored
                        self._adv_all_files = file_res2

                        def _apply_page1():
                            self._adv_page = 1
                            # Display ALL results immediately (no pagination)
                            all_cont  = self._adv_all_cont
                            all_files = self._adv_all_files
                            self._all_content_data = all_cont
                            self._all_files_data   = [r for r in all_files if str(r[0]).lower() != "folder"]
                            only_all_f   = [r for r in all_files if str(r[0]).lower() != "folder"]
                            only_all_fol = [r for r in all_files if str(r[0]).lower() == "folder"]
                            # Clear current tree and reload everything
                            for item in self.tree_c.get_children():   self.tree_c.delete(item)
                            for item in self.tree_f.get_children():   self.tree_f.delete(item)
                            for item in self.tree_fol.get_children(): self.tree_fol.delete(item)
                            for i, item in enumerate(all_cont):
                                try: self._insert_row(self.tree_c, item, i+1, 0)
                                except: pass
                            for i, item in enumerate(only_all_f):
                                try: self._insert_row(self.tree_f, item, i+1, 1)
                                except: pass
                            for i, item in enumerate(only_all_fol):
                                try: self._insert_row(self.tree_fol, item, i+1, 2)
                                except: pass
                            try:
                                if self._adv_search_btn and self._adv_search_btn.winfo_exists():
                                    self._adv_search_btn.config(state="normal", text=_adv_label(), fg="#90ee90", bg="#1a3a1a")
                                if self.content_filter_count_label:
                                    self.content_filter_count_label.config(text=f"{len(all_cont)} files")
                                if self.filter_count_label:
                                    self.filter_count_label.config(text=f"{len(only_all_f)} files")
                            except Exception: pass
                            # Open Advanced split INSIDE tab
                            self._adv_mode_active = True
                            self.root.after(0, self._open_adv_split_win)
                        self.root.after(0, _apply_page1)
                    except Exception as _e2:
                        print(f"[Advanced] Error: {_e2}")
                        try:
                            self.root.after(0, lambda: self._adv_search_btn.config(state="normal", text="Advanced", fg="#33363c", bg="#e6e8ec"))
                        except: pass

                threading.Thread(target=_run_full_search, daemon=True).start()
                return

            # Second click: reset to realtime (page 0) — show "Simple" mode
            self._adv_page = 0
            self._adv_mode_active = False
            # Restore single tree in File Content tab
            self._close_adv_split()
            # Clear trees and reload realtime results only
            rt_c  = self._last_bm25_cont_res or []
            rt_f  = self._last_bm25_file_res or []
            self._all_content_data = rt_c
            self._all_files_data   = [r for r in rt_f if str(r[0]).lower() != "folder"]
            only_rt_f   = [r for r in rt_f if str(r[0]).lower() != "folder"]
            only_rt_fol = [r for r in rt_f if str(r[0]).lower() == "folder"]
            for item in self.tree_c.get_children():   self.tree_c.delete(item)
            for item in self.tree_f.get_children():   self.tree_f.delete(item)
            for item in self.tree_fol.get_children(): self.tree_fol.delete(item)
            for i, item in enumerate(rt_c):
                try: self._insert_row(self.tree_c, item, i+1, 0)
                except: pass
            for i, item in enumerate(only_rt_f):
                try: self._insert_row(self.tree_f, item, i+1, 1)
                except: pass
            for i, item in enumerate(only_rt_fol):
                try: self._insert_row(self.tree_fol, item, i+1, 2)
                except: pass
            try:
                self._adv_search_btn.config(text="Advanced", fg="#33363c", bg="#e6e8ec", state="normal")
                if self.content_filter_count_label:
                    self.content_filter_count_label.config(text=f"{len(rt_c)} files")
                if self.filter_count_label:
                    self.filter_count_label.config(text=f"{len(only_rt_f)} files")
            except Exception: pass
            return

        self._adv_search_btn = tk.Button(
            search_bar, text="Advanced",
            font=("Segoe UI", 8, "bold"),
            bg="#e6e8ec", fg="#33363c",
            relief="raised", bd=2,
            padx=8, pady=2,
            activebackground="#d3d6dc", activeforeground="#111111",
            cursor="hand2",
            command=_on_advanced)
        self._adv_search_btn.pack(side="right", padx=(4, 2), pady=4)
        add_tooltip(self._adv_search_btn,
                    lambda: ("Back to simple results (top matches only)"
                             if getattr(self, "_adv_mode_active", False)
                             else "Show more results (paginate through all matches)"))
        # Apply current ramp-light status now that Advanced / AI Search / model
        # combobox all exist (they're rebuilt fresh every time show_results runs).
        self._sync_ai_adv_lock()
        # ──────────────────────────────────────────────────────────────────

        # ── Update DB button (same action as typing "--update data") ────────
        # Packed AFTER Advanced so it lands to Advanced's left — i.e. right next
        # to the Searchbox, between the Searchbox and the Advanced button.
        # v3.4: also reads the Searchbox for a bare tier expression — type
        # "tier 1,2" (or just "1,2") then click this button instead of typing
        # the full "--update data tier 1,2" command and pressing Enter.
        # If the box is empty or holds an ordinary search query, behaves
        # exactly as before (full Tier 1-4 scan).
        def _run_update_db(selected_tiers, selected_models=None, ocr_enabled=False):
            """Actually kick off the indexing_worker thread. Extracted out of
            the old _on_update_db so both the new dialog's 'Start Update'
            button and (if ever needed again) any other caller can trigger a
            scan the same way."""
            if self._update_db_running:
                return
            tier_desc = ("ALL (1-4)" if selected_tiers is None
                         else ",".join(str(t + 1) for t in sorted(selected_tiers)))
            try:
                self._update_db_btn.config(state="disabled", text="Updating...",
                                            fg="#ffcc00", bg="#2a2a1a")
            except Exception:
                pass
            self._update_db_running = True  # lock AI Search/model combo immediately, don't wait for the thread
            self._set_index_status(f"Updating (tier {tier_desc})...", "#ffcc00")
            self._ramp_blink_start()
            self.entry_var.set("")
            threading.Thread(target=self.indexing_worker,
                              args=(selected_tiers, selected_models, ocr_enabled), daemon=True).start()

        if not hasattr(self, "_model_abort_events"):
            self._model_abort_events = {}  # model_key -> threading.Event(), set on Abort click

        def _install_model_online(model_key, status_lbl, install_btn, dialog, cb=None, mv=None):
            """Download a model's weights from HuggingFace Hub into
            models/<dir_name>/ next to the app. Requires internet + the
            huggingface_hub package. Runs in a background thread so the
            dialog stays responsive; UI updates are marshalled back via
            self.root.after(0, ...) since Tk isn't thread-safe.
            v5.9: while downloading, the button becomes "Abort" (rather than
            just greyed out) — huggingface_hub's snapshot_download can't be
            killed mid-transfer from here, but clicking Abort immediately
            resets the button/status and flags the in-flight download as
            abandoned, so its eventual result (success or error) is silently
            ignored instead of overwriting the UI once it finally returns."""
            info = SEMANTIC_MODELS.get(model_key, {})
            label = info.get("label", model_key)
            repo = info.get("hf_repo")
            if not repo:
                messagebox.showerror("Install model", f"No download source configured for {label}.", parent=dialog)
                return
            if not messagebox.askyesno(
                "Install model — internet required",
                f"Download {label} now from HuggingFace ({repo})?\n\n"
                f"This requires internet access and may take a while /\n"
                f"use significant disk space (several GB for larger models).",
                parent=dialog):
                dialog.lift()
                return
            dest_dir = os.path.join(_MODEL_ROOT_CANDIDATES[0], info["dir_names"][0])
            abort_event = threading.Event()
            self._model_abort_events[model_key] = abort_event

            def _do_abort():
                abort_event.set()
                try:
                    status_lbl.config(text="cancelled", fg="#888888")
                    install_btn.config(state="normal", text="Install online...",
                                        command=lambda: _install_model_online(model_key, status_lbl, install_btn, dialog, cb, mv))
                except Exception:
                    pass

            try:
                install_btn.config(state="normal", text="Abort", command=_do_abort)
                status_lbl.config(text="downloading...", fg="#ffcc00")
            except Exception:
                pass

            def _bg():
                try:
                    from huggingface_hub import snapshot_download
                except ImportError:
                    def _fail():
                        if abort_event.is_set():
                            return
                        status_lbl.config(text="huggingface_hub not installed", fg="#ff6666")
                        install_btn.config(state="normal", text="Install online...",
                                            command=lambda: _install_model_online(model_key, status_lbl, install_btn, dialog, cb, mv))
                        messagebox.showerror(
                            "Install model",
                            "The 'huggingface_hub' package isn't installed.\n"
                            "Run: pip install huggingface_hub --break-system-packages",
                            parent=dialog)
                        dialog.lift()
                    self.root.after(0, _fail)
                    return
                try:
                    os.makedirs(dest_dir, exist_ok=True)
                    snapshot_download(repo_id=repo, local_dir=dest_dir)
                    if abort_event.is_set():
                        return  # user aborted meanwhile — UI already reset, don't touch it
                    def _ok():
                        SEMANTIC_MODEL_DIRS[model_key] = dest_dir
                        status_lbl.config(text="✓ installed", fg="#4caf50")
                        install_btn.config(state="normal", text="Reinstall...",
                                           command=lambda: _install_model_online(model_key, status_lbl, install_btn, dialog, cb, mv))
                        # v7.10: model just became installed -- un-grey the
                        # checkbox and tick it by default so the user doesn't
                        # have to close/reopen this dialog to use it this run.
                        if cb is not None:
                            try: cb.config(state="normal")
                            except Exception: pass
                        if mv is not None:
                            try: mv.set(True)
                            except Exception: pass
                        messagebox.showinfo("Install model", f"{label} installed successfully.", parent=dialog)
                        dialog.lift()
                    self.root.after(0, _ok)
                except Exception as e:
                    if abort_event.is_set():
                        return  # aborted — the exception is just the interrupted transfer, ignore it
                    err = str(e)
                    def _err():
                        status_lbl.config(text="download failed", fg="#ff6666")
                        install_btn.config(state="normal", text="Install online...",
                                            command=lambda: _install_model_online(model_key, status_lbl, install_btn, dialog, cb, mv))
                        messagebox.showerror("Install model", f"Download failed for {label}:\n{err}", parent=dialog)
                        dialog.lift()
                    self.root.after(0, _err)

            threading.Thread(target=_bg, daemon=True).start()

        def _install_vi_diacritics_online(status_lbl, install_btn, dialog):
            """Download the Vietnamese diacritics-restoration model (base
            model + LoRA adapter, two separate HF repos) into
            models/vi-diacritics/{base,adapter}/. Same safe
            snapshot_download-into-a-folder mechanism as
            _install_model_online, just for 2 repos instead of 1, and with
            ignore_patterns to skip redundant TensorFlow/Flax/ONNX weight
            copies that vinai/bartpho-syllable ships alongside the PyTorch
            ones we actually use (this is what made an earlier bundled-into-
            the-exe attempt balloon to several GB for no benefit)."""
            if not messagebox.askyesno(
                "Install Vietnamese diacritics restoration — internet required",
                "Download the Vietnamese diacritics-restoration model now?\n\n"
                "Improves AI Search when queries are typed without dấu "
                "(e.g. \"he thong kiem soat\" instead of \"hệ thống kiểm soát\").\n\n"
                "Requires internet access, ~1-2GB disk space.",
                parent=dialog):
                dialog.lift()
                return
            local_root = _vi_dia_local_dir()
            abort_event = threading.Event()

            def _do_abort():
                abort_event.set()
                try:
                    status_lbl.config(text="cancelled", fg="#888888")
                    install_btn.config(state="normal", text="Install online...",
                                        command=lambda: _install_vi_diacritics_online(status_lbl, install_btn, dialog))
                except Exception:
                    pass

            try:
                install_btn.config(state="normal", text="Abort", command=_do_abort)
                status_lbl.config(text="downloading...", fg="#ffcc00")
            except Exception:
                pass

            def _bg():
                try:
                    from huggingface_hub import snapshot_download
                except ImportError:
                    def _fail():
                        if abort_event.is_set():
                            return
                        status_lbl.config(text="huggingface_hub not installed", fg="#ff6666")
                        install_btn.config(state="normal", text="Install online...",
                                            command=lambda: _install_vi_diacritics_online(status_lbl, install_btn, dialog))
                        messagebox.showerror(
                            "Install model",
                            "The 'huggingface_hub' package isn't installed.\n"
                            "Run: pip install huggingface_hub --break-system-packages",
                            parent=dialog)
                        dialog.lift()
                    self.root.after(0, _fail)
                    return
                try:
                    base_dest = os.path.join(local_root, "base")
                    adapter_dest = os.path.join(local_root, "adapter")
                    os.makedirs(base_dest, exist_ok=True)
                    os.makedirs(adapter_dest, exist_ok=True)
                    snapshot_download(repo_id=VI_DIA_BASE_REPO, local_dir=base_dest,
                                       ignore_patterns=VI_DIA_IGNORE_PATTERNS)
                    if abort_event.is_set():
                        return
                    snapshot_download(repo_id=VI_DIA_ADAPTER_REPO, local_dir=adapter_dest,
                                       ignore_patterns=VI_DIA_IGNORE_PATTERNS)
                    if abort_event.is_set():
                        return
                    def _ok():
                        status_lbl.config(text="✓ installed", fg="#4caf50")
                        install_btn.config(state="normal", text="Reinstall...",
                                           command=lambda: _install_vi_diacritics_online(status_lbl, install_btn, dialog))
                        dialog.lift()
                    self.root.after(0, _ok)
                except Exception as e:
                    if abort_event.is_set():
                        return
                    err = str(e)
                    def _err():
                        status_lbl.config(text="download failed", fg="#ff6666")
                        install_btn.config(state="normal", text="Install online...",
                                            command=lambda: _install_vi_diacritics_online(status_lbl, install_btn, dialog))
                        messagebox.showerror("Install model", f"Download failed:\n{err}", parent=dialog)
                        dialog.lift()
                    self.root.after(0, _err)

            threading.Thread(target=_bg, daemon=True).start()

        def _on_update_db():
            """v5.9: clicking Update DB now opens an options dialog instead
            of scanning immediately — pick which file-type tiers to index,
            check/install AI Search models (Jina-v3 / BGE-Gemma2), and pick
            which model(s) actually get embeddings built this run."""
            if self._update_db_running:
                messagebox.showinfo("Update DB", "An update is already running.")
                return

            # Pre-fill tier checkboxes from a tier expression already typed
            # in the Searchbox (e.g. "tier 1,2") — preserves the old quick-
            # type shortcut as a convenience default; otherwise only Tier 1
            # starts checked (fastest default: primary office docs/PDF).
            raw_box_text = self.entry_var.get().strip()
            pre_tiers, is_tier_expr = self._parse_tier_filter(raw_box_text)
            default_on = set(pre_tiers) if (raw_box_text and is_tier_expr) else {0}

            dlg = tk.Toplevel(self.root)
            dlg.title("Update DB — Options")
            dlg.configure(bg=BG_COLOR)
            dlg.transient(self.root)
            # v5.9b: NOT topmost — this was actively unwanted (it was staying
            # pinned above unrelated apps like Chrome). A normal Toplevel
            # behaves like any other window: on top of self.root initially,
            # but can be covered by whatever the user clicks next, same as
            # any other dialog in the app.
            dlg.resizable(False, False)
            # placeholder position -- actually centered on self.root once every
            # widget below is packed (see the centering block near the bottom,
            # right after the Start/Cancel buttons)
            dlg.geometry(f"+{self.root.winfo_x()}+{self.root.winfo_y() + 40}")

            # ── Tier 1-4 file-type checkboxes ───────────────────────────────
            tiers_f = tk.LabelFrame(dlg, text="File types to index (Tiers)",
                                     bg=BG_COLOR, fg="#33363c",
                                     font=("Segoe UI", 9, "bold"), padx=10, pady=8)
            tiers_f.pack(fill="x", padx=12, pady=(12, 6))

            tier_labels = ["Tier 1  (Office / PDF)", "Tier 2  (Email / Text / Notes)",
                           "Tier 3  (Scripts / Config)", "Tier 4  (Markup / Query / Misc)"]
            tier_var_list = []
            for i, exts in enumerate(self._ALL_TIERS):
                v = tk.BooleanVar(value=(i in default_on))
                tier_var_list.append(v)
                ext_str = ", ".join(sorted(exts))
                cb = tk.Checkbutton(tiers_f, text=f"{tier_labels[i]}: {ext_str}",
                                     variable=v, bg=BG_COLOR, anchor="w",
                                     font=("Segoe UI", 9), justify="left",
                                     wraplength=460)
                cb.pack(fill="x", anchor="w", pady=2)

            # v9.13: OCR images checkbox — off by default. Separate from the
            # Tier checkboxes above (images don't belong to any of the 4
            # Tiers), but placed in the same frame since it's conceptually
            # the same kind of choice: "what content gets extracted".
            ocr_var = tk.BooleanVar(value=False)
            ocr_cb = tk.Checkbutton(
                tiers_f, text="OCR images (.jpg/.png/...) — slower, downloads OCR model on first use",
                variable=ocr_var, bg=BG_COLOR, anchor="w",
                font=("Segoe UI", 9), justify="left", wraplength=460)
            ocr_cb.pack(fill="x", anchor="w", pady=(6, 2))

            # ── AI Search models section (separate frame, same dialog) ─────
            # v5.9b: each model now has its own "build embeddings this run"
            # checkbox (checked by default = same as before/all models) --
            # unchecking a model skips its (slow) embedding pass entirely,
            # e.g. if the user only ever uses Jina-v3 and doesn't care about
            # BGE-Gemma2, unchecking it noticeably speeds up Update DB.
            ai_f = tk.LabelFrame(dlg, text="Search AI models",
                                  bg=BG_COLOR, fg="#33363c",
                                  font=("Segoe UI", 9, "bold"), padx=10, pady=8)
            ai_f.pack(fill="x", padx=12, pady=(0, 6))

            model_var_list = []  # (model_key, BooleanVar) — which models to embed this run
            for model_key, info in SEMANTIC_MODELS.items():
                row = tk.Frame(ai_f, bg=BG_COLOR)
                row.pack(fill="x", pady=3)
                mdir = _find_model_dir_for(model_key)
                installed = bool(mdir) and os.path.isfile(os.path.join(mdir, "config.json"))
                label = info.get("label", model_key)

                # v7.10: a model that isn't installed locally can't build
                # embeddings this run regardless of the checkbox state, so
                # start it unchecked + disabled (greyed out, checkbox AND
                # its "Build embeddings for <model>" label together) instead
                # of letting the user tick a box that does nothing. It's
                # re-enabled automatically the moment install succeeds (see
                # _install_model_online's _ok()) without needing to reopen
                # this dialog.
                mv = tk.BooleanVar(value=installed)
                model_var_list.append((model_key, mv))
                cb = tk.Checkbutton(row, text=f"Build embeddings for {label}",
                                     variable=mv, bg=BG_COLOR, font=("Segoe UI", 9),
                                     state=("normal" if installed else "disabled"),
                                     disabledforeground="#a0a0a0")
                cb.pack(side="left")

                status_txt = ("✓ installed" if installed else "not installed locally")
                status_fg = "#4caf50" if installed else "#888888"
                status_lbl = tk.Label(row, text=status_txt, bg=BG_COLOR, fg=status_fg,
                                       font=("Segoe UI", 9))
                status_lbl.pack(side="left", padx=(8, 0))

                # v5.9b: install button is NEVER greyed out / disabled, even
                # once "installed" — if a previous download was Aborted, the
                # model folder can exist but be incomplete, and a disabled
                # button would make it impossible to try installing again.
                # Always clickable so the user can (re)install / overwrite
                # at any time.
                install_btn = tk.Button(row, text=("Reinstall..." if installed else "Install online..."),
                                         font=("Segoe UI", 8), cursor="hand2", state="normal")
                install_btn.config(command=lambda mk=model_key, sl=status_lbl, ib=install_btn, ckb=cb, mvv=mv:
                                    _install_model_online(mk, sl, ib, dlg, ckb, mvv))
                install_btn.pack(side="right")


            tk.Label(ai_f, text="Installing requires internet access (downloads from HuggingFace).",
                     bg=BG_COLOR, fg="#888888", font=("Segoe UI", 8)).pack(anchor="w", pady=(4, 0))

            # ── Vietnamese diacritics restoration (optional, separate from
            # the embedding models above — this is a query-preprocessing
            # helper, not something you pick for search, so no "build
            # embeddings" checkbox here, just install status) ─────────────
            vidia_f = tk.LabelFrame(dlg, text="Vietnamese diacritics restoration (optional)",
                                     bg=BG_COLOR, fg="#33363c",
                                     font=("Segoe UI", 9, "bold"), padx=10, pady=8)
            vidia_f.pack(fill="x", padx=12, pady=(0, 6))
            vidia_row = tk.Frame(vidia_f, bg=BG_COLOR)
            vidia_row.pack(fill="x", pady=3)
            tk.Label(vidia_row, text="Improves AI Search for queries typed without dấu",
                     bg=BG_COLOR, font=("Segoe UI", 9)).pack(side="left")
            vidia_installed = _vi_dia_installed()
            vidia_status_lbl = tk.Label(
                vidia_row, text=("✓ installed" if vidia_installed else "not installed"),
                bg=BG_COLOR, fg=("#4caf50" if vidia_installed else "#888888"),
                font=("Segoe UI", 9))
            vidia_status_lbl.pack(side="left", padx=(8, 0))
            vidia_install_btn = tk.Button(
                vidia_row, text=("Reinstall..." if vidia_installed else "Install online..."),
                font=("Segoe UI", 8), cursor="hand2", state="normal")
            vidia_install_btn.config(
                command=lambda: _install_vi_diacritics_online(vidia_status_lbl, vidia_install_btn, dlg))
            vidia_install_btn.pack(side="right")

            # ── Action buttons ───────────────────────────────────────────────
            btn_row = tk.Frame(dlg, bg=BG_COLOR)
            btn_row.pack(fill="x", padx=12, pady=(6, 12))

            def _start():
                chosen = {i for i, v in enumerate(tier_var_list) if v.get()}
                if not chosen:
                    messagebox.showwarning("Update DB", "Select at least one tier.", parent=dlg)
                    return
                selected_tiers = None if chosen == {0, 1, 2, 3} else chosen
                chosen_models = [mk for mk, v in model_var_list if v.get()]
                # v8.4: no models ticked is now a valid, deliberate choice —
                # "Tier-only" update (BM25/content only, AI stage skipped
                # entirely this run). Confirm instead of blocking, since it's
                # easy to forget to tick a model and not realize AI won't be
                # touched.
                if not chosen_models:
                    if not messagebox.askyesno(
                        "Update DB",
                        "No AI model selected — this run will update BM25/content "
                        "search only. AI Search data will NOT be built or refreshed.\n\n"
                        "Continue with Tier-only update?", parent=dlg):
                        return
                    selected_models = []  # explicit empty = skip AI entirely (see indexing_worker)
                else:
                    selected_models = None if len(chosen_models) == len(model_var_list) else chosen_models
                dlg.destroy()
                _run_update_db(selected_tiers, selected_models, ocr_var.get())

            tk.Button(btn_row, text="Start Update", font=("Segoe UI", 9, "bold"),
                      bg="#2196f3", fg="white", activebackground="#1976d2",
                      cursor="hand2", padx=10, command=_start).pack(side="right", padx=(6, 0))
            tk.Button(btn_row, text="Cancel", font=("Segoe UI", 9), cursor="hand2",
                      padx=10, command=dlg.destroy).pack(side="right")

            # v7.10: center the dialog on the Search GUI (self.root) instead
            # of just offsetting from its top-left corner. Done last, now
            # that every widget above has been packed, so winfo_reqwidth/
            # reqheight below reflect the dialog's real final size.
            dlg.update_idletasks()
            rw, rh = self.root.winfo_width(), self.root.winfo_height()
            rx, ry = self.root.winfo_x(), self.root.winfo_y()
            dw, dh = dlg.winfo_reqwidth(), dlg.winfo_reqheight()
            cx = rx + (rw - dw) // 2
            cy = ry + (rh - dh) // 2
            # keep it fully on-screen
            sw, sh = dlg.winfo_screenwidth(), dlg.winfo_screenheight()
            cx = max(0, min(cx, sw - dw))
            cy = max(0, min(cy, sh - dh))
            dlg.geometry(f"+{cx}+{cy}")

        _btn_text = "Update DB"
        if self._update_db_running:
            # v9.11 fix: prefer the last real progress string (e.g. "AI 2/2:
            # 28%") over a bare "Updating..." — this button gets recreated
            # from scratch whenever the search box is minimized/reopened
            # mid-run, and previously always lost the percentage at that
            # point even though indexing was still happily running in the
            # background the whole time.
            _btn_text = self._last_index_status_text or "Updating..."
        self._update_db_btn = tk.Button(
            search_bar, text=_btn_text,
            font=("Segoe UI", 8, "bold"),
            bg="#e6e8ec", fg="#33363c",
            relief="raised", bd=2,
            padx=8, pady=2,
            activebackground="#d3d6dc", activeforeground="#111111",
            cursor="hand2",
            state=("disabled" if self._update_db_running else "normal"),
            command=_on_update_db)
        self._update_db_btn.pack(side="right", padx=(4, 2), pady=4)
        add_tooltip(self._update_db_btn,
                    "Open Update DB options (choose file-type tiers,\n"
                    "check/install AI Search models) then start the scan.")
        # ──────────────────────────────────────────────────────────────────


        self.c_tree_frame = tk.Frame(c_frame)
        self.c_tree_frame.pack(fill="both", expand=True)

        self.tree_c = ttk.Treeview(self.c_tree_frame, columns=("icon_name", "size", "mtime", "ftype", "full_path"), show="tree headings")
        self.tree_c.heading("#0", text="")
        self.tree_c.column("#0", width=40, minwidth=40, stretch=False, anchor="center")
        self.tree_c.heading("icon_name", text="File Location")
        self.tree_c.heading("size",      text="File Size")
        self.tree_c.heading("mtime",     text="Date Modified")
        self.tree_c.heading("ftype",     text="Type")
        self.tree_c.column("icon_name",  width=850, stretch=False)
        self.tree_c.column("size",       width=90,  anchor="center", stretch=False)
        self.tree_c.column("mtime",      width=130, anchor="center", stretch=False)
        self.tree_c.column("ftype",      width=55,  anchor="center", stretch=False)
        self.tree_c["displaycolumns"] = ("icon_name", "size", "mtime", "ftype")
        sb_c = ttk.Scrollbar(self.c_tree_frame, orient="vertical", command=self.tree_c.yview)
        sb_c_x = ttk.Scrollbar(self.c_tree_frame, orient="horizontal", command=self.tree_c.xview)
        self.tree_c.configure(yscrollcommand=sb_c.set, xscrollcommand=sb_c_x.set)
        sb_c.pack(side="right", fill="y")
        sb_c_x.pack(side="bottom", fill="x")
        self.tree_c.pack(side="left", fill="both", expand=True)

        # 2. TAB FILE NAME — added FIRST to notebook (MFT default tab)
        f_frame = tk.Frame(self.nb)
        self.nb.add(f_frame, text=" File Name ")

        # ── Smart Size Filter bar ──────────────────────────────────────────
        filter_bar = tk.Frame(f_frame, bg=BG_COLOR, pady=2)
        filter_bar.pack(fill="x", padx=6, pady=(4, 2))

        tk.Label(filter_bar, text="Size:", bg=BG_COLOR, fg=TEXT_COLOR,
                 font=("Segoe UI", 9)).pack(side="left", padx=(4, 2))

        op_menu = ttk.OptionMenu(filter_bar, self.size_op_var, self.size_op_var.get(),
                                  "Any", ">", ">=", "<", "<=", "=",
                                  command=lambda *_: self._apply_size_filter_to_tree())
        op_menu.config(width=4); op_menu.pack(side="left", padx=2)

        size_num_entry = tk.Entry(filter_bar, textvariable=self.size_num_var, width=7,
                                   font=("Segoe UI", 9), bg=ENTRY_BG, fg=TEXT_COLOR,
                                   insertbackground=TEXT_COLOR, bd=1, relief="flat")
        size_num_entry.pack(side="left", padx=2)
        self.size_num_var.trace_add("write", self._apply_size_filter_to_tree)

        unit_menu = ttk.OptionMenu(filter_bar, self.size_unit_var, self.size_unit_var.get(),
                                    "B", "KB", "MB", "GB",
                                    command=lambda *_: self._apply_size_filter_to_tree())
        unit_menu.config(width=4); unit_menu.pack(side="left", padx=2)

        def _clear_size_filter():
            self.size_op_var.set(">"); self.size_num_var.set(""); self.size_unit_var.set("MB")
            self._apply_size_filter_to_tree()

        # ── Extension filter ───────────────────────────────────────────────
        tk.Label(filter_bar, text="│", bg=BG_COLOR, fg="#555",
                 font=("Segoe UI", 11)).pack(side="left", padx=4)
        tk.Label(filter_bar, text="Ext:", bg=BG_COLOR, fg=TEXT_COLOR,
                 font=("Segoe UI", 9)).pack(side="left", padx=(0, 3))

        # Quick preset buttons
        _EXT_PRESETS = [
            ("All",   ""),
            ("PDF",   "pdf"),
            ("Word",  "doc,docx"),
            ("Excel", "xls,xlsx,csv"),
            ("PPT",   "ppt,pptx"),
            ("Img",   "png,jpg,jpeg,bmp,gif"),
            ("Txt",   "txt,log"),
        ]
        def _set_ext(val):
            self.ext_filter_var.set(val)   # trace fires _apply_size_filter_to_tree

        for _lbl, _val in _EXT_PRESETS:
            tk.Button(filter_bar, text=_lbl, font=("Segoe UI", 8),
                      bg="#e6e8ec", fg="#33363c", bd=0, padx=5, pady=1,
                      activebackground="#d3d6dc", cursor="hand2",
                      command=lambda v=_val: _set_ext(v)).pack(side="left", padx=1)

        tk.Label(filter_bar, text="│", bg=BG_COLOR, fg="#555",
                 font=("Segoe UI", 9)).pack(side="left", padx=3)

        ext_entry = tk.Entry(filter_bar, textvariable=self.ext_filter_var, width=14,
                              font=("Segoe UI", 9), bg=ENTRY_BG, fg="#888",
                              insertbackground=TEXT_COLOR, bd=1, relief="flat")
        ext_entry.pack(side="left", padx=(0, 2))

        def _ext_focus_in(e):
            if not self.ext_filter_var.get(): ext_entry.config(fg=TEXT_COLOR)
        def _ext_focus_out(e):
            if not self.ext_filter_var.get(): ext_entry.config(fg="#888")
        ext_entry.bind("<FocusIn>",  _ext_focus_in)
        ext_entry.bind("<FocusOut>", _ext_focus_out)
        # trace already wired through ext_filter_var → _apply_size_filter_to_tree
        self.ext_filter_var.trace_add("write", self._apply_size_filter_to_tree)

        tk.Button(filter_bar, text="✕", command=lambda: self.ext_filter_var.set(""),
                  bg="#c9ccd2", fg="#222222", font=("Segoe UI", 8), bd=0, padx=4,
                  activebackground="#b0b3ba", cursor="hand2").pack(side="left", padx=(0, 4))
        # ──────────────────────────────────────────────────────────────────

        # ── Whole word toggle ───────────────────────────────────────────────
        # v7.10: OFF by default (raw substring, original behavior). When ON,
        # a keyword only counts as a match when it's bounded by non-letter
        # characters on both sides -- e.g. searching "adas" still matches
        # "ADAS_systems.pdf" / "VDIM_0ADAS_ESP..." but no longer matches
        # buried-substring false positives like "readasync.xml" or
        # "ReadAStringExample.mlx". Applies to both the live MFT scan (File
        # Name / Folder Name, no DB needed) and the DB-backed search.
        tk.Label(filter_bar, text="│", bg=BG_COLOR, fg="#555",
                 font=("Segoe UI", 9)).pack(side="left", padx=3)
        _ww_cb = tk.Checkbutton(filter_bar, text="Whole word", variable=self.whole_word_var,
                                  bg=BG_COLOR, fg=TEXT_COLOR, selectcolor=ENTRY_BG,
                                  font=("Segoe UI", 9),
                                  command=self._rerun_current_search)
        _ww_cb.pack(side="left", padx=(0, 4))
        add_tooltip(_ww_cb,
                     "ON: \"adas\" matches ADAS_systems.pdf but not readasync.xml\n"
                     "OFF (default): \"adas\" matches anywhere, incl. inside other words")
        # ──────────────────────────────────────────────────────────────────

        # ── Name filter ────────────────────────────────────────────────────
        tk.Label(filter_bar, text="│", bg=BG_COLOR, fg="#555",
                 font=("Segoe UI", 9)).pack(side="left", padx=3)
        tk.Label(filter_bar, text="Name:", bg=BG_COLOR, fg=TEXT_COLOR,
                 font=("Segoe UI", 9)).pack(side="left", padx=(0, 2))
        name_entry = tk.Entry(filter_bar, textvariable=self.name_filter_var, width=18,
                               font=("Segoe UI", 9), bg=ENTRY_BG, fg="#888",
                               insertbackground=TEXT_COLOR, bd=1, relief="flat")
        name_entry.pack(side="left", padx=(0, 2))
        def _name_focus_in(e):
            if not self.name_filter_var.get(): name_entry.config(fg=TEXT_COLOR)
        def _name_focus_out(e):
            if not self.name_filter_var.get(): name_entry.config(fg="#888")
        name_entry.bind("<FocusIn>",  _name_focus_in)
        name_entry.bind("<FocusOut>", _name_focus_out)
        self.name_filter_var.trace_add("write", self._apply_size_filter_to_tree)
        tk.Button(filter_bar, text="✕", command=lambda: self.name_filter_var.set(""),
                  bg="#c9ccd2", fg="#222222", font=("Segoe UI", 8), bd=0, padx=4,
                  activebackground="#b0b3ba", cursor="hand2").pack(side="left", padx=(0, 4))
        # ──────────────────────────────────────────────────────────────────

        self.filter_count_label = None  # removed - counts shown in pane headers
        # ──────────────────────────────────────────────────────────────────

        self.f_tree_frame = tk.Frame(f_frame)
        self.f_tree_frame.pack(fill="both", expand=True)

        self.tree_f = ttk.Treeview(self.f_tree_frame, columns=("icon_name", "size", "mtime", "ftype", "location", "full_path"), show="tree headings")
        self.tree_f.heading("#0", text="")
        self.tree_f.column("#0", width=40, minwidth=40, stretch=False, anchor="center")
        self.tree_f.heading("icon_name", text="File Name")
        self.tree_f.heading("size",      text="File Size")
        self.tree_f.heading("mtime",     text="Date Modified")
        self.tree_f.heading("ftype",     text="Type")
        self.tree_f.heading("location",  text="File Location")
        self.tree_f.column("icon_name",  width=380, minwidth=150, stretch=False)
        self.tree_f.column("size",       width=90,  anchor="center", stretch=False, minwidth=70)
        self.tree_f.column("mtime",      width=130, anchor="center", stretch=False, minwidth=100)
        self.tree_f.column("ftype",      width=55,  anchor="center", stretch=False, minwidth=40)
        self.tree_f.column("location",   width=820, minwidth=300, stretch=False)
        self.tree_f["displaycolumns"] = ("icon_name", "size", "mtime", "ftype", "location")
        sb_f  = ttk.Scrollbar(self.f_tree_frame, orient="vertical",   command=self.tree_f.yview)
        sb_fx = ttk.Scrollbar(self.f_tree_frame, orient="horizontal", command=self.tree_f.xview)
        self.tree_f.configure(yscrollcommand=sb_f.set, xscrollcommand=sb_fx.set)
        # Pack order matters: scrollbars first, then tree fills remaining space
        sb_f.pack(side="right",  fill="y")
        sb_fx.pack(side="bottom", fill="x")
        self.tree_f.pack(side="left", fill="both", expand=True)

        # 3. TAB FOLDER NAME — added SECOND to notebook
        fol_frame = tk.Frame(self.nb)
        self.nb.add(fol_frame, text=" Folder Name ")
        self.fol_tree_frame = tk.Frame(fol_frame)
        self.fol_tree_frame.pack(fill="both", expand=True)
        self.tree_fol = ttk.Treeview(self.fol_tree_frame, columns=("icon_name", "mtime", "location", "full_path"), show="tree headings")
        self.tree_fol.heading("#0", text="")
        self.tree_fol.column("#0", width=40, minwidth=40, stretch=False, anchor="center")
        self.tree_fol.heading("icon_name", text="Folder Name")
        self.tree_fol.heading("mtime",     text="Date Modified")
        self.tree_fol.heading("location",  text="Folder Location")
        self.tree_fol.column("icon_name",  width=440, minwidth=150, stretch=False)
        self.tree_fol.column("mtime",      width=130, anchor="center", stretch=False, minwidth=100)
        self.tree_fol.column("location",   width=750, minwidth=300, stretch=False)
        self.tree_fol["displaycolumns"] = ("icon_name", "mtime", "location")
        sb_fol  = ttk.Scrollbar(self.fol_tree_frame, orient="vertical",   command=self.tree_fol.yview)
        sb_folx = ttk.Scrollbar(self.fol_tree_frame, orient="horizontal", command=self.tree_fol.xview)
        self.tree_fol.configure(yscrollcommand=sb_fol.set, xscrollcommand=sb_folx.set)
        sb_fol.pack(side="right",  fill="y")
        sb_folx.pack(side="bottom", fill="x")
        self.tree_fol.pack(side="left", fill="both", expand=True)

        # v2.3: File Content tab added LAST (rightmost) — only populated after --update data
        self.nb.add(c_frame, text=" File Content ")

        # v3.4: back to being a real Notebook tab (a floating toggle button
        # felt wrong/inconsistent with the other tabs). ttk.Notebook has no
        # way to pin a tab flush to the far-right edge -- tabs are always
        # left-aligned/sequential -- so as a fallback this uses a wide
        # disabled (unclickable) spacer tab to visually push Help further
        # away from File Content instead.
        _spacer_frame = tk.Frame(self.nb)
        self.nb.add(_spacer_frame, text=" " * 28)
        self.nb.tab(_spacer_frame, state="disabled")

        help_frame = tk.Frame(self.nb, bg=BG_COLOR)
        help_split = tk.PanedWindow(help_frame, orient="horizontal", bg=BG_COLOR,
                                     sashwidth=4, sashrelief="raised", bd=0)
        help_split.pack(fill="both", expand=True, padx=2, pady=2)
        # Bordered panels (grooved relief) so the split is easy to read even
        # though the surrounding window has no title bar / chrome of its own.
        help_left = tk.Frame(help_split, bg=BG_COLOR, relief="groove", bd=2)
        help_right = tk.Frame(help_split, bg=BG_COLOR, relief="groove", bd=2)
        help_split.add(help_left, minsize=250)
        help_split.add(help_right, minsize=250)
        tk.Label(help_left, text=" Help ", font=("Segoe UI", 9, "bold"),
                 bg=BG_COLOR, fg="#7ec8e3", anchor="w").pack(fill="x", padx=4, pady=(4, 0))
        tk.Label(help_right, text=" Search History ", font=("Segoe UI", 9, "bold"),
                 bg=BG_COLOR, fg="#7ec8e3", anchor="w").pack(fill="x", padx=4, pady=(4, 0))
        self._build_help_content(help_left)
        HistoryPanel(help_right, self)
        self.nb.add(help_frame, text=" Help ")

        # Default the Help/History split to 50%-50% once the pane has a real
        # width (can't compute % of width before it's actually laid out).
        # Only fires once -- if the user drags the sash afterwards, further
        # window resizes won't snap it back to 50/50.
        _sash_done = {"v": False}
        def _center_help_sash(event=None):
            if _sash_done["v"]:
                return
            try:
                total_w = help_split.winfo_width()
                if total_w > 20:
                    help_split.sash_place(0, total_w // 2, 0)
                    _sash_done["v"] = True
            except Exception:
                pass
        help_split.bind("<Configure>", _center_help_sash, add="+")

        self.tree_c.search_id = version
        self.tree_f.search_id = version
        self.tree_fol.search_id = version

        # Initialize pane-tree dicts so filters work from the very first pane
        self._c_pane_trees   = {"main": self.tree_c}
        self._f_pane_trees   = {"main": self.tree_f}
        self._fol_pane_trees = {"main": self.tree_fol}

        def bg_load_ui(tree, data, mode, current_vid):
            if not data or getattr(tree, 'search_id', 0) != current_vid: return
            CHUNK = 500
            first_chunk, remaining_chunk = data[:CHUNK], data[CHUNK:]
            offset = [0]
            for item in first_chunk:
                if getattr(tree, 'search_id', 0) != current_vid: return
                try:
                    offset[0] += 1
                    self._insert_row(tree, item, offset[0], mode)
                except: pass

            def _update_count():
                if mode == 1 and self.filter_count_label:
                    self.filter_count_label.config(text=f"{len(self._all_files_data)} files")
                elif mode == 0 and self.content_filter_count_label:
                    self.content_filter_count_label.config(text=f"{len(self._all_content_data)} files")

            def load_rest(step=0):
                if not self.active_result_win or not tk.Toplevel.winfo_exists(self.active_result_win): return
                if getattr(tree, 'search_id', 0) != current_vid: return
                sub_chunk = remaining_chunk[step: step + CHUNK]
                if sub_chunk:
                    for item in sub_chunk:
                        if getattr(tree, 'search_id', 0) != current_vid: return
                        try:
                            offset[0] += 1
                            self._insert_row(tree, item, offset[0], mode)
                        except: pass
                    self.root.after(5, lambda: load_rest(step + CHUNK))
                else:
                    _update_count()
            if remaining_chunk:
                self.root.after(10, lambda: load_rest(0))
            else:
                self.root.after(20, _update_count)

        bg_load_ui(self.tree_c, cont_res, 0, version)
        bg_load_ui(self.tree_f, only_files, 1, version)
        bg_load_ui(self.tree_fol, only_folders, 2, version)

        def on_tree_click(event, tree):
            row = tree.identify_row(event.y); col = tree.identify_column(event.x)
            if row and col:
                now = time.time()
                if getattr(tree, "_last_row", None) == row and getattr(tree, "_last_col", None) == col and (now - getattr(tree, "_last_time", 0)) > 0.4:
                    show_cell_edit(tree, row, col)
                tree._last_row, tree._last_col, tree._last_time = row, col, now

        def show_cell_edit(tree, row, col):
            if self.res_edit: self.res_edit.destroy()
            if col == "#0": return  # v5.6: icon+# combined column, nothing to edit
            idx = int(col[1:]) - 1
            val = str(tree.item(row, "values")[idx])
            if val.startswith(("📁 ", "📕 ", "📝 ", "📊 ", "📉 ", "🔮 ", "📄 ", "📎 ")): val = val[2:]
            bbox = tree.bbox(row, col)
            if not bbox: return  # cell not visible (row scrolled out or window resized)
            x, y, w, h = bbox
            self.res_edit = tk.Entry(tree, font=("Segoe UI", 9), bd=0); self.res_edit.insert(0, val)
            self.res_edit.place(x=x, y=y, width=w, height=h); self.res_edit.focus_set()
            add_only_copy_menu(self.res_edit); self.res_edit.bind("<FocusOut>", lambda e: self.res_edit.destroy())

        # Extensions that open directly (Office, PDF, text, logs, etc.)
        OPEN_DIRECT_EXTS = {
            '.docx', '.doc', '.docm', '.xlsx', '.xls', '.xlsm',
            '.pptx', '.ppt', '.pptm', '.pdf', '.msg', '.one',
            '.txt', '.csv', '.tsv', '.json', '.xml', '.yaml', '.yml',
            '.ini', '.cfg', '.conf', '.toml', '.md',
            '.log', '.html', '.htm', '.bat', '.ps1', '.sh',
            '.py', '.js', '.ts', '.cpp', '.c', '.h', '.java',
            '.fmu', '.scen', '.mdl', '.m', '.slx',
        }
        # Extensions to open with Notepad++ if available, else Notepad
        NOTEPAD_EXTS = {
            '.log', '.html', '.htm', '.bat', '.ps1', '.sh',
            '.py', '.js', '.ts', '.cpp', '.c', '.h', '.java',
            '.ini', '.cfg', '.conf', '.toml', '.yaml', '.yml',
            '.json', '.xml', '.md', '.txt', '.csv', '.tsv',
        }
        # Find Notepad++ once
        _npp_paths = [
            r"C:\Program Files\Notepad++\notepad++.exe",
            r"C:\Program Files (x86)\Notepad++\notepad++.exe",
        ]
        NOTEPAD_PP = next((p for p in _npp_paths if os.path.isfile(p)), None)

        def _normalize_path(raw):
            return os.path.abspath(os.path.normpath(raw.replace("¥", "\\")))

        def open_file(tree):
            """Double-click: open file directly if possible, else open folder."""
            sel = tree.selection()
            if not sel: return
            try:
                values = tree.item(sel[0])["values"]
                raw_path = str(values[-1])
                if raw_path.startswith(("📁 ", "📕 ", "📝 ", "📊 ", "📉 ", "🔮 ", "📄 ", "📎 ")):
                    raw_path = raw_path[2:]
                abs_p = _normalize_path(raw_path)
                ext = os.path.splitext(abs_p)[1].lower()
                has_ext = bool(ext)
                is_file = has_ext and not os.path.isdir(abs_p)
                if is_file and ext in OPEN_DIRECT_EXTS:
                    if ext in NOTEPAD_EXTS:
                        editor = NOTEPAD_PP or "notepad.exe"
                        subprocess.Popen([editor, abs_p])
                    else:
                        try:
                            os.startfile(abs_p)
                        except:
                            subprocess.Popen(f'cmd /c start "" "{abs_p}"', shell=True)
                    return
                # Not openable directly → fall through to open folder
                open_explorer(tree)
            except Exception as ex:
                print(f"Open File Error: {ex}")

        def open_explorer(tree, is_content=False):
            sel = tree.selection()
            if not sel: return
            try:
                values = tree.item(sel[0])["values"]
                raw_path = str(values[-1])
                if raw_path.startswith(("📁 ", "📕 ", "📝 ", "📊 ", "📉 ", "🔮 ", "📄 ", "📎 ")):
                    raw_path = raw_path[2:]
                abs_p = _normalize_path(raw_path)
                has_ext = bool(os.path.splitext(abs_p)[1])
                is_file = has_ext and not os.path.isdir(abs_p)
                folder  = os.path.dirname(abs_p) if is_file else abs_p
                if is_file and len(abs_p) <= 250:
                    try:
                        subprocess.Popen(f'cmd /c explorer /select,"{abs_p}"', shell=True)
                        return
                    except:
                        pass
                try:
                    os.startfile(folder)
                except:
                    subprocess.Popen(f'cmd /c explorer "{folder}"', shell=True)
            except Exception as ex:
                print(f"Explorer Failure: {ex}")

        def show_context_menu(e, tree, is_folder=False, is_content=False):
            row = tree.identify_row(e.y)
            if not row:
                return
            tree.selection_set(row)
            val = list(tree.item(row)["values"])
            name_val = str(val[1])
            if name_val.startswith(("📁 ", "📕 ", "📝 ", "📊 ", "📉 ", "🔮 ", "📄 ", "📎 ")):
                name_val = name_val[2:]
            full_path = str(val[-1])
            ext = os.path.splitext(full_path)[1].lower()
            # Use self._result_win so this works from any pane (split or not)
            ctx_win = self._result_win if hasattr(self, '_result_win') and self._result_win.winfo_exists() else win
            m = tk.Menu(ctx_win, tearoff=0)
            if is_content and ext in OPEN_DIRECT_EXTS:
                m.add_command(label="Open File",   command=lambda t=tree: open_file(t))
                m.add_command(label="Open Folder", command=lambda t=tree: open_explorer(t))
            else:
                m.add_command(label="Open Folder", command=lambda t=tree: open_explorer(t))
            m.add_separator()
            lbl_name = "Copy Folder Name" if is_folder else "Copy File Name"
            lbl_path = "Copy Folder Path" if is_folder else "Copy File Path"
            m.add_command(label=lbl_name, command=lambda v=name_val:  (ctx_win.clipboard_clear(), ctx_win.clipboard_append(v)))
            m.add_command(label=lbl_path, command=lambda v=full_path: (ctx_win.clipboard_clear(), ctx_win.clipboard_append(v)))
            m.post(e.x_root, e.y_root)

        def dbl_select(tree):
            row = tree.identify_row(tree.winfo_pointery() - tree.winfo_rooty())
            if row: tree.selection_set(row); tree.focus(row)

        # ── Store handlers as instance attrs so _wire_tree_clicks can use them
        self._fn_open_file     = open_file
        self._fn_open_explorer = open_explorer
        self._fn_show_ctx_menu = show_context_menu
        self._fn_on_tree_click = on_tree_click
        self._fn_dbl_select    = dbl_select

        # ── Sortable columns ──────────────────────────────────────────────
        # _sort_tree / _make_sortable have been promoted to class methods (see
        # above, near _build_tree_*) so all Advanced/AI panes created after split are also
        # automatically sortable. Only need to wire the 3 root trees (Default) here.
        self._make_sortable(self.tree_c,   ["size", "mtime", "ftype"])
        self._make_sortable(self.tree_f,   ["size", "mtime", "ftype"])
        self._make_sortable(self.tree_fol, ["mtime"])
        # ─────────────────────────────────────────────────────────────────

        # Wire click events on the initial (main) trees
        self._wire_tree_clicks(self.tree_f,   "f")
        self._wire_tree_clicks(self.tree_c,   "c")
        self._wire_tree_clicks(self.tree_fol, "fol")

        # v2.6: the Notebook (self.nb) was created/packed AFTER the resize_grip,
        # so it stacks visually on top and would swallow the grip's clicks.
        # Re-raise the grip now that every widget in this window has been built.
        try:
            resize_grip.lift()
        except Exception:
            pass

    # ── Smart Size + Extension Filter ────────────────────────────────────
    def _get_size_filter_bytes(self):
        """Return (op, bytes) from dropdown vars, or (None, None) when inactive."""
        op = self.size_op_var.get()
        if op == "Any":
            return None, None
        try:
            val = float(self.size_num_var.get())
        except ValueError:
            return None, None
        unit = self.size_unit_var.get()
        mult = {'B': 1, 'KB': 1024, 'MB': 1024**2, 'GB': 1024**3}.get(unit, 1)
        return op, int(val * mult)

    def _get_ext_filter(self):
        """Parse ext_filter_var → frozenset of .ext strings, e.g. {'.pdf','.docx'}, or None."""
        raw = self.ext_filter_var.get().strip()
        if not raw:
            return None
        exts = set()
        for e in raw.replace(',', ' ').split():
            e = e.strip().lower().lstrip('.')
            if e:
                exts.add('.' + e)
        return frozenset(exts) if exts else None

    def _update_pane_label(self, tree, shown, total):
        """Update the header label of a pane tree with filtered count."""
        try:
            lbl = getattr(tree, '_header_label', None)
            base = getattr(tree, '_base_label', '')
            if not lbl or not lbl.winfo_exists():
                return
            # Extract base text without count: e.g. '📄 Default (33)' → '📄 Default'
            import re as _re
            base_no_count = _re.sub(r'\s*\(.*\)\s*$', '', base).strip()
            if shown < total:
                lbl.config(text=f"{base_no_count} ({shown}/{total})")
            else:
                lbl.config(text=f"{base_no_count} ({total})")
        except Exception:
            pass

    def _apply_size_filter_to_tree(self, *_):
        """Client-side filter: size, extension, AND name substring.
        Applies to ALL visible panes in File Name tab, each with its own data source."""
        if not hasattr(self, 'tree_f') or not self.active_result_win:
            return
        try:
            if not tk.Toplevel.winfo_exists(self.active_result_win):
                return
            if not self.tree_f.winfo_exists():
                return
        except Exception:
            return

        # Data source per pane key — each pane must only show its own results
        # "main" pane = MFT (realtime scan) blended with BM25 (DB), deduplicated by path.
        # MFT goes first (user already saw it stream in), BM25 fills in anything extra.
        mft_files  = [r for r in (self._mft_file_res or [])        if str(r[0]).lower() != "folder"]
        bm25_files = [r for r in (self._last_bm25_file_res or []) if str(r[0]).lower() != "folder"]
        adv_files  = [r for r in (self._adv_all_files or [])      if str(r[0]).lower() != "folder"]
        ai_files   = [r for r in (self._ai_file_res or [])        if str(r[0]).lower() != "folder"]

        _seen_paths = set()
        main_files = []
        for r in (mft_files + bm25_files):
            p = r[2]
            if p in _seen_paths:
                continue
            _seen_paths.add(p)
            main_files.append(r)

        # Priority sort (office/pdf/msg first, log/code last) — same sort used
        # by the MFT live renderer, so manual filter edits don't undo it.
        main_files = self._sort_priority(main_files, 1)
        adv_files  = self._sort_priority(adv_files, 1)
        ai_files   = self._sort_priority(ai_files, 1)

        pane_sources = {"main": main_files, "adv": adv_files, "ai": ai_files}

        total_shown = 0
        for pane_key, tree in self._f_pane_trees.items():
            try:
                if not tree.winfo_exists(): continue
            except Exception:
                continue
            source = pane_sources.get(pane_key, main_files)
            filtered = self._filter_file_rows(source)
            x0 = self._save_xview(tree)
            widths = self._save_col_widths(tree)
            for row in tree.get_children(): tree.delete(row)
            for rn, item in enumerate(filtered, start=1):
                try:
                    img = get_tree_icon_image(item[2], is_folder=False)
                    name_txt = item[1] if img else get_file_icon(item[2], is_folder=False) + item[1]
                    readable_sz = format_size(get_live_size(item[2], item[3]))
                    mtime = get_live_mtime(item[2])
                    ftype = get_file_type(item[2])
                    tree.insert("", tk.END, text="", **({"image": img} if img else {}), values=(name_txt, readable_sz, mtime, ftype, os.path.dirname(item[2]), item[2]))
                except Exception:
                    pass
            self._restore_col_widths(tree, widths)
            self._restore_xview(tree, x0)
            self._update_pane_label(tree, len(filtered), len(source))
            if pane_key == "main":
                total_shown = len(filtered)
    # ─────────────────────────────────────────────────────────────────────

    def _get_content_size_filter(self):
        op = self.c_size_op_var.get()
        if op == "Any":
            return None, None
        try:
            val = float(self.c_size_num_var.get())
        except ValueError:
            return None, None
        unit = self.c_size_unit_var.get()
        mult = {'B': 1, 'KB': 1024, 'MB': 1024**2, 'GB': 1024**3}.get(unit, 1)
        return op, int(val * mult)

    def _get_content_ext_filter(self):
        raw = self.c_ext_filter_var.get().strip()
        if not raw:
            return None
        exts = set()
        for e in raw.replace(',', ' ').split():
            e = e.strip().lower().lstrip('.')
            if e:
                exts.add('.' + e)
        return frozenset(exts) if exts else None

    def _apply_content_filter_to_tree(self, *_):
        """Client-side filter for File Content tab: size, extension, AND name substring.
        Applies to ALL visible panes, each with its own data source."""
        if not hasattr(self, 'tree_c') or not self.active_result_win:
            return
        try:
            if not tk.Toplevel.winfo_exists(self.active_result_win):
                return
            if not self.tree_c.winfo_exists():
                return
        except Exception:
            return

        op, size_bytes = self._get_content_size_filter()
        ext_filter = self._get_content_ext_filter()
        name_needle = _norm_txt(self.c_name_filter_var.get().strip().lower())

        # Data source per pane key — each pane must only show its own results
        bm25_cont = self._last_bm25_cont_res or []
        adv_cont  = self._adv_all_cont  or []
        ai_cont   = self._ai_cont_res   or []
        pane_sources = {"main": bm25_cont, "adv": adv_cont, "ai": ai_cont}

        def _filter_content(data):
            if op is None and ext_filter is None and not name_needle:
                return data
            cmp_fn = {'>':  lambda a, b: a >  b, '>=': lambda a, b: a >= b,
                      '<':  lambda a, b: a <  b, '<=': lambda a, b: a <= b,
                      '=':  lambda a, b: a == b}.get(op) if op else None
            result = []
            for item in data:
                fpath = item[0]
                if cmp_fn is not None:
                    if not cmp_fn(get_live_size(fpath, 0), size_bytes):
                        continue
                if ext_filter is not None:
                    _, ext = os.path.splitext(fpath)
                    if ext.lower() not in ext_filter:
                        continue
                if name_needle:
                    fname = _norm_txt(os.path.basename(fpath).lower())
                    if name_needle not in fname:
                        continue
                result.append(item)
            return result

        total_shown = 0
        total_src   = 0
        for pane_key, tree in self._c_pane_trees.items():
            try:
                if not tree.winfo_exists(): continue
            except Exception:
                continue
            source = pane_sources.get(pane_key, bm25_cont)
            filtered = _filter_content(source)
            for row in tree.get_children(): tree.delete(row)
            for rn, item in enumerate(filtered, start=1):
                try:
                    img = get_tree_icon_image(item[0], is_folder=False)
                    name_txt = item[0] if img else get_file_icon(item[0], is_folder=False) + item[0]
                    readable_sz = format_size(get_live_size(item[0], 0))
                    mtime = get_live_mtime(item[0])
                    ftype = get_file_type(item[0])
                    tree.insert("", tk.END, text="", **({"image": img} if img else {}), values=(name_txt, readable_sz, mtime, ftype, item[0]))
                except Exception:
                    pass
            self._update_pane_label(tree, len(filtered), len(source))
            if pane_key == "main":
                total_shown = len(filtered)
                total_src   = len(source)
    # ─────────────────────────────────────────────────────────────────────

    def _save_hist(self, q):
        try:
            # v7.10 FIX: this used to connect to DB_FILE (search_data.db)
            # directly, which meant simply performing a search -- even
            # with --update data NEVER run -- created/grew search_data.db
            # on disk. On the next launch the ramp light saw that file
            # "exists" and jumped from Red to Yellow, even though nothing
            # was ever actually indexed. History now lives in its own
            # HISTORY_DB_FILE so searching can never touch search_data.db.
            conn = sqlite3.connect(HISTORY_DB_FILE); c = conn.cursor()
            # v7.4 FIX: history table used to only get created inside
            # indexing_worker (--update data), so on a fresh/never-updated
            # DB (search_data.db = 0KB) every save silently failed with
            # "no such table: history". Ensure it exists here too.
            c.execute("CREATE TABLE IF NOT EXISTS history (id INTEGER PRIMARY KEY, query TEXT, date TEXT)")
            c.execute("INSERT INTO history (query, date) VALUES (?,?)", (q, datetime.now().strftime("%Y-%m-%d %H:%M:%S"))); conn.commit(); conn.close()
            self._last_saved_hist_query = q
        except Exception:
            pass

    def _maybe_save_hist_idle(self, snapshot):
        """Fired ~10s after the user stops typing without pressing Enter.
        The realtime preview means Enter is rarely needed to SEE results,
        but that also meant most searches never got saved to History at
        all. Save the query anyway once the user has genuinely paused on
        it -- but only if the box still holds that exact text (they didn't
        keep typing/edit since) and it isn't a duplicate of what's already
        the most recently saved entry (e.g. they already pressed Enter, or
        this same idle-save already fired once for this text)."""
        self._hist_idle_timer = None
        if self.entry_var.get().strip() != snapshot:
            return  # stale — text changed since this timer was scheduled
        if not snapshot or snapshot == self._last_saved_hist_query:
            return
        if snapshot.lower().startswith("-"):
            return  # commands like --update data / --exit aren't search queries
        self._save_hist(snapshot)

    def resolve_id(self, text):
        if re.match(r"^[a-z]{3}\d*$", text, re.I): return f"{BASE_DS}{text.lower()}"
        patterns = {
            "SR": r"^SR\d{8}(?:-\d{2})?$", "NTSR": r"^NTSR\d{9}(?:-\d{2})?$", "QA": r"^QA\d{11}$", 
            "BR": r"^BR\d{11}$", "ER": r"^ER\d{6}$", "IR": r"^IR-\d+-3DEXPERIENCER\d{4}x$", 
            "CRIT": r"^CRITSR\d{8}(?:-\d{2})?$", "JIRA": r"^(?:SPCK-)?\d{5}$", "SIT": r"^\d{15}$", "CTC": r"^[Cc]\d{11}$"
        }
        for key, reg in patterns.items():
            if re.match(reg, text, re.I):
                t = text.upper()
                if key == "SR": return f"{BASE_SR}{t if len(t)==13 else t+'-01'}"
                if key == "NTSR": return f"{BASE_NTSR}{t if len(t)==16 else t+'-01'}"
                if key in ["QA", "BR"]: return f"{BASE_QA}{t}"
                if key == "ER": return f"{BASE_ER}{t}"
                if key == "IR": return f"{BASE_IR}{text}"
                if key == "CRIT": return f"{BASE_CRIT}{t if '-' in t else t+'-01'}"
                if key == "JIRA": return f"{BASE_JIRA}{t if t.startswith('SPCK-') else 'SPCK-'+t}"
                if key == "SIT": return f"{BASE_SIT}{t}"
                if key == "CTC": return f"{BASE_CTC}{t}"
        return None

    def _build_help_content(self, parent):
        """Left pane of the Help tab — static usage text. Used to be a
        separate popup opened via typing '--h'/'--help'; now it's just
        always-visible content inside the Help tab."""
        txt = tk.Text(parent, bg=BG_COLOR, fg=TEXT_COLOR, font=("Segoe UI", 10), padx=15, pady=15, bd=0)
        txt.pack(fill="both", expand=True)
        help_content = """
 [ Shortcuts / Commands ]
  - abaqus cae / abq2026 cae   : Launch Abaqus CAE
  - --update data               : Manually update the database (search_data.db) — all Tier 1-4
  - --update data tier 1        : Update Tier 1 only (office/pdf)
  - --update data tier 1,2      : Update Tier 1 and Tier 2 only (comma/space separated, multiple OK)
  - Update DB button             : Type "tier 1,2" (or just "1,2") into the Searchbox first,
                                   then click — updates only that Tier (empty box = all Tier 1-4)
  - ESC, Alt+F4, --exit, --quit  : Exit

 [ Quick Search (URL redirect) ]
  - Technical SR          : e.g. SR01403940-01~09 / SR01403940
  - NonTechnical SR       : e.g. NTSR000187735-01~09 / NTSR000187735
  - CRITSIT               : e.g. CRITSR01403947-01~09 / CRITSR01403947
  - QA/BR                 : e.g. QA00000329750 / BR10000406277
  - ER                    : e.g. ER140863
  - IR                    : e.g. IR-1424613-3DEXPERIENCER2022x
  - Simpack JIRA          : e.g. SPCK-73341 / 73341
  - DS user Trigram       : e.g. ddh, ddh1~19...
  - SiteID/ContactID      : e.g. 200000000112873 / C00007548951

 [ File / Folder Name Search ]
  - keyword               : Search file content / file name / folder name
  - space                 : Multiple keywords (AND search)

 [ Search Results ]
  - Double-click           : Open the file's location in Explorer
  - Second click            : Select/copy text
  - Right-click             : Open Folder, Copy File Name, Copy File Path

 [ History ]
  - Help tab (right side)   : Search History is always visible here
  - Right-click              : "Search again" to re-run a past search
  - Export Excel           : Save history to Excel
  - Double-click             : Select/copy text

 [ Ramp Light ]
  - ● Red                    : search_data.db not created yet (run --update data)
  - ● Yellow                 : search_data.db exists but content not fully indexed yet
  - ● Yellow (blinking)      : Update DB is currently scanning/extracting content
  - ● Green                  : BM25 content ready — Search/Advanced usable (AI may still be building)
  - ● Blue                   : BM25 content AND AI embeddings both fully up to date

 [ File Content Search ]
  - keyword               : Text files, Outlook, OneNote, Office & PDF files

 [ AI Search Model ]
  - Model dropdown : Inside the AI search result window, next to the 🤖 AI Search button
  - Options: Jina-v3 / BGE-Gemma2
  - Each model uses its own data table — switching does not delete the other model's data
  - Run --update data once per model to build its index (can switch freely after that)
"""
        txt.insert("1.0", help_content.strip()); txt.config(state="disabled")

if __name__ == "__main__":
    root = tk.Tk()
    try:
        ttk.Style(root).configure("Treeview", rowheight=20)  # v4.9: room for 16px icons
    except Exception:
        pass
    app = RealtimeSmartSearchApp(root); root.mainloop()
